from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# SAP color palette
SAP_BLUE = RGBColor(0x00, 0x6D, 0x9E)       # SAP primary blue
SAP_GOLD = RGBColor(0xF0, 0xAB, 0x00)        # SAP gold/yellow
SAP_DARK = RGBColor(0x1A, 0x2E, 0x44)        # dark navy
SAP_LIGHT = RGBColor(0xF5, 0xF6, 0xF7)       # light background
SAP_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SAP_GRAY = RGBColor(0x60, 0x6A, 0x76)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_size=12, bold=False, color=SAP_DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=None):
    """Dark top bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.4, fill_color=SAP_DARK)
    add_rect(slide, 0, 1.4, 13.33, 0.07, fill_color=SAP_GOLD)
    add_text(slide, title, 0.5, 0.15, 12, 0.7, font_size=28, bold=True,
             color=SAP_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.82, 12, 0.5, font_size=14, bold=False,
                 color=SAP_GOLD, align=PP_ALIGN.LEFT)


def footer_bar(slide, note="SAP HANA PLA Analysis Process  |  SWAT Team"):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=SAP_DARK)
    add_text(slide, note, 0.4, 7.13, 12.5, 0.3, font_size=9,
             color=SAP_GRAY, align=PP_ALIGN.LEFT)


# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.33, 7.5, fill_color=SAP_DARK)
add_rect(s1, 0, 5.8, 13.33, 0.18, fill_color=SAP_GOLD)

# Decorative accent
add_rect(s1, 0, 0, 0.25, 7.5, fill_color=SAP_GOLD)

add_text(s1, "PLA HANA", 0.6, 1.4, 12, 1.1, font_size=52, bold=True,
         color=SAP_WHITE, align=PP_ALIGN.LEFT)
add_text(s1, "ANALYSIS PROCESS", 0.6, 2.4, 12, 1.0, font_size=38, bold=False,
         color=SAP_GOLD, align=PP_ALIGN.LEFT)
add_text(s1, "End-to-End Workflow for SAP HANA Health Check & Reporting",
         0.6, 3.5, 11, 0.6, font_size=16, italic=True,
         color=SAP_LIGHT, align=PP_ALIGN.LEFT)
add_text(s1, "SWAT Team  |  2026", 0.6, 6.0, 8, 0.4, font_size=12,
         color=SAP_GRAY, align=PP_ALIGN.LEFT)


# ── SLIDE 2: PROCESS OVERVIEW ─────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.33, 7.5, fill_color=SAP_LIGHT)
header_bar(s2, "Process Overview", "Four-phase end-to-end workflow")
footer_bar(s2)

steps = [
    ("1", "Pre-Data\nCollection", SAP_BLUE),
    ("2", "Health Check\nAnalysis", RGBColor(0x00, 0x8A, 0x5E)),
    ("3", "AI Analysis\nof Report", RGBColor(0x7B, 0x3F, 0xA0)),
    ("4", "Copy Results\nto PLA File", RGBColor(0xD0, 0x6A, 0x00)),
]

box_w, box_h = 2.4, 2.2
gap = 0.45
start_x = 0.6
top_y = 2.0

for i, (num, label, color) in enumerate(steps):
    x = start_x + i * (box_w + gap)

    # Shadow
    add_rect(s2, x + 0.07, top_y + 0.07, box_w, box_h, fill_color=RGBColor(0xCC, 0xCC, 0xCC))
    # Box
    add_rect(s2, x, top_y, box_w, box_h, fill_color=color)
    # Number badge
    add_rect(s2, x + 0.1, top_y + 0.1, 0.5, 0.5, fill_color=SAP_GOLD)
    add_text(s2, num, x + 0.1, top_y + 0.08, 0.5, 0.5, font_size=18, bold=True,
             color=SAP_DARK, align=PP_ALIGN.CENTER)
    # Label
    add_text(s2, label, x + 0.08, top_y + 0.72, box_w - 0.16, 1.0,
             font_size=15, bold=True, color=SAP_WHITE, align=PP_ALIGN.CENTER)

    # Arrow (except after last)
    if i < len(steps) - 1:
        ax = x + box_w + 0.05
        ay = top_y + box_h / 2 - 0.18
        add_text(s2, "▶", ax, ay, gap, 0.4, font_size=20, bold=True,
                 color=SAP_GOLD, align=PP_ALIGN.CENTER)

add_text(s2,
         "Each phase produces a defined outcome that feeds the next phase.",
         0.6, 4.55, 12, 0.4, font_size=12, italic=True,
         color=SAP_GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 3: PHASE 1 ──────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.33, 7.5, fill_color=SAP_LIGHT)
header_bar(s3, "Phase 1 — Pre-Data Collection", "Gather all required credentials and identifiers before starting")
footer_bar(s3)

# Phase badge
add_rect(s3, 0.5, 1.65, 0.55, 0.55, fill_color=SAP_GOLD)
add_text(s3, "1", 0.5, 1.63, 0.55, 0.55, font_size=22, bold=True,
         color=SAP_DARK, align=PP_ALIGN.CENTER)

# Gather card
add_rect(s3, 1.2, 1.65, 5.5, 3.6, fill_color=SAP_WHITE,
         line_color=SAP_BLUE, line_width=1.5)
add_rect(s3, 1.2, 1.65, 5.5, 0.5, fill_color=SAP_BLUE)
add_text(s3, "  GATHER", 1.2, 1.65, 5.5, 0.5, font_size=13, bold=True,
         color=SAP_WHITE)

items = [
    "▸  External ID information",
    "      (provided by Team Lead or Project Manager)",
    "▸  CAM credentials",
    "▸  Host name",
    "▸  Port",
    "▸  Tenant database name",
]
for j, item in enumerate(items):
    add_text(s3, item, 1.35, 2.28 + j * 0.42, 5.2, 0.4,
             font_size=12, color=SAP_DARK)

# Outcome card
add_rect(s3, 7.2, 1.65, 5.5, 1.6, fill_color=SAP_WHITE,
         line_color=SAP_GOLD, line_width=2)
add_rect(s3, 7.2, 1.65, 5.5, 0.5, fill_color=SAP_GOLD)
add_text(s3, "  OUTCOME", 7.2, 1.65, 5.5, 0.5, font_size=13, bold=True,
         color=SAP_DARK)
add_text(s3, "General information data package\nready for analysis phase",
         7.35, 2.28, 5.1, 0.9, font_size=13, color=SAP_DARK)

# Arrow between cards
add_text(s3, "▶", 6.6, 2.6, 0.6, 0.5, font_size=22, bold=True,
         color=SAP_GOLD, align=PP_ALIGN.CENTER)


# ── SLIDE 4: PHASE 2 ──────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.33, 7.5, fill_color=SAP_LIGHT)
header_bar(s4, "Phase 2 — Health Check Analysis Execution",
           "Run diagnostics on the HANA host")
footer_bar(s4)

add_rect(s4, 0.5, 1.65, 0.55, 0.55, fill_color=RGBColor(0x00, 0x8A, 0x5E))
add_text(s4, "2", 0.5, 1.63, 0.55, 0.55, font_size=22, bold=True,
         color=SAP_WHITE, align=PP_ALIGN.CENTER)

tasks_color = RGBColor(0x00, 0x8A, 0x5E)

# Tasks card
add_rect(s4, 1.2, 1.65, 5.5, 2.8, fill_color=SAP_WHITE,
         line_color=tasks_color, line_width=1.5)
add_rect(s4, 1.2, 1.65, 5.5, 0.5, fill_color=tasks_color)
add_text(s4, "  TASKS", 1.2, 1.65, 5.5, 0.5, font_size=13, bold=True,
         color=SAP_WHITE)

task_items = [
    "▸  Run health_check analysis on host",
    "▸  Run Parameter checks on host",
]
for j, item in enumerate(task_items):
    add_text(s4, item, 1.35, 2.28 + j * 0.55, 5.2, 0.5,
             font_size=13, color=SAP_DARK)

# Outcome card
add_rect(s4, 7.2, 1.65, 5.5, 2.0, fill_color=SAP_WHITE,
         line_color=SAP_GOLD, line_width=2)
add_rect(s4, 7.2, 1.65, 5.5, 0.5, fill_color=SAP_GOLD)
add_text(s4, "  OUTCOME", 7.2, 1.65, 5.5, 0.5, font_size=13, bold=True,
         color=SAP_DARK)

outcome_lines = [
    "Plain text report containing:",
    "  • HealthCheck Analysis results",
    "  • Non-standard Parameters found",
]
for j, line in enumerate(outcome_lines):
    add_text(s4, line, 7.35, 2.28 + j * 0.42, 5.1, 0.4,
             font_size=12, color=SAP_DARK, bold=(j == 0))

add_text(s4, "▶", 6.6, 2.4, 0.6, 0.5, font_size=22, bold=True,
         color=SAP_GOLD, align=PP_ALIGN.CENTER)


# ── SLIDE 5: PHASE 3 ──────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.33, 7.5, fill_color=SAP_LIGHT)
header_bar(s5, "Phase 3 — AI Analysis of HealthCheck Report",
           "Leverage AI to transform raw output into an actionable report")
footer_bar(s5)

ai_color = RGBColor(0x7B, 0x3F, 0xA0)
add_rect(s5, 0.5, 1.65, 0.55, 0.55, fill_color=ai_color)
add_text(s5, "3", 0.5, 1.63, 0.55, 0.55, font_size=22, bold=True,
         color=SAP_WHITE, align=PP_ALIGN.CENTER)

# Input
add_rect(s5, 1.2, 1.65, 3.4, 1.5, fill_color=SAP_WHITE,
         line_color=SAP_GRAY, line_width=1)
add_rect(s5, 1.2, 1.65, 3.4, 0.45, fill_color=SAP_GRAY)
add_text(s5, "  INPUT", 1.2, 1.65, 3.4, 0.45, font_size=12, bold=True, color=SAP_WHITE)
add_text(s5, "Plain text HealthCheck\n& Parameters report", 1.35, 2.18, 3.1, 0.85,
         font_size=12, color=SAP_DARK)

add_text(s5, "▶", 4.7, 2.15, 0.6, 0.5, font_size=22, bold=True,
         color=ai_color, align=PP_ALIGN.CENTER)

# AI Process
add_rect(s5, 5.35, 1.65, 3.2, 1.5, fill_color=ai_color)
add_text(s5, "  AI ENGINE", 5.35, 1.65, 3.2, 0.45, font_size=12, bold=True, color=SAP_GOLD)
add_text(s5, "Intelligent analysis\n& report generation", 5.5, 2.18, 2.9, 0.85,
         font_size=12, color=SAP_WHITE)

add_text(s5, "▶", 8.65, 2.15, 0.6, 0.5, font_size=22, bold=True,
         color=ai_color, align=PP_ALIGN.CENTER)

# Outcome
add_rect(s5, 9.3, 1.65, 3.5, 1.5, fill_color=SAP_WHITE,
         line_color=SAP_GOLD, line_width=2)
add_rect(s5, 9.3, 1.65, 3.5, 0.45, fill_color=SAP_GOLD)
add_text(s5, "  OUTCOME", 9.3, 1.65, 3.5, 0.45, font_size=12, bold=True, color=SAP_DARK)
add_text(s5, "Interactive Analysis Report\nin HTML format", 9.45, 2.18, 3.2, 0.85,
         font_size=13, bold=True, color=SAP_DARK)

# Benefits
add_rect(s5, 1.2, 3.55, 11.6, 2.8, fill_color=SAP_WHITE,
         line_color=ai_color, line_width=1)
add_rect(s5, 1.2, 3.55, 11.6, 0.45, fill_color=ai_color)
add_text(s5, "  KEY BENEFITS OF AI-GENERATED REPORT", 1.2, 3.55, 11.6, 0.45,
         font_size=12, bold=True, color=SAP_WHITE)

benefits = [
    ("Interactive", "Navigate findings directly in the HTML report"),
    ("Prioritized", "Issues ranked by severity and impact"),
    ("Actionable", "Clear recommendations per finding"),
    ("Shareable", "Self-contained HTML, no additional tools required"),
]
for j, (title, desc) in enumerate(benefits):
    bx = 1.4 + j * 2.9
    add_rect(s5, bx, 4.1, 2.6, 1.85, fill_color=SAP_LIGHT,
             line_color=ai_color, line_width=0.5)
    add_text(s5, title, bx + 0.1, 4.18, 2.4, 0.4, font_size=12, bold=True, color=ai_color)
    add_text(s5, desc, bx + 0.1, 4.58, 2.4, 1.2, font_size=10, color=SAP_DARK)


# ── SLIDE 6: PHASE 4 ──────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 13.33, 7.5, fill_color=SAP_LIGHT)
header_bar(s6, "Phase 4 — Copy Results to PLA File",
           "Consolidate outputs into the standardized PLA deliverable")
footer_bar(s6)

pla_color = RGBColor(0xD0, 0x6A, 0x00)
add_rect(s6, 0.5, 1.65, 0.55, 0.55, fill_color=pla_color)
add_text(s6, "4", 0.5, 1.63, 0.55, 0.55, font_size=22, bold=True,
         color=SAP_WHITE, align=PP_ALIGN.CENTER)

# Database Tab card
add_rect(s6, 1.2, 1.65, 5.3, 4.5, fill_color=SAP_WHITE,
         line_color=pla_color, line_width=1.5)
add_rect(s6, 1.2, 1.65, 5.3, 0.5, fill_color=pla_color)
add_text(s6, "  DATABASE TAB", 1.2, 1.65, 5.3, 0.5, font_size=13, bold=True,
         color=SAP_WHITE)

db_items = [
    ("Screenshots", "Paste screenshots into the PLA file Database tab"),
    ("HC Cover", "Include Health Check Analysis Report cover page"),
    ("HC Link", "Add link to the actual full HTML report"),
]
for j, (label, desc) in enumerate(db_items):
    iy = 2.28 + j * 1.1
    add_rect(s6, 1.35, iy, 4.95, 0.9, fill_color=SAP_LIGHT,
             line_color=pla_color, line_width=0.5)
    add_text(s6, label, 1.5, iy + 0.05, 1.2, 0.38, font_size=11, bold=True, color=pla_color)
    add_text(s6, desc, 1.5, iy + 0.42, 4.6, 0.42, font_size=11, color=SAP_DARK)

# Action Items card
add_rect(s6, 6.9, 1.65, 5.9, 4.5, fill_color=SAP_WHITE,
         line_color=SAP_BLUE, line_width=1.5)
add_rect(s6, 6.9, 1.65, 5.9, 0.5, fill_color=SAP_BLUE)
add_text(s6, "  ACTION ITEMS TAB", 6.9, 1.65, 5.9, 0.5, font_size=13, bold=True,
         color=SAP_WHITE)

ai_items = [
    ("Excel Format", "Paste action items in standard Excel format"),
    ("HANA Only", "Scope: HANA database action items exclusively"),
    ("Dedicated Tab", "Place in the Action Items tab of PLA file"),
]
for j, (label, desc) in enumerate(ai_items):
    iy = 2.28 + j * 1.1
    add_rect(s6, 7.05, iy, 5.55, 0.9, fill_color=SAP_LIGHT,
             line_color=SAP_BLUE, line_width=0.5)
    add_text(s6, label, 7.2, iy + 0.05, 1.5, 0.38, font_size=11, bold=True, color=SAP_BLUE)
    add_text(s6, desc, 7.2, iy + 0.42, 5.2, 0.42, font_size=11, color=SAP_DARK)

# Note
add_rect(s6, 1.2, 6.3, 11.6, 0.5, fill_color=SAP_GOLD)
add_text(s6, "  NOTE: All outputs must be consolidated in the PLA file before closing the engagement.",
         1.2, 6.3, 11.6, 0.5, font_size=11, bold=True, color=SAP_DARK)


# ── SLIDE 7: SUMMARY ──────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, 13.33, 7.5, fill_color=SAP_DARK)
add_rect(s7, 0, 0, 0.25, 7.5, fill_color=SAP_GOLD)
add_rect(s7, 0, 5.95, 13.33, 0.12, fill_color=SAP_GOLD)

add_text(s7, "Process Summary", 0.6, 0.3, 12, 0.8, font_size=30, bold=True,
         color=SAP_WHITE)

summary = [
    ("Phase 1", "Pre-Data Collection",
     "External ID · CAM credentials · Host / Port / Tenant DB",
     "General information data"),
    ("Phase 2", "Health Check Execution",
     "health_check analysis · Parameter checks",
     "Plain text HC & Parameters report"),
    ("Phase 3", "AI Analysis",
     "AI-powered interpretation of HC output",
     "Interactive HTML report"),
    ("Phase 4", "PLA File Update",
     "Screenshots · HC cover + link · Action items (HANA)",
     "Completed PLA deliverable"),
]

row_colors = [SAP_BLUE, RGBColor(0x00, 0x8A, 0x5E),
              RGBColor(0x7B, 0x3F, 0xA0), RGBColor(0xD0, 0x6A, 0x00)]

# Table header
for ci, (col, cw, cx) in enumerate([
    ("Phase", 1.2, 0.5), ("Name", 2.4, 1.8), ("Key Activities", 5.0, 4.3), ("Outcome", 3.3, 9.4)
]):
    add_rect(s7, cx, 1.25, cw, 0.45, fill_color=SAP_GOLD)
    add_text(s7, col, cx + 0.08, 1.28, cw - 0.1, 0.38, font_size=11, bold=True, color=SAP_DARK)

for i, (phase, name, activities, outcome) in enumerate(summary):
    ry = 1.75 + i * 1.1
    rc = row_colors[i]

    add_rect(s7, 0.5, ry, 1.2, 1.0, fill_color=rc)
    add_text(s7, phase, 0.55, ry + 0.25, 1.1, 0.5, font_size=12, bold=True,
             color=SAP_WHITE, align=PP_ALIGN.CENTER)

    add_rect(s7, 1.8, ry, 2.4, 1.0, fill_color=RGBColor(0x22, 0x3A, 0x55))
    add_text(s7, name, 1.9, ry + 0.22, 2.2, 0.55, font_size=11, bold=True, color=SAP_WHITE)

    add_rect(s7, 4.3, ry, 5.0, 1.0, fill_color=RGBColor(0x1A, 0x2E, 0x44))
    add_text(s7, activities, 4.4, ry + 0.15, 4.8, 0.7, font_size=10, color=SAP_LIGHT)

    add_rect(s7, 9.4, ry, 3.3, 1.0, fill_color=rc)
    add_text(s7, outcome, 9.5, ry + 0.15, 3.1, 0.7, font_size=10, bold=True, color=SAP_WHITE)

add_text(s7, "SWAT Team  |  SAP HANA Practice  |  2026",
         0.5, 6.1, 12, 0.35, font_size=10, color=SAP_GRAY)

# ── SAVE ──────────────────────────────────────────────────────────────────────
out = r"C:\Users\I522148\OneDrive - SAP SE\SWAT\AI\Claude\working directory\PLA_HANA_Analysis_Process.pptx"
prs.save(out)
print(f"Saved: {out}")
