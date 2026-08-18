"""
SAP HANA Full Analysis Report Generator v5.0
- All v4 features retained
- New: "Export to PowerPoint" button enabled after HTML report is generated
  Parses the generated HTML and builds a 16:9 deck (one slide per section)
  using SAP colors, Segoe UI font, and the same design language as the HTML
"""

import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
import threading
import os
import re
import json
import base64
from datetime import datetime
from pathlib import Path
import anthropic

# ── PPTX availability ─────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

# ── Config ────────────────────────────────────────────────────────────────────
WORKING_DIR   = Path(r"C:\Users\I522148\OneDrive - SAP SE\SWAT\AI\Claude\working directory")
TEMPLATE_FILE = WORKING_DIR / "HANA_HealthCheck_Prompt_Template.md"
RELEASE_FILE  = WORKING_DIR / "HANA_latest_release.txt"
OUTPUT_DIR    = WORKING_DIR / "Results"
LOGO_PATH     = WORKING_DIR / "SAP_LOGO.png"
EDGE_PATH     = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
SETTINGS_FILE = Path(r"C:\Users\I522148\.claude\settings.json")

MAX_FILE_KB         = 2500
EXPECTED_OUT_TOKENS = 40000
MAX_OUTPUT_TOKENS   = 64000
APP_VERSION         = "v5.0"

# ── UI Colors ─────────────────────────────────────────────────────────────────
PRIMARY_BLUE = "#0057a8"
DARK_BLUE    = "#003366"
ACCENT_CYAN  = "#00b4e6"
PAGE_BG      = "#e8f4fd"
WHITE        = "#ffffff"
LIGHT_BLUE   = "#f0f7ff"
OK_GREEN     = "#5cb85c"
WARN_ORANGE  = "#f0ad4e"
DANGER_RED   = "#d9534f"
PURPLE       = "#6f42c1"

# ── Enhanced CSS (injected into every generated report) ───────────────────────
ENHANCED_CSS = """
/* ── Enhanced v4 additions ─── */
.consequence{background:#fff3f3;border-left:4px solid #d9534f;padding:8px 14px;
  margin:2px 0 6px 0;border-radius:4px;font-size:0.85rem;color:#555}
.esc-banner{background:#d9534f;color:#fff;padding:12px 20px;border-radius:6px;
  font-weight:700;font-size:1rem;margin-bottom:16px;display:flex;align-items:center;gap:10px}
.esc-banner.warn{background:#e65c00}
.exec-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:16px;margin:16px 0}
.exec-card{background:#fff;border-radius:8px;padding:16px;box-shadow:0 2px 6px rgba(0,0,0,.1)}
.exec-card.critical{border-top:4px solid #d9534f}
.exec-card.warning{border-top:4px solid #f0ad4e}
.exec-card.monthly{border-top:4px solid #0057a8}
.exec-card.ok{border-top:4px solid #5cb85c}
.exec-card h4{margin:0 0 6px 0;font-size:0.95rem;color:#003366}
.exec-card .biz-impact{font-size:0.82rem;color:#d9534f;font-weight:600;margin:4px 0}
.exec-card .biz-action{font-size:0.82rem;color:#555;margin:4px 0}
.health-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(160px,1fr));gap:10px;margin:12px 0}
.health-item{background:#fff;border-radius:6px;padding:10px 14px;display:flex;
  align-items:center;gap:10px;box-shadow:0 1px 4px rgba(0,0,0,.08)}
.health-item .hi-icon{font-size:1.4rem}
.health-item .hi-label{font-size:0.8rem;color:#666}
.health-item .hi-val{font-size:0.85rem;font-weight:700}
.ap-section{margin:12px 0}
.ap-section h3{font-size:0.9rem;color:#003366;margin:0 0 8px 0;
  padding:6px 0;border-bottom:2px solid #e8f0f8}
.ap-card{background:#fff;border-radius:8px;padding:14px 18px;margin:8px 0;
  box-shadow:0 2px 5px rgba(0,0,0,.08);border-left:5px solid #0057a8}
.ap-card.p1{border-left-color:#d9534f}
.ap-card.p2{border-left-color:#f0ad4e}
.ap-card.p3{border-left-color:#5cb85c}
.ap-card h4{margin:0 0 6px 0;font-size:0.95rem;color:#003366}
.ap-card p{margin:0;font-size:0.85rem;color:#555}
.ap-tags{display:flex;flex-wrap:wrap;gap:6px;margin-top:8px}
.ap-tag{font-size:0.72rem;font-weight:600;padding:2px 8px;border-radius:10px;
  background:#e8f4fd;color:#003366}
.ap-tag.urgent{background:#fde8e8;color:#d9534f}
.ap-tag.owner{background:#e8f8e8;color:#3a7a3a}
.ap-tag.effort{background:#fff8e1;color:#856404}
.ap-tag.prevents{background:#f0e8ff;color:#5a3a8a}
.kpi-critical{border-top:4px solid #d9534f !important}
.kpi-critical .kpi-val{color:#d9534f !important}
/* ── End enhanced v4 ─── */
"""


# ═════════════════════════════════════════════════════════════════════════════
# HTML → PPTX: parsing helpers
# ═════════════════════════════════════════════════════════════════════════════

def _strip_tags(s: str) -> str:
    """Strip all HTML tags and decode common entities."""
    s = re.sub(r"<[^>]+>", " ", s)
    for ent, ch in [("&amp;", "&"), ("&lt;", "<"), ("&gt;", ">"),
                    ("&nbsp;", " "), ("&#8203;", ""), ("&apos;", "'")]:
        s = s.replace(ent, ch)
    return " ".join(s.split()).strip()


def _parse_report_meta(html: str) -> dict:
    """Extract report-level metadata from the footer and header."""
    meta = {
        "sid": "UNKNOWN", "tenant": "", "host": "", "instance": "",
        "version": "", "system_type": "", "generated": "", "script": "",
    }
    # Footer: "Health Check Report | SID: HP4 (BHP) / Tenant: PS4 | Generated: 2026-08-12 | ..."
    footer_m = re.search(r"<footer[^>]*>(.*?)</footer>", html, re.DOTALL)
    if footer_m:
        ft = _strip_tags(footer_m.group(1))
        for pat, key in [
            (r"SID:\s*([^|]+)",       "sid"),
            (r"Tenant:\s*([^|]+)",    "tenant"),
            (r"Generated:\s*([^|]+)", "generated"),
            (r"Script:\s*([^|]+)",    "script"),
        ]:
            m = re.search(pat, ft, re.IGNORECASE)
            if m:
                meta[key] = m.group(1).strip()

    # Header metadata div: host, instance, version, type
    hdr_m = re.search(r'class=["\'][^"\']*header-meta[^"\']*["\'][^>]*>(.*?)</div>',
                      html, re.DOTALL)
    if hdr_m:
        hdr_text = _strip_tags(hdr_m.group(1))
        for pat, key in [
            (r"Host[:\s]+([^\s|·]+)",       "host"),
            (r"Instance[:\s]+([^\s|·]+)",   "instance"),
            (r"(\d+\.\d+\.\d+\.\d+)",       "version"),
            (r"(MDC[^|·<\n]*)",             "system_type"),
        ]:
            m = re.search(pat, hdr_text, re.IGNORECASE)
            if m:
                meta[key] = m.group(1).strip()

    return meta


def _parse_kpi_cards(html: str) -> list[tuple[str, str, str]]:
    """Return list of (value, label, color) from KPI summary cards."""
    results = []
    kpi_section_m = re.search(
        r'<div[^>]+class=["\'][^"\']*kpi-grid[^"\']*["\'][^>]*>(.*?)</div>\s*</div>',
        html, re.DOTALL)
    if not kpi_section_m:
        # fallback: find all kpi-card divs in the first 8000 chars
        kpi_section_m = re.search(r'((?:<div[^>]+class="kpi-card[^>]*>.*?</div>\s*){1,10})',
                                   html[:8000], re.DOTALL)
    if not kpi_section_m:
        return results

    chunk = kpi_section_m.group(1)
    for card_m in re.finditer(
        r'<div[^>]+class=["\'][^"\']*kpi-(?:val|value)[^"\']*["\'][^>]*'
        r'(?:style="[^"]*color:\s*([^;"]+)[^"]*")?[^>]*>(.*?)</div>'
        r'.*?<div[^>]+class=["\'][^"\']*kpi-label["\'][^>]*>(.*?)</div>',
        chunk, re.DOTALL
    ):
        color = (card_m.group(1) or "#0057a8").strip()
        val   = _strip_tags(card_m.group(2))
        lbl   = _strip_tags(card_m.group(3))
        if val and lbl:
            results.append((val, lbl, color))
    return results[:8]


def _parse_sections(html: str) -> list[dict]:
    """
    Split HTML by section divs and extract per-section data:
      id, icon, title, badge_text, badge_color, h4s, bullets, ap_tiers

    Robust: attribute order doesn't matter, handles class="section open",
    class="section collapsible", id before class, single or double quotes.
    """
    sections = []

    # Scan every <div ...> opening tag — check for class containing "section"
    # and a non-body id.  Attribute order and extra classes are irrelevant.
    starts = []
    for m in re.finditer(r'<div\b[^>]*>', html, re.IGNORECASE):
        tag = m.group(0)
        id_m  = re.search(r'\bid\s*=\s*["\']([^"\']+)["\']', tag, re.IGNORECASE)
        cls_m = re.search(r'\bclass\s*=\s*["\']([^"\']*)["\']', tag, re.IGNORECASE)
        if not id_m or not cls_m:
            continue
        sec_id  = id_m.group(1).strip()
        classes = cls_m.group(1).split()
        # Must have "section" as one of its classes
        # Must not be a body-container id (body-sx, body-s0, …)
        if "section" in classes and not sec_id.startswith("body-"):
            starts.append((m.start(), sec_id))

    if not starts:
        return sections

    for i, (pos, sid) in enumerate(starts):
        end = starts[i + 1][0] if i + 1 < len(starts) else len(html)
        chunk = html[pos:end]

        # Icon + title
        icon_m  = re.search(r'class=["\'][^"\']*section-icon[^"\']*["\'][^>]*>(.*?)</span>',  chunk, re.DOTALL)
        title_m = re.search(r'class=["\'][^"\']*section-title[^"\']*["\'][^>]*>(.*?)</span>', chunk, re.DOTALL)
        icon    = _strip_tags(icon_m.group(1))  if icon_m  else ""
        title   = _strip_tags(title_m.group(1)) if title_m else sid

        # Badge — first pill-XXX span anywhere in the section header area
        badge_m = re.search(
            r'class=["\'][^"\']*pill[^"\']*pill-([a-z]+)[^"\']*["\'][^>]*>(.*?)</span>',
            chunk, re.DOTALL)
        if not badge_m:
            # alternate order: pill-red pill (some generators swap them)
            badge_m = re.search(
                r'class=["\'][^"\']*pill-([a-z]+)[^"\']*["\'][^>]*>(.*?)</span>',
                chunk, re.DOTALL)
        badge_color = badge_m.group(1).strip() if badge_m else "blue"
        badge_text  = _strip_tags(badge_m.group(2)) if badge_m else ""

        # Body content — everything after the section-header div
        body_m = re.search(r'class=["\'][^"\']*section-body[^"\']*["\'][^>]*>(.*)', chunk, re.DOTALL)
        body   = body_m.group(1) if body_m else chunk

        # h4 headings (key findings)
        h4s = [_strip_tags(h) for h in re.findall(r"<h4[^>]*>(.*?)</h4>", body, re.DOTALL)]
        h4s = [h for h in h4s if h][:5]

        # Action-list li items
        bullets = [_strip_tags(li)
                   for li in re.findall(r"<li[^>]*>(.*?)</li>", body, re.DOTALL)]
        bullets = [b for b in bullets if len(b) > 5][:8]

        # Fallback for table-only sections (e.g. Additional Risks):
        # extract each <tr> as "first-cell  —  last-cell"
        if not h4s and not bullets:
            for row_m in re.finditer(r"<tr\b[^>]*>(.*?)</tr>", body, re.DOTALL):
                cells = [_strip_tags(c)
                         for c in re.findall(r"<td[^>]*>(.*?)</td>", row_m.group(1), re.DOTALL)]
                cells = [c for c in cells if c]
                if len(cells) >= 2:
                    bullets.append(f"{cells[0]}  —  {cells[-1]}")
            bullets = bullets[:8]

        # Action plan tiers (for sap section)
        ap_tiers: dict[str, list[str]] = {"p1": [], "p2": [], "p3": []}
        for tier in ("p1", "p2", "p3"):
            for card_m in re.finditer(
                rf'<div\b[^>]*class=["\'][^"\']*\bap-card\b[^"\']*\b{tier}\b[^"\']*["\'][^>]*>'
                r'.*?<h4[^>]*>(.*?)</h4>.*?<p[^>]*>(.*?)</p>',
                body, re.DOTALL
            ):
                ap_tiers[tier].append(
                    _strip_tags(card_m.group(1)) + " — " + _strip_tags(card_m.group(2))
                )

        sections.append({
            "id": sid, "icon": icon, "title": title,
            "badge_color": badge_color, "badge_text": badge_text,
            "h4s": h4s, "bullets": bullets, "ap_tiers": ap_tiers,
        })

    return sections


# ═════════════════════════════════════════════════════════════════════════════
# PPTX builder helpers
# ═════════════════════════════════════════════════════════════════════════════

# Slide dimensions 16:9
_SW = Inches(13.33)
_SH = Inches(7.5)

# PPTX SAP colors
_PC = {
    "dark":   RGBColor(0x00, 0x33, 0x66),
    "blue":   RGBColor(0x00, 0x57, 0xa8),
    "cyan":   RGBColor(0x00, 0xb4, 0xe6),
    "red":    RGBColor(0xd9, 0x53, 0x4f),
    "orange": RGBColor(0xf0, 0xad, 0x4e),
    "green":  RGBColor(0x5c, 0xb8, 0x5c),
    "grey":   RGBColor(0x66, 0x66, 0x66),
    "white":  RGBColor(0xff, 0xff, 0xff),
    "bg":     RGBColor(0xe8, 0xf4, 0xfd),
    "card":   RGBColor(0xf0, 0xf7, 0xff),
    "body":   RGBColor(0x21, 0x25, 0x29),
    "edge":   RGBColor(0xe0, 0xec, 0xf8),
    "red_bg": RGBColor(0xff, 0xe8, 0xe8),
    "org_bg": RGBColor(0xff, 0xfd, 0xf0),
    "grn_bg": RGBColor(0xf0, 0xff, 0xf4),
}

def _badge_color_rgb(css_class: str) -> RGBColor:
    return {
        "red": _PC["red"], "orange": _PC["orange"],
        "green": _PC["green"], "blue": _PC["blue"],
    }.get(css_class.split()[0], _PC["blue"])


def _R(slide, x, y, w, h, bg: RGBColor, border: RGBColor = None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = bg
    if border:
        s.line.color.rgb = border; s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s


def _T(slide, text, x, y, w, h, sz=10, bold=False,
       color: RGBColor = None, align=PP_ALIGN.LEFT, italic=False):
    color = color or _PC["body"]
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Segoe UI"
    return tb


def _slide_header(slide, icon: str, title: str, badge: str, badge_col: RGBColor):
    _R(slide, 0, 0, _SW, Inches(0.88), _PC["dark"])
    _T(slide, f"{icon}  {title}",
       Inches(0.3), Inches(0.1), Inches(10.5), Inches(0.7),
       sz=18, bold=True, color=_PC["white"])
    bw = Inches(2.1)
    _R(slide, _SW - bw - Inches(0.22), Inches(0.19), bw, Inches(0.52), badge_col)
    _T(slide, badge,
       _SW - bw - Inches(0.22), Inches(0.19), bw, Inches(0.52),
       sz=10, bold=True, color=_PC["white"], align=PP_ALIGN.CENTER)


def _info_card(slide, title: str, lines: list[str],
               x, y, w, h, accent: RGBColor = None):
    accent = accent or _PC["blue"]
    bg_map = {
        id(_PC["red"]): _PC["red_bg"],
        id(_PC["orange"]): _PC["org_bg"],
        id(_PC["green"]): _PC["grn_bg"],
    }
    bg = bg_map.get(id(accent), _PC["card"])
    _R(slide, x, y, w, h, bg, accent)
    _R(slide, x, y, Inches(0.08), h, accent)
    _T(slide, title,
       x + Inches(0.16), y + Inches(0.08), w - Inches(0.24), Inches(0.28),
       sz=9.5, bold=True, color=_PC["dark"])
    body = "\n".join(f"\u25B6  {l}" for l in lines if l)
    _T(slide, body,
       x + Inches(0.16), y + Inches(0.38), w - Inches(0.24), h - Inches(0.44),
       sz=8.5, color=_PC["body"])


def _kpi_box(slide, val: str, lbl: str, x, y, w, h, vc: RGBColor = None):
    vc = vc or _PC["blue"]
    _R(slide, x, y, w, h, _PC["white"], _PC["edge"])
    _R(slide, x, y, w, Inches(0.06), vc)
    _T(slide, val, x, y + Inches(0.1), w, Inches(0.44),
       sz=20, bold=True, color=vc, align=PP_ALIGN.CENTER)
    _T(slide, lbl, x, y + Inches(0.54), w, Inches(0.26),
       sz=7.5, color=_PC["grey"], align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# PPTX slide generators
# ═════════════════════════════════════════════════════════════════════════════

def _slide_title(prs, meta: dict, logo_path: Path):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _R(sl, 0, 0, _SW, _SH, _PC["dark"])
    _R(sl, 0, _SH - Inches(0.1), _SW, Inches(0.1), _PC["cyan"])
    if logo_path and logo_path.exists():
        try:
            sl.shapes.add_picture(str(logo_path), Inches(0.45), Inches(0.35), height=Inches(0.55))
        except Exception:
            pass
    _T(sl, "SAP HANA Database  \u00B7  Health Check Report",
       Inches(0.45), Inches(1.1), Inches(12), Inches(0.65), sz=18, color=_PC["cyan"])
    sid_line = meta.get("sid", "UNKNOWN")
    tenant = meta.get("tenant", "")
    if tenant:
        sid_line += f"  \u2014  {tenant}"
    _T(sl, sid_line,
       Inches(0.45), Inches(1.82), Inches(12), Inches(1.0),
       sz=38, bold=True, color=_PC["white"])
    _R(sl, Inches(0.45), Inches(2.96), Inches(7.5), Inches(0.04), _PC["cyan"])
    rows = [
        ("Host",         meta.get("host", "")),
        ("Instance",     meta.get("instance", "")),
        ("HANA Version", meta.get("version", "")),
        ("System Type",  meta.get("system_type", "")),
        ("Generated",    meta.get("generated", "")),
        ("Script",       meta.get("script", "")),
    ]
    for i, (lbl, val) in enumerate(rows):
        if not val:
            continue
        ry = Inches(3.12 + i * 0.5)
        _T(sl, lbl, Inches(0.45), ry, Inches(2.4), Inches(0.42),
           sz=10, bold=True, color=RGBColor(0xa8, 0xd0, 0xf0))
        _T(sl, val, Inches(2.95), ry, Inches(9.8), Inches(0.42),
           sz=10, color=_PC["white"])
    _T(sl, f"Generated by SAP HANA Analysis Report Generator {APP_VERSION}",
       Inches(0.45), _SH - Inches(0.54), Inches(12), Inches(0.38),
       sz=8.5, color=RGBColor(0x60, 0x88, 0xb8))


def _slide_section(prs, sec: dict):
    """
    Generic section slide: header + up to 3 content cards.
    Layout adapts based on how many h4s / bullets the section has.
    """
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _R(sl, 0, 0, _SW, _SH, _PC["bg"])

    badge_col = _badge_color_rgb(sec["badge_color"])
    _slide_header(sl, sec["icon"], sec["title"], sec["badge_text"], badge_col)

    h4s     = sec["h4s"]
    bullets = sec["bullets"]
    ap      = sec["ap_tiers"]

    # ── Action Plan section: 3-tier card layout ───────────────────────────
    if sec["id"] == "sap" and any(ap.values()):
        _render_action_plan(sl, ap)
        return

    # ── Generic content layout ────────────────────────────────────────────
    # Split bullets into "findings" (h4-level, usually short) and "actions"
    findings = h4s[:3]
    actions  = bullets[:6]

    y_cursor = Inches(1.0)
    available = _SH - Inches(1.08)  # total content height

    if findings and actions:
        # Two cards: findings top, actions below
        fh = min(available * 0.48, Inches(2.8))
        ah = min(available - fh - Inches(0.14), Inches(3.2))
        _info_card(sl, "Key Findings",  findings, Inches(0.28), y_cursor,
                   _SW - Inches(0.56), fh, badge_col)
        _info_card(sl, "Recommended Actions", actions,
                   Inches(0.28), y_cursor + fh + Inches(0.14),
                   _SW - Inches(0.56), ah)
    elif findings:
        # Full-height findings card
        fh = min(available, Inches(5.8))
        _info_card(sl, "Key Findings", findings,
                   Inches(0.28), y_cursor, _SW - Inches(0.56), fh, badge_col)
    elif actions:
        # Full-height actions card
        ah = min(available, Inches(5.8))
        _info_card(sl, "Details & Actions", actions,
                   Inches(0.28), y_cursor, _SW - Inches(0.56), ah, badge_col)
    else:
        # Empty section placeholder
        _info_card(sl, "Section", ["No content extracted — see HTML report for details."],
                   Inches(0.28), y_cursor, _SW - Inches(0.56), Inches(1.5))


def _render_action_plan(sl, ap: dict):
    """Render the 3-tier action plan layout."""
    tier_cfg = [
        ("p1", "\U0001F534  Priority 1 \u2014 Immediate (\u226448h)", _PC["red"]),
        ("p2", "\U0001F7E1  Priority 2 \u2014 This Week",              _PC["orange"]),
        ("p3", "\U0001F7E2  Priority 3 \u2014 This Month",             _PC["green"]),
    ]
    bg_map = {
        id(_PC["red"]):    _PC["red_bg"],
        id(_PC["orange"]): _PC["org_bg"],
        id(_PC["green"]):  _PC["grn_bg"],
    }
    y = Inches(1.0)
    for tier_key, hdr_txt, col in tier_cfg:
        items = ap.get(tier_key, [])
        if not items:
            continue
        # Tier header strip
        _R(sl, Inches(0.28), y, _SW - Inches(0.56), Inches(0.34), col)
        _T(sl, hdr_txt, Inches(0.42), y + Inches(0.03), Inches(12.5), Inches(0.26),
           sz=10, bold=True, color=_PC["white"])
        y += Inches(0.34)

        # Item cards — 2 per row
        cw = Inches(6.2); gap = Inches(0.2); ch = Inches(0.72)
        for idx, item in enumerate(items[:4]):
            cx = Inches(0.28) + (idx % 2) * (cw + gap)
            cy = y + (idx // 2) * (ch + Inches(0.08))
            _R(sl, cx, cy, cw, ch, bg_map.get(id(col), _PC["card"]), col)
            _R(sl, cx, cy, Inches(0.07), ch, col)
            parts = item.split(" \u2014 ", 1)
            title_t = parts[0][:60]
            desc_t  = parts[1][:120] if len(parts) > 1 else ""
            _T(sl, title_t, cx + Inches(0.16), cy + Inches(0.05),
               cw - Inches(0.24), Inches(0.28), sz=9, bold=True, color=_PC["dark"])
            if desc_t:
                _T(sl, desc_t, cx + Inches(0.16), cy + Inches(0.34),
                   cw - Inches(0.24), Inches(0.32), sz=8, color=_PC["grey"])
        row_count = (len(items[:4]) + 1) // 2
        y += row_count * (ch + Inches(0.08)) + Inches(0.12)
        if y > _SH - Inches(0.3):
            break


# ═════════════════════════════════════════════════════════════════════════════
# Main PPTX builder
# ═════════════════════════════════════════════════════════════════════════════

def build_pptx_from_html(html_path: Path, logo_path: Path) -> Path:
    """
    Parse html_path, build a 16:9 PowerPoint deck, save next to the HTML file
    with a .pptx extension, and return the output path.
    """
    if not PPTX_OK:
        raise RuntimeError(
            "python-pptx is not installed.\n\n"
            "Run:  pip install python-pptx"
        )

    html = html_path.read_text(encoding="utf-8", errors="replace")

    meta     = _parse_report_meta(html)
    sections = _parse_sections(html)

    prs = Presentation()
    prs.slide_width  = _SW
    prs.slide_height = _SH

    _slide_title(prs, meta, logo_path)

    for sec in sections:
        _slide_section(prs, sec)

    output = html_path.with_suffix(".pptx")
    prs.save(str(output))
    return output


# ═════════════════════════════════════════════════════════════════════════════
# v4 helpers (unchanged)
# ═════════════════════════════════════════════════════════════════════════════

def load_text(path: Path) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def load_proxy_config() -> tuple[str, str, str]:
    try:
        cfg = json.loads(SETTINGS_FILE.read_text(encoding="utf-8"))
        env = cfg.get("env", {})
        base_url = env.get("ANTHROPIC_BASE_URL", "https://api.anthropic.com")
        token    = env.get("ANTHROPIC_AUTH_TOKEN", "")
        model    = env.get("ANTHROPIC_MODEL", "claude-sonnet-latest")
        return base_url, token, model
    except Exception:
        return "https://api.anthropic.com", "", "claude-sonnet-latest"


def get_sap_logo_b64() -> str:
    if LOGO_PATH.exists():
        with open(LOGO_PATH, "rb") as f:
            return base64.b64encode(f.read()).decode("utf-8")
    return ""


def extract_sid(text: str) -> str:
    for pat in [r"\[environment\].*?SID[:\s]+([A-Z0-9]{3})",
                r"SID[:\s]+([A-Z0-9]{3})"]:
        m = re.search(pat, text, re.IGNORECASE)
        if m:
            return m.group(1).upper()
    return "UNKNOWN"


def build_prompt(template: str, release_info: str, health_check_text: str) -> str:
    prompt = template.replace(
        "<<< REPLACE THIS LINE WITH THE FULL HANA HEALTH CHECK SCRIPT OUTPUT >>>",
        health_check_text,
    )
    if release_info:
        prompt += f"\n\n[HANA Latest Release Reference]\n{release_info}\n"
    return prompt


# ── Post-processing pipeline ──────────────────────────────────────────────────

def _strip_code_fences(html: str) -> str:
    html = re.sub(r"^```html\s*", "", html.strip())
    return re.sub(r"\s*```$", "", html).strip()


def _inject_enhanced_css(html: str) -> str:
    if "Enhanced v4 additions" in html:
        return html
    return html.replace("</style>", ENHANCED_CSS + "\n</style>", 1)


def _inject_sap_logo(html: str, logo_b64: str) -> str:
    if not logo_b64:
        return html
    img_left = (
        f'<img src="data:image/png;base64,{logo_b64}" '
        'height="32" alt="SAP" style="vertical-align:middle">'
    )
    html = re.sub(
        r'(<div[^>]+class=["\'][^"\']*header-logo[^"\']*["\'][^>]*>).*?(</div>)',
        rf'\1{img_left}\2', html, flags=re.DOTALL,
    )
    img_right = (
        f'\n<div style="margin-left:auto;display:flex;align-items:center;'
        f'flex-shrink:0;padding-right:6px">'
        f'<img src="data:image/png;base64,{logo_b64}" height="26" alt="SAP" '
        f'style="opacity:0.85;vertical-align:middle"></div>'
    )
    m = re.search(r'(</div>)\s*<div[^>]+class=["\'][^"\']*header-meta', html)
    if not m:
        m = re.search(r'(</div>)\s*</div>\s*<div[^>]+class=["\'][^"\']*nav-bar', html)
    if m:
        insert_at = m.start(1)
        html = html[:insert_at] + img_right + html[insert_at:]
    return html


def _inject_footer_version(html: str) -> str:
    today = datetime.now().strftime("%Y-%m-%d")
    marker = f"Enhanced: {today}"
    if marker in html:
        return html
    return re.sub(r"(Generated:\s*[^<|]+)", rf"\1  |  {marker}", html, count=1)


def _ensure_js(html: str) -> str:
    if "<script" in html.lower():
        return html
    FIX_JS = """
<script>
function toggleSection(id) {
  var body = document.getElementById('body-' + id);
  var chev = document.getElementById('chev-' + id);
  if (!body) return;
  var isOpen = body.classList.contains('open');
  body.classList.toggle('open', !isOpen);
  if (chev) chev.textContent = isOpen ? '\u25bc' : '\u25b2';
}
document.addEventListener('DOMContentLoaded', function() {
  document.querySelectorAll('.nav-pill').forEach(function(pill) {
    pill.addEventListener('click', function() {
      var m = (pill.getAttribute('onclick') || '').match(/'([^']+)'/);
      var el = m && document.getElementById(m[1]);
      if (el) el.scrollIntoView({behavior:'smooth'});
      document.querySelectorAll('.nav-pill').forEach(function(p){ p.classList.remove('active'); });
      pill.classList.add('active');
    });
  });
  var tip = document.getElementById('tooltip');
  if (tip) {
    document.querySelectorAll('.term').forEach(function(el) {
      el.addEventListener('mousemove', function(e) {
        var def = window.GLOSSARY && window.GLOSSARY[el.dataset.term];
        if (!def) return;
        tip.innerHTML = '<strong style="color:#a8d0f0">' + el.dataset.term + '</strong><br>' + def;
        tip.style.display = 'block';
        tip.style.left = (e.clientX + 14) + 'px';
        tip.style.top  = (e.clientY + 14) + 'px';
      });
      el.addEventListener('mouseleave', function() { tip.style.display = 'none'; });
    });
  }
});
</script>
"""
    if "</body>" in html:
        return html.replace("</body>", FIX_JS + "</body>", 1)
    return html + "\n" + FIX_JS + "\n</body>\n</html>"


def post_process_html(html: str, logo_b64: str) -> str:
    html = _strip_code_fences(html)
    html = _inject_enhanced_css(html)
    html = _inject_sap_logo(html, logo_b64)
    html = _inject_footer_version(html)
    html = _ensure_js(html)
    return html


# ═════════════════════════════════════════════════════════════════════════════
# App
# ═════════════════════════════════════════════════════════════════════════════

class AnalysisApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title(f"SAP HANA Full Analysis Report Generator {APP_VERSION}")
        self.configure(bg=PAGE_BG)
        self.minsize(920, 700)
        self.resizable(True, True)

        self._selected_file:    Path | None = None
        self._last_output_file: Path | None = None
        self._status_var   = tk.StringVar(value="Ready — select a health check file to begin.")
        self._progress_var = tk.IntVar(value=0)
        self._pct_var      = tk.StringVar(value="")
        self._cancel_event = threading.Event()
        self._logo_b64     = get_sap_logo_b64()

        base_url, token, model = load_proxy_config()
        self._base_url = base_url
        self._token    = token
        self._model    = model

        self._build_ui()

    # ── UI construction ───────────────────────────────────────────────────────

    def _build_ui(self):
        self._build_header()
        self._build_content()
        self._build_footer()

    def _load_logo_tk_image(self, target_height: int = 42):
        if not self._logo_b64:
            return None
        try:
            from PIL import Image, ImageTk
            import io
            raw = base64.b64decode(self._logo_b64)
            img = Image.open(io.BytesIO(raw))
            w, h = img.size
            new_w = max(1, int(w * target_height / h))
            img = img.resize((new_w, target_height), Image.LANCZOS)
            return ImageTk.PhotoImage(img)
        except Exception:
            pass
        try:
            photo = tk.PhotoImage(data=self._logo_b64)
            factor = max(1, photo.height() // target_height)
            return photo.subsample(factor, factor) if factor > 1 else photo
        except Exception:
            return None

    def _build_header(self):
        hdr = tk.Frame(self, bg=DARK_BLUE)
        hdr.pack(fill="x")

        left = tk.Frame(hdr, bg=DARK_BLUE)
        left.pack(side="left", padx=20, pady=12)
        tk.Label(left, text="SAP HANA Full Analysis Report Generator",
                 bg=DARK_BLUE, fg=WHITE, font=("Segoe UI", 15, "bold")).pack(anchor="w")
        tk.Label(left, text=f"Powered by Anthropic SDK  |  {APP_VERSION}  |  {self._model}",
                 bg=DARK_BLUE, fg="#a8d0f0", font=("Segoe UI", 8)).pack(anchor="w")

        right = tk.Frame(hdr, bg=DARK_BLUE)
        right.pack(side="right", padx=16, pady=8)
        tk.Label(right, text="RDE LAC", bg="#e65c00", fg=WHITE,
                 font=("Segoe UI", 9, "bold"), padx=12, pady=4).pack(side="right", padx=(10, 0))
        self._tk_logo = self._load_logo_tk_image(target_height=42)
        if self._tk_logo:
            tk.Label(right, image=self._tk_logo, bg=DARK_BLUE, bd=0).pack(side="right")
        else:
            tk.Label(right, text="SAP_LOGO.png not found", bg=DARK_BLUE, fg=WARN_ORANGE,
                     font=("Segoe UI", 7)).pack(side="right")

    def _build_content(self):
        content = tk.Frame(self, bg=PAGE_BG, padx=20, pady=14)
        content.pack(fill="both", expand=True)

        # File selection
        fframe = self._lframe(content, "Health Check Input File")
        fframe.pack(fill="x", pady=(0, 10))
        self._file_label = tk.Label(fframe, text="No file selected",
                                    bg=WHITE, fg="#888", font=("Segoe UI", 9), anchor="w")
        self._file_label.pack(side="left", fill="x", expand=True)
        self._size_badge = tk.Label(fframe, text="", bg=LIGHT_BLUE, fg=PRIMARY_BLUE,
                                    font=("Segoe UI", 8, "bold"), padx=8, pady=2)
        self._size_badge.pack(side="left", padx=(6, 8))
        tk.Button(fframe, text="Browse...", command=self._on_browse,
                  bg=PRIMARY_BLUE, fg=WHITE, font=("Segoe UI", 9, "bold"),
                  relief="flat", padx=14, cursor="hand2").pack(side="right")

        # Preview
        pframe = self._lframe(content, "File Preview  (first 60 lines)")
        pframe.pack(fill="both", expand=True, pady=(0, 10))
        self._preview_text = scrolledtext.ScrolledText(
            pframe, height=10, font=("Consolas", 8),
            bg="#f8fbff", fg="#333", relief="flat", state="disabled", wrap="none")
        self._preview_text.pack(fill="both", expand=True)

        # Progress
        pgframe = self._lframe(content, "Generation Progress")
        pgframe.pack(fill="x", pady=(0, 10))
        row = tk.Frame(pgframe, bg=WHITE)
        row.pack(fill="x")
        self._progress = ttk.Progressbar(row, variable=self._progress_var,
                                          maximum=100, mode="determinate", length=100)
        self._progress.pack(side="left", fill="x", expand=True, padx=(0, 10))
        tk.Label(row, textvariable=self._pct_var, bg=WHITE, fg=PRIMARY_BLUE,
                 font=("Segoe UI", 10, "bold"), width=6).pack(side="left")
        self._phase_label = tk.Label(pgframe, text="", bg=WHITE, fg="#555",
                                     font=("Segoe UI", 8), anchor="w")
        self._phase_label.pack(fill="x", pady=(4, 0))

        # Log
        lframe = self._lframe(content, "Generation Log")
        lframe.pack(fill="x", pady=(0, 10))
        self._log_text = scrolledtext.ScrolledText(
            lframe, height=6, font=("Consolas", 8),
            bg="#f0f7ff", fg="#333", relief="flat", state="disabled", wrap="word")
        self._log_text.pack(fill="both", expand=True)

        # Buttons
        btn_frame = tk.Frame(content, bg=PAGE_BG)
        btn_frame.pack(fill="x")

        self._run_btn = tk.Button(
            btn_frame, text="  Generate Full Analysis Report",
            command=self._on_generate, bg=PRIMARY_BLUE, fg=WHITE,
            font=("Segoe UI", 11, "bold"), relief="flat",
            padx=24, pady=10, cursor="hand2", state="disabled")
        self._run_btn.pack(side="left")

        self._cancel_btn = tk.Button(
            btn_frame, text="  Cancel",
            command=self._on_cancel, bg=DANGER_RED, fg=WHITE,
            font=("Segoe UI", 10), relief="flat",
            padx=14, pady=10, cursor="hand2", state="disabled")
        self._cancel_btn.pack(side="left", padx=(8, 0))

        self._open_btn = tk.Button(
            btn_frame, text="  Open in Edge",
            command=self._on_open_last, bg=OK_GREEN, fg=WHITE,
            font=("Segoe UI", 10), relief="flat",
            padx=16, pady=10, cursor="hand2", state="disabled")
        self._open_btn.pack(side="left", padx=(8, 0))

        # ── NEW v5: PowerPoint export button ─────────────────────────────
        pptx_label = "  Export to PowerPoint"
        pptx_state = "disabled"
        self._pptx_btn = tk.Button(
            btn_frame, text=pptx_label,
            command=self._on_pptx_generate, bg=PURPLE, fg=WHITE,
            font=("Segoe UI", 10), relief="flat",
            padx=16, pady=10, cursor="hand2", state=pptx_state)
        self._pptx_btn.pack(side="left", padx=(8, 0))

        if not PPTX_OK:
            self._pptx_btn.config(
                text="  python-pptx not installed",
                bg="#888", cursor="arrow")

    def _build_footer(self):
        footer = tk.Frame(self, bg=DARK_BLUE)
        footer.pack(fill="x", side="bottom")
        tk.Label(footer, textvariable=self._status_var, bg=DARK_BLUE, fg="#a8d0f0",
                 font=("Segoe UI", 8), anchor="w", padx=12).pack(
                     side="left", pady=4, fill="x", expand=True)
        tk.Label(footer, text="RDE LAC", bg=DARK_BLUE, fg=ACCENT_CYAN,
                 font=("Segoe UI", 8, "bold"), padx=16).pack(side="right", pady=4)

    def _lframe(self, parent, title):
        return tk.LabelFrame(parent, text=title, bg=WHITE, fg=DARK_BLUE,
                              font=("Segoe UI", 9, "bold"), padx=10, pady=8,
                              relief="flat", highlightbackground="#c8ddf0",
                              highlightthickness=1)

    # ── Browse / generate ─────────────────────────────────────────────────────

    def _on_browse(self):
        initial = str(WORKING_DIR / "HANA Health Check Reports")
        path = filedialog.askopenfilename(
            title="Select HANA Health Check Output File",
            initialdir=initial if os.path.isdir(initial) else str(WORKING_DIR),
            filetypes=[("Text files", "*.txt"), ("All files", "*.*")],
        )
        if not path:
            return
        self._selected_file = Path(path)
        size_kb = self._selected_file.stat().st_size / 1024
        approx_tokens = int(size_kb * 1024 / 4)
        self._file_label.config(text=self._selected_file.name, fg=DARK_BLUE)
        self._size_badge.config(
            text=f"{size_kb:.0f} KB  ~{approx_tokens:,} tokens",
            bg=WARN_ORANGE if size_kb > MAX_FILE_KB * 0.75 else LIGHT_BLUE,
            fg=DARK_BLUE if size_kb > MAX_FILE_KB * 0.75 else PRIMARY_BLUE,
        )
        self._load_preview()
        if size_kb > MAX_FILE_KB:
            self._run_btn.config(state="disabled")
            self._set_status(f"File too large ({size_kb:.0f} KB) — max {MAX_FILE_KB} KB.")
            messagebox.showwarning("File Too Large",
                f"{self._selected_file.name}\n{size_kb:.0f} KB (~{approx_tokens:,} tokens)\n\n"
                f"Exceeds {MAX_FILE_KB} KB limit. Please select a smaller file.")
        else:
            self._run_btn.config(state="normal")
            self._set_status(f"Ready  |  {self._selected_file.name}  ({size_kb:.0f} KB)")

    def _load_preview(self):
        try:
            lines = load_text(self._selected_file).splitlines()[:60]
            preview = "\n".join(lines)
        except Exception as e:
            preview = f"Could not read file: {e}"
        self._preview_text.config(state="normal")
        self._preview_text.delete("1.0", "end")
        self._preview_text.insert("end", preview)
        self._preview_text.config(state="disabled")

    def _on_generate(self):
        if not self._selected_file or not self._selected_file.exists():
            messagebox.showerror("No File", "Please select a valid health check file.")
            return
        self._cancel_event.clear()
        self._run_btn.config(state="disabled")
        self._cancel_btn.config(state="normal")
        self._open_btn.config(state="disabled")
        self._pptx_btn.config(state="disabled")  # disable during generation
        self._log_clear()
        self._set_progress(0, "Starting...")
        self._set_status("Generating report via streaming API...")
        threading.Thread(target=self._run_generation, daemon=True).start()

    def _on_cancel(self):
        self._cancel_event.set()
        self._cancel_btn.config(state="disabled")
        self._set_status("Cancelling...")

    def _on_open_last(self):
        if self._last_output_file and self._last_output_file.exists():
            self._open_in_edge(self._last_output_file)
        else:
            messagebox.showinfo("No Report", "No report has been generated yet.")

    # ── HTML generation thread ────────────────────────────────────────────────

    def _run_generation(self):
        try:
            health_check_text = load_text(self._selected_file)
            template          = load_text(TEMPLATE_FILE)
            release_info      = load_text(RELEASE_FILE) if RELEASE_FILE.exists() else ""

            sid = extract_sid(health_check_text)
            if sid == "UNKNOWN":
                sid = self._selected_file.stem.split("_")[0].upper()

            ts          = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = OUTPUT_DIR / f"{sid}_HANA_Analysis_Report_{ts}.html"
            OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

            prompt = build_prompt(template, release_info, health_check_text)

            self.after(0, self._log, f"File    : {self._selected_file.name}\n")
            self.after(0, self._log, f"Prompt  : {len(prompt):,} chars (~{len(prompt)//4:,} tokens)\n")
            self.after(0, self._log, f"Model   : {self._model}\n")
            self.after(0, self._log, f"Output  : {output_file.name}\n")
            self.after(0, self._log, f"Logo    : {'embedded' if self._logo_b64 else 'not found'}\n")
            self.after(0, self._set_progress, 5, "Connecting to proxy...")

            client = anthropic.Anthropic(base_url=self._base_url, api_key=self._token)

            self.after(0, self._log, "Streaming response...\n")
            self.after(0, self._set_progress, 10, "Claude is generating the HTML report...")

            html_chunks   = []
            output_tokens = 0
            start_time    = datetime.now()

            with client.messages.stream(
                model=self._model,
                max_tokens=MAX_OUTPUT_TOKENS,
                messages=[{"role": "user", "content": prompt}],
            ) as stream:
                for text_chunk in stream.text_stream:
                    if self._cancel_event.is_set():
                        break
                    html_chunks.append(text_chunk)
                    output_tokens += len(text_chunk) // 4
                    pct = min(10 + int(output_tokens / EXPECTED_OUT_TOKENS * 85), 95)
                    elapsed = int((datetime.now() - start_time).total_seconds())
                    self.after(0, self._set_progress, pct,
                               f"Generating...  ~{output_tokens:,} tokens  "
                               f"[{elapsed // 60}m {elapsed % 60:02d}s]")

            if self._cancel_event.is_set():
                self.after(0, self._on_generation_error, "Cancelled by user.")
                return

            elapsed     = int((datetime.now() - start_time).total_seconds())
            self.after(0, self._log,
                       f"Stream complete in {elapsed // 60}m {elapsed % 60:02d}s\n")
            self.after(0, self._set_progress, 96, "Post-processing HTML...")

            raw_html = "".join(html_chunks)
            raw_html = _strip_code_fences(raw_html)

            if not raw_html.startswith("<"):
                raise ValueError(
                    "Response is not valid HTML.\n\nFirst 300 chars:\n" + raw_html[:300]
                )

            self.after(0, self._log, "Applying enhancements: CSS / logo / footer...\n")
            full_html = post_process_html(raw_html, self._logo_b64)

            self.after(0, self._set_progress, 98, "Saving...")
            output_file.write_text(full_html, encoding="utf-8")
            size_kb = output_file.stat().st_size // 1024
            self.after(0, self._log, f"Saved: {output_file.name}  ({size_kb} KB)\n")

            self._last_output_file = output_file
            self.after(0, self._on_generation_complete, output_file)

        except anthropic.APIStatusError as e:
            self.after(0, self._on_generation_error,
                       f"API error {e.status_code}: {e.message}")
        except Exception as e:
            self.after(0, self._on_generation_error, str(e))

    def _on_generation_complete(self, output_file: Path):
        self._set_progress(100, "Complete!")
        self._run_btn.config(state="normal")
        self._cancel_btn.config(state="disabled")
        self._open_btn.config(state="normal")
        # Enable PPTX button only if python-pptx is available
        if PPTX_OK:
            self._pptx_btn.config(state="normal")
        self._log(f"\n  Report ready: {output_file}\n")
        self._set_status(f"Done — {output_file.name}")
        self._open_in_edge(output_file)

    def _on_generation_error(self, msg: str):
        self._set_progress(0, "")
        self._run_btn.config(state="normal")
        self._cancel_btn.config(state="disabled")
        self._log(f"\n  Error: {msg}\n")
        self._set_status("Error — see log.")
        messagebox.showerror("Generation Failed", msg)

    # ── PowerPoint export (v5 addition) ──────────────────────────────────────

    def _on_pptx_generate(self):
        if not PPTX_OK:
            messagebox.showerror(
                "python-pptx not installed",
                "Install it with:\n\n    pip install python-pptx\n\nthen restart the app.")
            return
        if not self._last_output_file or not self._last_output_file.exists():
            messagebox.showinfo("No Report", "Generate an HTML report first.")
            return
        self._pptx_btn.config(state="disabled")
        self._set_status("Building PowerPoint presentation...")
        self._log("\nBuilding PowerPoint...\n")
        threading.Thread(
            target=self._run_pptx_generation,
            args=(self._last_output_file,),
            daemon=True,
        ).start()

    def _run_pptx_generation(self, html_path: Path):
        try:
            self.after(0, self._log, f"Parsing: {html_path.name}\n")
            pptx_path = build_pptx_from_html(html_path, LOGO_PATH)
            size_kb = pptx_path.stat().st_size // 1024
            self.after(0, self._log,
                       f"Saved: {pptx_path.name}  ({size_kb} KB)\n")
            self.after(0, self._on_pptx_complete, pptx_path)
        except Exception as e:
            self.after(0, self._on_pptx_error, str(e))

    def _on_pptx_complete(self, pptx_path: Path):
        self._pptx_btn.config(state="normal")
        self._set_status(f"PowerPoint ready — {pptx_path.name}")
        self._log(f"  PowerPoint ready: {pptx_path}\n")
        import subprocess
        try:
            os.startfile(str(pptx_path))
        except Exception:
            try:
                subprocess.Popen(["explorer", "/select,", str(pptx_path)])
            except Exception:
                pass
        messagebox.showinfo(
            "PowerPoint Ready",
            f"Saved to:\n{pptx_path}\n\nOpening file...")

    def _on_pptx_error(self, msg: str):
        self._pptx_btn.config(state="normal")
        self._set_status("PowerPoint generation failed — see log.")
        self._log(f"\n  PPTX Error: {msg}\n")
        messagebox.showerror("PowerPoint Failed", msg)

    # ── Utilities ─────────────────────────────────────────────────────────────

    def _open_in_edge(self, path: Path):
        import subprocess
        try:
            subprocess.Popen([EDGE_PATH, str(path)])
            self._set_status(f"Opened in Edge: {path.name}")
        except Exception as e:
            messagebox.showerror("Could not open Edge", str(e))

    def _set_progress(self, pct: int, phase: str):
        self._progress_var.set(pct)
        self._pct_var.set(f"{pct}%" if pct > 0 else "")
        self._phase_label.config(text=phase)

    def _log(self, msg: str):
        self._log_text.config(state="normal")
        self._log_text.insert("end", msg)
        self._log_text.see("end")
        self._log_text.config(state="disabled")

    def _log_clear(self):
        self._log_text.config(state="normal")
        self._log_text.delete("1.0", "end")
        self._log_text.config(state="disabled")

    def _set_status(self, msg: str):
        self._status_var.set(msg)


if __name__ == "__main__":
    app = AnalysisApp()
    app.mainloop()
