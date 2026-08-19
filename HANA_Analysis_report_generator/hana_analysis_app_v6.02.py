"""
SAP HANA Full Analysis Report Generator v6.02
- All v6.01 features retained
- Minicheck discovery: scans health check input and dynamically injects ALL
  minicheck sections found (not just the 7 hardcoded in the template)
- PPTX export limited to 4 key sections:
    Executive Overview, Environment Summary/Overview,
    Additional Risks, Consolidated Action Plans
- Style matches SAP 360 Caretaker Engagement Review template:
    Font  : "72" / "72 Brand Medium" (SAP brand font)
    Colors: #0070F2 (SAP Blue), #C35500 (Mango), #131E29 (text)
    Layout: rounded-rect section pills, separator lines, footer note
"""

import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
import threading
import os
import re
import json
import base64
from datetime import datetime
from pathlib import Path
import anthropic

# ── PPTX availability ─────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

# ── Config ────────────────────────────────────────────────────────────────────
WORKING_DIR   = Path(r"C:\Users\I522148\OneDrive - SAP SE\SWAT\AI\Claude\working directory")
TEMPLATE_FILE = WORKING_DIR / "HANA_HealthCheck_Prompt_Template.md"
RELEASE_FILE  = WORKING_DIR / "HANA_latest_release.txt"
OUTPUT_DIR    = WORKING_DIR / "Results"
LOGO_PATH     = WORKING_DIR / "SAP_LOGO.png"
EDGE_PATH     = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
SETTINGS_FILE = Path(r"C:\Users\I522148\.claude\settings.json")

MAX_FILE_KB         = 2500
EXPECTED_OUT_TOKENS = 40000
MAX_OUTPUT_TOKENS   = 64000
APP_VERSION         = "v6.02"

# ── UI Colors ─────────────────────────────────────────────────────────────────
PRIMARY_BLUE = "#0057a8"
DARK_BLUE    = "#003366"
ACCENT_CYAN  = "#00b4e6"
PAGE_BG      = "#e8f4fd"
WHITE        = "#ffffff"
LIGHT_BLUE   = "#f0f7ff"
OK_GREEN     = "#5cb85c"
WARN_ORANGE  = "#f0ad4e"
DANGER_RED   = "#d9534f"
PURPLE       = "#6f42c1"

# ── Enhanced CSS (injected into every generated report) ───────────────────────
ENHANCED_CSS = """
/* ── Enhanced v4 additions ─── */
.consequence{background:#fff3f3;border-left:4px solid #d9534f;padding:8px 14px;
  margin:2px 0 6px 0;border-radius:4px;font-size:0.85rem;color:#555}
.esc-banner{background:#d9534f;color:#fff;padding:12px 20px;border-radius:6px;
  font-weight:700;font-size:1rem;margin-bottom:16px;display:flex;align-items:center;gap:10px}
.esc-banner.warn{background:#e65c00}
.exec-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:16px;margin:16px 0}
.exec-card{background:#fff;border-radius:8px;padding:16px;box-shadow:0 2px 6px rgba(0,0,0,.1)}
.exec-card.critical{border-top:4px solid #d9534f}
.exec-card.warning{border-top:4px solid #f0ad4e}
.exec-card.monthly{border-top:4px solid #0057a8}
.exec-card.ok{border-top:4px solid #5cb85c}
.exec-card h4{margin:0 0 6px 0;font-size:0.95rem;color:#003366}
.exec-card .biz-impact{font-size:0.82rem;color:#d9534f;font-weight:600;margin:4px 0}
.exec-card .biz-action{font-size:0.82rem;color:#555;margin:4px 0}
.health-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(160px,1fr));gap:10px;margin:12px 0}
.health-item{background:#fff;border-radius:6px;padding:10px 14px;display:flex;
  align-items:center;gap:10px;box-shadow:0 1px 4px rgba(0,0,0,.08)}
.health-item .hi-icon{font-size:1.4rem}
.health-item .hi-label{font-size:0.8rem;color:#666}
.health-item .hi-val{font-size:0.85rem;font-weight:700}
.ap-section{margin:12px 0}
.ap-section h3{font-size:0.9rem;color:#003366;margin:0 0 8px 0;
  padding:6px 0;border-bottom:2px solid #e8f0f8}
.ap-card{background:#fff;border-radius:8px;padding:14px 18px;margin:8px 0;
  box-shadow:0 2px 5px rgba(0,0,0,.08);border-left:5px solid #0057a8}
.ap-card.p1{border-left-color:#d9534f}
.ap-card.p2{border-left-color:#f0ad4e}
.ap-card.p3{border-left-color:#5cb85c}
.ap-card h4{margin:0 0 6px 0;font-size:0.95rem;color:#003366}
.ap-card p{margin:0;font-size:0.85rem;color:#555}
.ap-tags{display:flex;flex-wrap:wrap;gap:6px;margin-top:8px}
.ap-tag{font-size:0.72rem;font-weight:600;padding:2px 8px;border-radius:10px;
  background:#e8f4fd;color:#003366}
.ap-tag.urgent{background:#fde8e8;color:#d9534f}
.ap-tag.owner{background:#e8f8e8;color:#3a7a3a}
.ap-tag.effort{background:#fff8e1;color:#856404}
.ap-tag.prevents{background:#f0e8ff;color:#5a3a8a}
.kpi-critical{border-top:4px solid #d9534f !important}
.kpi-critical .kpi-val{color:#d9534f !important}
/* ── End enhanced v4 ─── */
"""


# =============================================================================
# HTML parsing helpers
# =============================================================================

def _strip_tags(s):
    s = re.sub(r"<[^>]+>", " ", s)
    for ent, ch in [("&amp;", "&"), ("&lt;", "<"), ("&gt;", ">"),
                    ("&nbsp;", " "), ("&#8203;", ""), ("&apos;", "'")]:
        s = s.replace(ent, ch)
    return " ".join(s.split()).strip()


def _parse_report_meta(html):
    meta = {
        "sid": "UNKNOWN", "tenant": "", "host": "", "instance": "",
        "version": "", "system_type": "", "generated": "", "script": "",
    }
    footer_m = re.search(r"<footer[^>]*>(.*?)</footer>", html, re.DOTALL)
    if footer_m:
        ft = _strip_tags(footer_m.group(1))
        for pat, key in [
            (r"SID:\s*([^|]+)",       "sid"),
            (r"Tenant:\s*([^|]+)",    "tenant"),
            (r"Generated:\s*([^|]+)", "generated"),
            (r"Script:\s*([^|]+)",    "script"),
        ]:
            m = re.search(pat, ft, re.IGNORECASE)
            if m:
                meta[key] = m.group(1).strip()

    hdr_m = re.search(r'class=["\'][^"\']*header-meta[^"\']*["\'][^>]*>(.*?)</div>',
                      html, re.DOTALL)
    if hdr_m:
        hdr_text = _strip_tags(hdr_m.group(1))
        for pat, key in [
            (r"SID[:\s]+(\S+)",              "sid"),
            (r"Tenant.*?:\s*(\S+)",          "tenant"),
            (r"Host[:\s]+([^\s|·]+)",        "host"),
            (r"Instance[:\s]+([^\s|·]+)",    "instance"),
            (r"(\d+\.\d+\.\d+\.\d+)",       "version"),
            (r"(MDC[^|·<\n]*)",             "system_type"),
        ]:
            m = re.search(pat, hdr_text, re.IGNORECASE)
            if m:
                meta[key] = m.group(1).strip()

    return meta


def _parse_sections(html):
    sections = []
    starts   = []

    # Build nav id→title map — handles both scrollTo('id') buttons and <a href="#id"> anchors
    nav_map = {}
    for nm in re.finditer(r"scrollTo\([\"'](\w+)[\"']\)[^>]*>(.*?)</button>", html, re.DOTALL):
        raw = _strip_tags(nm.group(2)).strip()
        nav_map[nm.group(1)] = re.sub(r'^[^\w]+', '', raw).strip()
    for nm in re.finditer(r'<a\b[^>]+href=["\']#(\w+)["\'][^>]*>(.*?)</a>', html, re.DOTALL):
        sid = nm.group(1)
        if sid not in nav_map:
            raw = _strip_tags(nm.group(2)).strip()
            clean = re.sub(r'^[^\w]+', '', raw).strip()
            if clean:
                nav_map[sid] = clean

    for m in re.finditer(r'<div\b[^>]*>', html, re.IGNORECASE):
        tag   = m.group(0)
        id_m  = re.search(r'\bid\s*=\s*["\']([^"\']+)["\']',    tag, re.IGNORECASE)
        cls_m = re.search(r'\bclass\s*=\s*["\']([^"\']*)["\']', tag, re.IGNORECASE)
        if not id_m or not cls_m:
            continue
        sec_id  = id_m.group(1).strip()
        classes = cls_m.group(1).split()
        if (any("section" in c for c in classes)
                and "section-body" not in classes
                and not sec_id.startswith("body-")):
            starts.append((m.start(), sec_id))

    for i, (pos, sid) in enumerate(starts):
        end   = starts[i + 1][0] if i + 1 < len(starts) else len(html)
        chunk = html[pos:end]

        # icon
        icon_m = re.search(
            r'class=["\'][^"\']*section-icon(?![-\w])[^"\']*["\'][^>]*>(.*?)</span>',
            chunk, re.DOTALL)
        icon = _strip_tags(icon_m.group(1)) if icon_m else ""

        # title — nav_map is authoritative (consistent across both HTML formats)
        title = nav_map.get(sid, "")
        if not title:
            # old format: <span class="section-title">text</span>
            title_m = re.search(
                r'class=["\'][^"\']*section-title(?![-\w])[^"\']*["\'][^>]*>(.*?)</span>',
                chunk, re.DOTALL)
            if title_m:
                t = _strip_tags(title_m.group(1)).strip()
                title = t if len(t) >= 3 else ""
        if not title:
            # new format: <div class="section-title"><span …icon…</span> Title text</div>
            div_m = re.search(
                r'class=["\']section-title["\'][^>]*>(.*?)</div>',
                chunk, re.DOTALL)
            if div_m:
                raw = re.sub(r'^[^\w]+', '', _strip_tags(div_m.group(1))).strip()
                title = raw
        if not title:
            title = sid

        # badge
        badge_m = re.search(
            r'class=["\'][^"\']*pill[^"\']*pill-([a-z]+)[^"\']*["\'][^>]*>(.*?)</span>',
            chunk, re.DOTALL)
        if not badge_m:
            badge_m = re.search(
                r'class=["\'][^"\']*pill-([a-z]+)[^"\']*["\'][^>]*>(.*?)</span>',
                chunk, re.DOTALL)
        badge_color = badge_m.group(1).strip() if badge_m else "blue"
        badge_text  = _strip_tags(badge_m.group(2)) if badge_m else ""

        # body
        body_m = re.search(
            r'class=["\'][^"\']*section-body[^"\']*["\'][^>]*>(.*)',
            chunk, re.DOTALL)
        body = body_m.group(1) if body_m else chunk

        # h4s and bullets
        h4s     = [_strip_tags(h) for h in re.findall(r"<h4[^>]*>(.*?)</h4>",    body, re.DOTALL)]
        h4s     = [h for h in h4s if h][:5]
        bullets = [_strip_tags(li) for li in re.findall(r"<li[^>]*>(.*?)</li>",  body, re.DOTALL)]
        bullets = [b for b in bullets if len(b) > 5][:8]

        # table-row fallback
        if not h4s and not bullets:
            for row_m in re.finditer(r"<tr\b[^>]*>(.*?)</tr>", body, re.DOTALL):
                cells = [_strip_tags(c)
                         for c in re.findall(r"<td[^>]*>(.*?)</td>", row_m.group(1), re.DOTALL)]
                cells = [c for c in cells if c]
                if len(cells) >= 2:
                    bullets.append(cells[0] + "  \u2014  " + cells[-1])
            bullets = bullets[:10]

        # env-item key/value pairs (Environment Summary section)
        env_items = []
        for em in re.finditer(
            r'class=["\']env-item["\'][^>]*>.*?'
            r'class=["\']env-item-label["\'][^>]*>(.*?)</div>.*?'
            r'class=["\']env-item-val["\'][^>]*>(.*?)</div>',
            body, re.DOTALL
        ):
            lbl = _strip_tags(em.group(1))
            val = _strip_tags(em.group(2))
            if lbl:
                env_items.append((lbl, val))
        # Append bar-label utilisation rows (Memory, Disk, etc.)
        for bm in re.finditer(
            r'<h4[^>]*>(.*?)</h4>.*?class=["\']bar-label["\'][^>]*>(.*?)</span>',
            body, re.DOTALL
        ):
            h4txt  = _strip_tags(bm.group(1))
            bartxt = _strip_tags(bm.group(2))
            if h4txt and bartxt:
                env_items.append((h4txt, bartxt))

        # action plan tiers
        ap_tiers = {"p1": [], "p2": [], "p3": []}
        for tier in ("p1", "p2", "p3"):
            for card_m in re.finditer(
                r'<div\b[^>]*class=["\'][^"\']*\bap-card\b[^"\']*\b' + tier
                + r'\b[^"\']*["\'][^>]*>.*?<h4[^>]*>(.*?)</h4>.*?<p[^>]*>(.*?)</p>',
                body, re.DOTALL
            ):
                ap_tiers[tier].append(
                    _strip_tags(card_m.group(1)) + " \u2014 " + _strip_tags(card_m.group(2))
                )

        sections.append({
            "id": sid, "icon": icon, "title": title,
            "badge_color": badge_color, "badge_text": badge_text,
            "h4s": h4s, "bullets": bullets,
            "ap_tiers": ap_tiers, "env_items": env_items,
        })

    return sections


# =============================================================================
# PPTX v6 — SAP 360 style
# =============================================================================

# Only these 4 sections go into the PPTX
_PPTX_INCLUDE = [
    "executive overview",
    "executive summary",
    "environment",          # covers "Environment Overview", "Environment Summary", etc.
    "additional risk",
    "consolidated action",
    "action plan",
]

def _pptx_include(sec):
    t = sec["title"].lower()
    return any(kw in t for kw in _PPTX_INCLUDE)


# Slide dimensions 16:9
_SW = Inches(13.33) if PPTX_OK else None
_SH = Inches(7.5)   if PPTX_OK else None

# SAP 2023 brand palette
_C = {
    "white":    RGBColor(0xFF, 0xFF, 0xFF) if PPTX_OK else None,
    "blue":     RGBColor(0x00, 0x70, 0xF2) if PPTX_OK else None,   # SAP Blue 7
    "mango":    RGBColor(0xC3, 0x55, 0x00) if PPTX_OK else None,   # SAP Mango 7
    "text":     RGBColor(0x13, 0x1E, 0x29) if PPTX_OK else None,   # body
    "text2":    RGBColor(0x1D, 0x2D, 0x3E) if PPTX_OK else None,   # secondary
    "p1":       RGBColor(0xEE, 0x39, 0x39) if PPTX_OK else None,
    "p2":       RGBColor(0xE7, 0x65, 0x00) if PPTX_OK else None,
    "p3":       RGBColor(0x36, 0xA4, 0x1D) if PPTX_OK else None,
    "line":     RGBColor(0xC8, 0xD6, 0xE8) if PPTX_OK else None,
    "row_alt":  RGBColor(0xEF, 0xF4, 0xFB) if PPTX_OK else None,
    "p1_bg":    RGBColor(0xFF, 0xF0, 0xF0) if PPTX_OK else None,
    "p2_bg":    RGBColor(0xFF, 0xFB, 0xF0) if PPTX_OK else None,
    "p3_bg":    RGBColor(0xF0, 0xFB, 0xF0) if PPTX_OK else None,
    "dark_blue":RGBColor(0x00, 0x2A, 0x86) if PPTX_OK else None,
}

_FONT     = "72"
_FONT_MED = "72 Brand Medium"


# ── Low-level shape helpers ───────────────────────────────────────────────────

def _rect(sl, x, y, w, h, fill, line=None):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s


def _rrect(sl, x, y, w, h, fill, line=None):
    """Rounded rectangle (shape type 5 = ROUNDED_RECTANGLE)."""
    s = sl.shapes.add_shape(5, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s


def _hline(sl, y, color=None, thickness=1):
    color = color or _C["line"]
    s = sl.shapes.add_shape(1, 0, y, _SW, Pt(thickness))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()


def _txt(sl, text, x, y, w, h, sz=11, bold=False, color=None,
         align=PP_ALIGN.LEFT, italic=False, font=None):
    color = color or _C["text"]
    font  = font  or _FONT
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return tb


def _badge_col(css_name):
    name = (css_name or "blue").split()[0]
    return {"red": _C["p1"], "orange": _C["p2"],
            "green": _C["p3"], "blue": _C["blue"]}.get(name, _C["blue"])


# ── Shared slide header ───────────────────────────────────────────────────────

def _section_header(sl, title, badge_text, badge_rgb):
    """Rounded-rect section pill + optional badge chip + separator line."""
    pill_w = max(Inches(3.8), min(Inches(7.5), Inches(len(title) * 0.22 + 0.9)))
    s = _rrect(sl, Inches(0.45), Inches(0.45), pill_w, Inches(0.46), badge_rgb)
    tf = s.text_frame; tf.word_wrap = False
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r  = p.add_run()
    r.text = "  " + title; r.font.name = _FONT_MED
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = _C["white"]

    if badge_text:
        bw = Inches(2.2); bx = _SW - bw - Inches(0.45)
        bs = _rrect(sl, bx, Inches(0.45), bw, Inches(0.46), badge_rgb)
        tf2 = bs.text_frame; tf2.word_wrap = False
        p2  = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        r2  = p2.add_run()
        r2.text = badge_text; r2.font.name = _FONT; r2.font.size = Pt(10)
        r2.font.bold = True; r2.font.color.rgb = _C["white"]

    _hline(sl, Inches(1.07), _C["line"])

    # Footer note bottom-right (matching reference template)
    _txt(sl, "SAP HANA Health Check Report",
         Inches(7.5), _SH - Inches(0.33), Inches(5.5), Inches(0.26),
         sz=9, color=_C["text2"], align=PP_ALIGN.RIGHT, font=_FONT)


# ── Title slide ───────────────────────────────────────────────────────────────

def _slide_title(prs, meta, logo_path):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, _SW, _SH, _C["white"])
    _rect(sl, 0, 0, Inches(0.1), _SH, _C["blue"])                      # left bar
    _rect(sl, Inches(0.1), 0, _SW - Inches(0.1), Pt(3), _C["blue"])    # top line
    _rect(sl, 0, _SH - Pt(3), _SW, Pt(3), _C["blue"])                  # bottom line

    if logo_path and logo_path.exists():
        try:
            sl.shapes.add_picture(str(logo_path),
                                  _SW - Inches(1.85), Inches(0.28), height=Inches(0.5))
        except Exception:
            pass

    _txt(sl, "SAP HANA Database  \u00B7  Health Check Report",
         Inches(0.35), Inches(0.82), Inches(10.5), Inches(0.55),
         sz=15, color=_C["blue"], font=_FONT_MED)

    sid_line = meta.get("sid", "UNKNOWN")
    tenant   = meta.get("tenant", "")
    if tenant:
        sid_line += "   \u00B7   " + tenant
    _txt(sl, sid_line,
         Inches(0.35), Inches(1.42), Inches(12.6), Inches(0.9),
         sz=36, bold=True, color=_C["text"], font=_FONT_MED)

    _hline(sl, Inches(2.52), _C["line"])

    meta_rows = [(l, v) for l, v in [
        ("Host",         meta.get("host",        "")),
        ("Instance",     meta.get("instance",    "")),
        ("HANA Version", meta.get("version",     "")),
        ("System Type",  meta.get("system_type", "")),
        ("Generated",    meta.get("generated",   "")),
        ("Script",       meta.get("script",      "")),
    ] if v]

    for ci, col_items in enumerate([meta_rows[:3], meta_rows[3:]]):
        cx = Inches(0.35) if ci == 0 else Inches(6.9)
        for ri, (lbl, val) in enumerate(col_items):
            ry = Inches(2.72 + ri * 0.52)
            _txt(sl, lbl + ":", cx, ry, Inches(2.3), Inches(0.44),
                 sz=10, bold=True, color=_C["text2"], font=_FONT)
            _txt(sl, val, cx + Inches(2.35), ry, Inches(4.1), Inches(0.44),
                 sz=10, color=_C["text"], font=_FONT)

    _txt(sl, "Generated by SAP HANA Analysis Report Generator " + APP_VERSION,
         Inches(0.35), _SH - Inches(0.38), Inches(9), Inches(0.3),
         sz=8, color=RGBColor(0xA0, 0xB4, 0xCC), font=_FONT)


# ── Content renderers ─────────────────────────────────────────────────────────

def _render_exec_overview(sl, sec):
    """Executive Overview: finding cards + key findings list."""
    h4s     = sec["h4s"][:5]
    bullets = sec["bullets"][:8]
    y = Inches(1.22)

    if h4s:
        n     = len(h4s)
        kpi_w = min(Inches(2.55), (_SW - Inches(1.0)) / max(1, n) - Inches(0.1))
        kpi_h = Inches(1.15)
        for i, kpi in enumerate(h4s):
            kx = Inches(0.45) + i * (kpi_w + Inches(0.1))
            _rect(sl, kx, y, kpi_w, kpi_h, _C["row_alt"])
            _rect(sl, kx, y, kpi_w, Pt(4), _C["blue"])
            if ":" in kpi:
                val_part, lbl_part = kpi.split(":", 1)
            else:
                val_part, lbl_part = kpi, ""
            _txt(sl, val_part.strip(),
                 kx + Inches(0.12), y + Inches(0.12), kpi_w - Inches(0.2), Inches(0.52),
                 sz=12, bold=True, color=_C["text"], font=_FONT_MED)
            if lbl_part.strip():
                _txt(sl, lbl_part.strip(),
                     kx + Inches(0.12), y + Inches(0.64), kpi_w - Inches(0.2), Inches(0.44),
                     sz=10.5, color=_C["text2"], font=_FONT)
        y += kpi_h + Inches(0.22)

    if bullets:
        _hline(sl, y, _C["line"])
        y += Inches(0.18)
        _txt(sl, "Key Findings",
             Inches(0.45), y, Inches(5), Inches(0.38),
             sz=13, bold=True, color=_C["blue"], font=_FONT_MED)
        y += Inches(0.44)
        for item in bullets:
            if y > _SH - Inches(0.55):
                break
            _txt(sl, "\u25B6  " + item[:110],
                 Inches(0.55), y, _SW - Inches(1.1), Inches(0.38),
                 sz=11.5, color=_C["text"], font=_FONT)
            y += Inches(0.42)


def _render_environment(sl, sec):
    """Environment Summary: property/value table from env-item divs."""
    items = sec.get("env_items", [])

    if not items:
        for h in sec["h4s"]:
            if ":" in h:
                lbl, val = h.split(":", 1)
                items.append((lbl.strip(), val.strip()))
            else:
                items.append((h.strip(), ""))
        for b in sec["bullets"]:
            if ":" in b:
                lbl, val = b.split(":", 1)
                items.append((lbl.strip(), val.strip()))

    if not items:
        _txt(sl, "No environment data found. Refer to the HTML report.",
             Inches(0.45), Inches(1.3), Inches(12), Inches(0.45),
             sz=12, color=_C["text2"])
        return

    n     = min(len(items), 14)
    row_h = min(Inches(0.42), Inches(5.7) / (n + 1))

    _rect(sl, Inches(0.45), Inches(1.22), Inches(4.5),          row_h, _C["blue"])
    _rect(sl, Inches(5.05), Inches(1.22), _SW - Inches(5.5),    row_h, _C["blue"])
    _txt(sl, "Property",
         Inches(0.58), Inches(1.22), Inches(4.3), row_h,
         sz=11.5, bold=True, color=_C["white"], font=_FONT_MED)
    _txt(sl, "Value",
         Inches(5.18), Inches(1.22), _SW - Inches(5.6), row_h,
         sz=11.5, bold=True, color=_C["white"], font=_FONT_MED)

    for ri, (lbl, val) in enumerate(items[:n]):
        ry = Inches(1.22) + (ri + 1) * row_h
        bg = _C["row_alt"] if ri % 2 == 0 else _C["white"]
        _rect(sl, Inches(0.45), ry, Inches(4.5),       row_h, bg)
        _rect(sl, Inches(5.05), ry, _SW - Inches(5.5), row_h, bg)
        _hline(sl, ry + row_h - Pt(0.5), _C["line"])
        _txt(sl, lbl,
             Inches(0.58), ry + Inches(0.04), Inches(4.3), row_h - Inches(0.06),
             sz=11, bold=True, color=_C["text"], font=_FONT)
        _txt(sl, val[:120] if val else "\u2014",
             Inches(5.18), ry + Inches(0.04), _SW - Inches(5.6), row_h - Inches(0.06),
             sz=11, color=_C["text"], font=_FONT)


def _render_risks(sl, sec):
    """Additional Risks: risk area / recommended action table."""
    rows = []
    for b in sec["bullets"]:
        split_done = False
        # Try various separator characters (em dash, en dash, arrow, etc.)
        for sep in ["\u2014", "\u2192", "\u2013", "  -  "]:
            if sep in b:
                p1, p2 = b.split(sep, 1)
                rows.append((p1.strip(), p2.strip()))
                split_done = True
                break
        if not split_done:
            m = re.match(r'^(.+?)\s{2,}(.+)$', b.strip())
            if m:
                rows.append((m.group(1).strip(), m.group(2).strip()))
            else:
                rows.append((b.strip(), ""))

    if not rows:
        _txt(sl, "No additional risks identified.",
             Inches(0.45), Inches(1.3), Inches(12), Inches(0.45),
             sz=12, color=_C["text2"])
        return

    n     = min(len(rows), 10)
    row_h = Inches(0.55)

    _rect(sl, Inches(0.45), Inches(1.22), Inches(3.5),          row_h, _C["mango"])
    _rect(sl, Inches(4.05), Inches(1.22), _SW - Inches(4.5),    row_h, _C["mango"])
    _txt(sl, "Risk Area",
         Inches(0.58), Inches(1.22), Inches(3.3), row_h,
         sz=11.5, bold=True, color=_C["white"], font=_FONT_MED)
    _txt(sl, "Recommended Action",
         Inches(4.18), Inches(1.22), _SW - Inches(4.6), row_h,
         sz=11.5, bold=True, color=_C["white"], font=_FONT_MED)

    for ri, (risk, action) in enumerate(rows[:n]):
        ry = Inches(1.22) + (ri + 1) * row_h
        if ry + row_h > _SH - Inches(0.45):
            break
        bg = _C["row_alt"] if ri % 2 == 0 else _C["white"]
        _rect(sl, Inches(0.45), ry, Inches(3.5),       row_h, bg)
        _rect(sl, Inches(4.05), ry, _SW - Inches(4.5), row_h, bg)
        _hline(sl, ry + row_h - Pt(0.5), _C["line"])
        _txt(sl, risk[:60],
             Inches(0.58), ry + Inches(0.05), Inches(3.3), row_h - Inches(0.08),
             sz=11, bold=True, color=_C["text"], font=_FONT)
        _txt(sl, action[:150],
             Inches(4.18), ry + Inches(0.05), _SW - Inches(4.6), row_h - Inches(0.08),
             sz=10.5, color=_C["text"], font=_FONT)


def _render_action_plan(sl, ap_tiers):
    """Consolidated Action Plans: tier header strips + 2-per-row item cards."""
    TIERS = [
        ("p1", "Priority 1  \u2014  Immediate (\u226448h)", _C["p1"], _C["p1_bg"]),
        ("p2", "Priority 2  \u2014  This Week",              _C["p2"], _C["p2_bg"]),
        ("p3", "Priority 3  \u2014  This Month",             _C["p3"], _C["p3_bg"]),
    ]
    y = Inches(1.22)

    for tier_key, tier_label, col, bg in TIERS:
        items = ap_tiers.get(tier_key, [])
        if not items:
            continue

        _rect(sl, Inches(0.45), y, _SW - Inches(0.9), Inches(0.38), col)
        _txt(sl, "  " + tier_label,
             Inches(0.45), y + Inches(0.03), _SW - Inches(0.9), Inches(0.32),
             sz=12, bold=True, color=_C["white"], font=_FONT_MED)
        y += Inches(0.38)

        cw      = (_SW - Inches(1.12)) / 2 - Inches(0.06)
        ch      = Inches(0.76)
        n_items = min(len(items), 4)

        for idx, item in enumerate(items[:4]):
            cx = Inches(0.45) + (idx % 2) * (cw + Inches(0.12))
            cy = y + (idx // 2) * (ch + Inches(0.07))
            _rect(sl, cx, cy, cw, ch, bg, col)
            _rect(sl, cx, cy, Inches(0.07), ch, col)
            parts   = item.split(" \u2014 ", 1)
            title_t = parts[0][:70]
            desc_t  = parts[1][:130] if len(parts) > 1 else ""
            _txt(sl, title_t,
                 cx + Inches(0.16), cy + Inches(0.07), cw - Inches(0.24), Inches(0.32),
                 sz=11, bold=True, color=_C["text"], font=_FONT_MED)
            if desc_t:
                _txt(sl, desc_t,
                     cx + Inches(0.16), cy + Inches(0.42), cw - Inches(0.24), Inches(0.28),
                     sz=10, color=_C["text2"], font=_FONT)

        row_count = (n_items + 1) // 2
        y += row_count * (ch + Inches(0.07)) + Inches(0.12)

        if y > _SH - Inches(0.55):
            break


def _slide_section(prs, sec):
    """One content slide per section, SAP 360 style."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, _SW, _SH, _C["white"])

    title_lower = sec["title"].lower()
    hdr_col = _C["mango"] if "risk" in title_lower else _badge_col(sec["badge_color"])
    _section_header(sl, sec["title"], sec["badge_text"], hdr_col)

    ap = sec["ap_tiers"]
    if "executive" in title_lower:
        _render_exec_overview(sl, sec)
    elif "environment" in title_lower:
        _render_environment(sl, sec)
    elif "risk" in title_lower:
        _render_risks(sl, sec)
    elif "action" in title_lower or sec["id"] == "sap":
        if any(ap.values()):
            _render_action_plan(sl, ap)
        else:
            _render_exec_overview(sl, sec)
    else:
        _render_exec_overview(sl, sec)


# =============================================================================
# Main PPTX builder
# =============================================================================

def build_pptx_from_html(html_path, logo_path):
    if not PPTX_OK:
        raise RuntimeError("python-pptx is not installed.\n\nRun:  pip install python-pptx")

    html         = html_path.read_text(encoding="utf-8", errors="replace")
    meta         = _parse_report_meta(html)
    all_sections = _parse_sections(html)
    sections     = [s for s in all_sections if _pptx_include(s)]

    prs = Presentation()
    prs.slide_width  = _SW
    prs.slide_height = _SH

    _slide_title(prs, meta, logo_path)
    for sec in sections:
        _slide_section(prs, sec)

    output = html_path.with_suffix(".pptx")
    prs.save(str(output))
    return output


# =============================================================================
# v4/v5 helpers (unchanged)
# =============================================================================

def load_text(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def load_proxy_config():
    try:
        cfg      = json.loads(SETTINGS_FILE.read_text(encoding="utf-8"))
        env      = cfg.get("env", {})
        base_url = env.get("ANTHROPIC_BASE_URL",   "https://api.anthropic.com")
        token    = env.get("ANTHROPIC_AUTH_TOKEN",  "")
        model    = env.get("ANTHROPIC_MODEL",       "claude-sonnet-latest")
        return base_url, token, model
    except Exception:
        return "https://api.anthropic.com", "", "claude-sonnet-latest"


def get_sap_logo_b64():
    if LOGO_PATH.exists():
        with open(LOGO_PATH, "rb") as f:
            return base64.b64encode(f.read()).decode("utf-8")
    return ""


def extract_sid(text):
    for pat in [r"\[environment\].*?SID[:\s]+([A-Z0-9]{3})",
                r"SID[:\s]+([A-Z0-9]{3})"]:
        m = re.search(pat, text, re.IGNORECASE)
        if m:
            return m.group(1).upper()
    return "UNKNOWN"


_DISK_SPACE_GUIDANCE = """

[CRITICAL — DISK SPACE ANALYSIS RULES — READ BEFORE GENERATING ANY DISK-RELATED CONTENT]

The health check output contains TWO fundamentally different types of disk data.
You MUST distinguish them and NEVER conflate them:

TYPE A — HANA Internal Data Volume (from HANA_Disks_Data_Partitions / HANA_DISK_DATA_OVERVIEW):
  - ALLOC_GB = the current on-disk size of HANA's pre-allocated data files (*.dat files on the filesystem)
  - USED_GB  = how much of that pre-allocated space is actually written by HANA internally
  - TOTAL_GB in this view = same as ALLOC_GB (the file size), NOT the filesystem total capacity
  - These numbers say NOTHING about how much free space the underlying filesystem has.
  - HANA automatically extends (grows) its data files onto the filesystem when ALLOC is exhausted.
    This is normal behaviour and controlled by auto-extension configuration (enabled by default in ECS).
  - A high USED_GB/ALLOC_GB ratio means HANA is near its current pre-allocation boundary and will
    soon auto-extend. This is an INFORMATIONAL / MONITORING finding — it is NOT a disk risk.
  - NEVER compute "filesystem usage %" from ALLOC_GB and call it a disk space warning.

TYPE B — OS Filesystem Free Space (from df -h output — may NOT be present in the health check):
  - The actual bytes free on the Linux filesystem where HANA data lives (e.g. /hana/data/<SID>).
  - This is the ONLY data that can determine true disk exhaustion risk.

RULES FOR DISK ASSESSMENT IN THIS REPORT:
  1. Only raise a DISK SPACE WARNING or "Data Volume Near Capacity" action item if OS-level
     filesystem free space is confirmed to be critically low (< 15% free OR < 100 GB absolute free).
  2. If the health check output contains NO df/filesystem data (this is common), do NOT assume
     the filesystem is full or near capacity. State only that filesystem-level monitoring is recommended.
  3. When HANA_Disks_Data_Partitions data is present, report ALLOC_GB and USED_GB as informational
     metrics showing HANA's internal pre-allocation status — label it clearly as such.
  4. If USED_GB is very close to ALLOC_GB (> 90%), note it as: "HANA data volume pre-allocation is
     near its current boundary; HANA will auto-extend onto available filesystem space. Monitor
     filesystem free space on /hana/data to ensure sufficient room for growth."
  5. Never add a red/orange disk alert, kpi-critical card, or esc-banner based solely on
     ALLOC_GB data without confirmed OS filesystem free space being low.
"""


_MINICHECK_PATTERN = re.compile(
    r'\bHANA_[A-Za-z]+(?:_[A-Za-z]+)*_\d+\.\d+\.\d+\+?\b'
)

def _discover_minichecks(health_check_text):
    """Return a deduplicated ordered list of minicheck names found in the input."""
    seen = set()
    ordered = []
    for match in _MINICHECK_PATTERN.finditer(health_check_text):
        name = match.group()
        if name not in seen:
            seen.add(name)
            ordered.append(name)
    return ordered


def build_prompt(template, release_info, health_check_text):
    prompt = template.replace(
        "<<< REPLACE THIS LINE WITH THE FULL HANA HEALTH CHECK SCRIPT OUTPUT >>>",
        health_check_text,
    )
    if release_info:
        prompt += "\n\n[HANA Latest Release Reference]\n" + release_info + "\n"

    minichecks = _discover_minichecks(health_check_text)
    if minichecks:
        mc_list = "\n".join(f"  - {m}" for m in minichecks)
        prompt += f"""

[DYNAMIC MINICHECK DISCOVERY — OVERRIDE HARDCODED LIST]
The following minicheck sections were detected in this specific health check input.
Analyze ALL of them in Section 10 (Minichecks). Do NOT limit analysis to the 7
scripts listed in the template — cover every script name below.
For each one: identify what is abnormal, concerning, or critical; create a graph or
visual summary; and suggest concrete action items.

Detected minichecks ({len(minichecks)} total):
{mc_list}
"""

    prompt += _DISK_SPACE_GUIDANCE
    return prompt


# ── Post-processing ───────────────────────────────────────────────────────────

def _strip_code_fences(html):
    html = re.sub(r"^```html\s*", "", html.strip())
    return re.sub(r"\s*```$", "", html).strip()


def _inject_enhanced_css(html):
    if "Enhanced v4 additions" in html:
        return html
    return html.replace("</style>", ENHANCED_CSS + "\n</style>", 1)


def _inject_sap_logo(html, logo_b64):
    if not logo_b64:
        return html
    img_left = (
        '<img src="data:image/png;base64,' + logo_b64 + '" '
        'height="32" alt="SAP" style="vertical-align:middle">'
    )
    html = re.sub(
        r'(<div[^>]+class=["\'][^"\']*header-logo[^"\']*["\'][^>]*>).*?(</div>)',
        r'\g<1>' + img_left + r'\g<2>', html, flags=re.DOTALL,
    )
    return html


def _inject_footer_version(html):
    today  = datetime.now().strftime("%Y-%m-%d")
    marker = "Enhanced: " + today
    if marker in html:
        return html
    return re.sub(r"(Generated:\s*[^<|]+)", r"\1  |  " + marker, html, count=1)


def _ensure_js(html):
    if "<script" in html.lower():
        return html
    FIX_JS = """
<script>
function toggleSection(id) {
  var body = document.getElementById('body-' + id);
  var chev = document.getElementById('chev-' + id);
  if (!body) return;
  var isOpen = body.classList.contains('open');
  body.classList.toggle('open', !isOpen);
  if (chev) chev.textContent = isOpen ? '\u25bc' : '\u25b2';
}
document.addEventListener('DOMContentLoaded', function() {
  document.querySelectorAll('.nav-pill').forEach(function(pill) {
    pill.addEventListener('click', function() {
      var m = (pill.getAttribute('onclick') || '').match(/'([^']+)'/);
      var el = m && document.getElementById(m[1]);
      if (el) el.scrollIntoView({behavior:'smooth'});
      document.querySelectorAll('.nav-pill').forEach(function(p){ p.classList.remove('active'); });
      pill.classList.add('active');
    });
  });
  var tip = document.getElementById('tooltip');
  if (tip) {
    document.querySelectorAll('.term').forEach(function(el) {
      el.addEventListener('mousemove', function(e) {
        var def = window.GLOSSARY && window.GLOSSARY[el.dataset.term];
        if (!def) return;
        tip.innerHTML = '<strong style="color:#a8d0f0">' + el.dataset.term + '</strong><br>' + def;
        tip.style.display = 'block';
        tip.style.left = (e.clientX + 14) + 'px';
        tip.style.top  = (e.clientY + 14) + 'px';
      });
      el.addEventListener('mouseleave', function() { tip.style.display = 'none'; });
    });
  }
});
</script>
"""
    if "</body>" in html:
        return html.replace("</body>", FIX_JS + "</body>", 1)
    return html + "\n" + FIX_JS + "\n</body>\n</html>"


def post_process_html(html, logo_b64):
    html = _strip_code_fences(html)
    html = _inject_enhanced_css(html)
    html = _inject_sap_logo(html, logo_b64)
    html = _inject_footer_version(html)
    html = _ensure_js(html)
    return html


# =============================================================================
# Tkinter App
# =============================================================================

class AnalysisApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("SAP HANA Full Analysis Report Generator " + APP_VERSION)
        self.configure(bg=PAGE_BG)
        self.minsize(920, 700)
        self.resizable(True, True)

        self._selected_file    = None
        self._last_output_file = None
        self._status_var   = tk.StringVar(value="Ready \u2014 select a health check file to begin.")
        self._progress_var = tk.IntVar(value=0)
        self._pct_var      = tk.StringVar(value="")
        self._cancel_event = threading.Event()
        self._logo_b64     = get_sap_logo_b64()

        base_url, token, model = load_proxy_config()
        self._base_url = base_url
        self._token    = token
        self._model    = model

        self._build_ui()

    def _build_ui(self):
        self._build_header()
        self._build_content()
        self._build_footer()

    def _load_logo_tk_image(self, target_height=42):
        if not self._logo_b64:
            return None
        try:
            from PIL import Image, ImageTk
            import io
            raw = base64.b64decode(self._logo_b64)
            img = Image.open(io.BytesIO(raw))
            w, h = img.size
            new_w = max(1, int(w * target_height / h))
            img = img.resize((new_w, target_height), Image.LANCZOS)
            return ImageTk.PhotoImage(img)
        except Exception:
            pass
        try:
            photo  = tk.PhotoImage(data=self._logo_b64)
            factor = max(1, photo.height() // target_height)
            return photo.subsample(factor, factor) if factor > 1 else photo
        except Exception:
            return None

    def _build_header(self):
        hdr = tk.Frame(self, bg=DARK_BLUE)
        hdr.pack(fill="x")

        left = tk.Frame(hdr, bg=DARK_BLUE)
        left.pack(side="left", padx=20, pady=12)
        tk.Label(left, text="SAP HANA Full Analysis Report Generator",
                 bg=DARK_BLUE, fg=WHITE, font=("Segoe UI", 15, "bold")).pack(anchor="w")
        tk.Label(left, text="Powered by Anthropic SDK  |  " + APP_VERSION + "  |  " + self._model,
                 bg=DARK_BLUE, fg="#a8d0f0", font=("Segoe UI", 8)).pack(anchor="w")

        right = tk.Frame(hdr, bg=DARK_BLUE)
        right.pack(side="right", padx=16, pady=8)
        tk.Label(right, text="RDE LAC", bg="#e65c00", fg=WHITE,
                 font=("Segoe UI", 9, "bold"), padx=12, pady=4).pack(side="right", padx=(10, 0))
        self._tk_logo = self._load_logo_tk_image(target_height=42)
        if self._tk_logo:
            tk.Label(right, image=self._tk_logo, bg=DARK_BLUE, bd=0).pack(side="right")
        else:
            tk.Label(right, text="SAP_LOGO.png not found", bg=DARK_BLUE, fg=WARN_ORANGE,
                     font=("Segoe UI", 7)).pack(side="right")

    def _build_content(self):
        content = tk.Frame(self, bg=PAGE_BG, padx=20, pady=14)
        content.pack(fill="both", expand=True)

        fframe = self._lframe(content, "Health Check Input File")
        fframe.pack(fill="x", pady=(0, 10))
        self._file_label = tk.Label(fframe, text="No file selected",
                                    bg=WHITE, fg="#888", font=("Segoe UI", 9), anchor="w")
        self._file_label.pack(side="left", fill="x", expand=True)
        self._size_badge = tk.Label(fframe, text="", bg=LIGHT_BLUE, fg=PRIMARY_BLUE,
                                    font=("Segoe UI", 8, "bold"), padx=8, pady=2)
        self._size_badge.pack(side="left", padx=(6, 8))
        tk.Button(fframe, text="Browse...", command=self._on_browse,
                  bg=PRIMARY_BLUE, fg=WHITE, font=("Segoe UI", 9, "bold"),
                  relief="flat", padx=14, cursor="hand2").pack(side="right")

        pframe = self._lframe(content, "File Preview  (first 60 lines)")
        pframe.pack(fill="both", expand=True, pady=(0, 10))
        self._preview_text = scrolledtext.ScrolledText(
            pframe, height=10, font=("Consolas", 8),
            bg="#f8fbff", fg="#333", relief="flat", state="disabled", wrap="none")
        self._preview_text.pack(fill="both", expand=True)

        pgframe = self._lframe(content, "Generation Progress")
        pgframe.pack(fill="x", pady=(0, 10))
        row = tk.Frame(pgframe, bg=WHITE)
        row.pack(fill="x")
        self._progress = ttk.Progressbar(row, variable=self._progress_var,
                                          maximum=100, mode="determinate", length=100)
        self._progress.pack(side="left", fill="x", expand=True, padx=(0, 10))
        tk.Label(row, textvariable=self._pct_var, bg=WHITE, fg=PRIMARY_BLUE,
                 font=("Segoe UI", 10, "bold"), width=6).pack(side="left")
        self._phase_label = tk.Label(pgframe, text="", bg=WHITE, fg="#555",
                                     font=("Segoe UI", 8), anchor="w")
        self._phase_label.pack(fill="x", pady=(4, 0))

        lframe = self._lframe(content, "Generation Log")
        lframe.pack(fill="x", pady=(0, 10))
        self._log_text = scrolledtext.ScrolledText(
            lframe, height=6, font=("Consolas", 8),
            bg="#f0f7ff", fg="#333", relief="flat", state="disabled", wrap="word")
        self._log_text.pack(fill="both", expand=True)

        btn_frame = tk.Frame(content, bg=PAGE_BG)
        btn_frame.pack(fill="x")

        self._run_btn = tk.Button(
            btn_frame, text="  Generate Full Analysis Report",
            command=self._on_generate, bg=PRIMARY_BLUE, fg=WHITE,
            font=("Segoe UI", 11, "bold"), relief="flat",
            padx=24, pady=10, cursor="hand2", state="disabled")
        self._run_btn.pack(side="left")

        self._cancel_btn = tk.Button(
            btn_frame, text="  Cancel",
            command=self._on_cancel, bg=DANGER_RED, fg=WHITE,
            font=("Segoe UI", 10), relief="flat",
            padx=14, pady=10, cursor="hand2", state="disabled")
        self._cancel_btn.pack(side="left", padx=(8, 0))

        self._open_btn = tk.Button(
            btn_frame, text="  Open in Edge",
            command=self._on_open_last, bg=OK_GREEN, fg=WHITE,
            font=("Segoe UI", 10), relief="flat",
            padx=16, pady=10, cursor="hand2", state="disabled")
        self._open_btn.pack(side="left", padx=(8, 0))

        self._pptx_btn = tk.Button(
            btn_frame, text="  Export to PowerPoint",
            command=self._on_pptx_generate, bg=PURPLE, fg=WHITE,
            font=("Segoe UI", 10), relief="flat",
            padx=16, pady=10, cursor="hand2", state="disabled")
        # self._pptx_btn.pack(side="left", padx=(8, 0))  # hidden

        if not PPTX_OK:
            self._pptx_btn.config(text="  python-pptx not installed", bg="#888", cursor="arrow")

    def _build_footer(self):
        footer = tk.Frame(self, bg=DARK_BLUE)
        footer.pack(fill="x", side="bottom")
        tk.Label(footer, textvariable=self._status_var, bg=DARK_BLUE, fg="#a8d0f0",
                 font=("Segoe UI", 8), anchor="w", padx=12).pack(
                     side="left", pady=4, fill="x", expand=True)
        tk.Label(footer, text="RDE LAC", bg=DARK_BLUE, fg=ACCENT_CYAN,
                 font=("Segoe UI", 8, "bold"), padx=16).pack(side="right", pady=4)

    def _lframe(self, parent, title):
        return tk.LabelFrame(parent, text=title, bg=WHITE, fg=DARK_BLUE,
                              font=("Segoe UI", 9, "bold"), padx=10, pady=8,
                              relief="flat", highlightbackground="#c8ddf0",
                              highlightthickness=1)

    def _on_browse(self):
        initial = str(WORKING_DIR / "HANA Health Check Reports")
        path = filedialog.askopenfilename(
            title="Select HANA Health Check Output File",
            initialdir=initial if os.path.isdir(initial) else str(WORKING_DIR),
            filetypes=[("Text files", "*.txt"), ("All files", "*.*")],
        )
        if not path:
            return
        self._selected_file = Path(path)
        size_kb       = self._selected_file.stat().st_size / 1024
        approx_tokens = int(size_kb * 1024 / 4)
        self._file_label.config(text=self._selected_file.name, fg=DARK_BLUE)
        self._size_badge.config(
            text="{0:.0f} KB  ~{1:,} tokens".format(size_kb, approx_tokens),
            bg=WARN_ORANGE if size_kb > MAX_FILE_KB * 0.75 else LIGHT_BLUE,
            fg=DARK_BLUE   if size_kb > MAX_FILE_KB * 0.75 else PRIMARY_BLUE,
        )
        self._load_preview()
        if size_kb > MAX_FILE_KB:
            self._run_btn.config(state="disabled")
            self._set_status("File too large ({0:.0f} KB) \u2014 max {1} KB.".format(
                size_kb, MAX_FILE_KB))
            messagebox.showwarning("File Too Large",
                "{0}\n{1:.0f} KB (~{2:,} tokens)\n\nExceeds {3} KB limit.".format(
                    self._selected_file.name, size_kb, approx_tokens, MAX_FILE_KB))
        else:
            self._run_btn.config(state="normal")
            self._set_status("Ready  |  {0}  ({1:.0f} KB)".format(
                self._selected_file.name, size_kb))

    def _load_preview(self):
        try:
            lines   = load_text(self._selected_file).splitlines()[:60]
            preview = "\n".join(lines)
        except Exception as e:
            preview = "Could not read file: " + str(e)
        self._preview_text.config(state="normal")
        self._preview_text.delete("1.0", "end")
        self._preview_text.insert("end", preview)
        self._preview_text.config(state="disabled")

    def _on_generate(self):
        if not self._selected_file or not self._selected_file.exists():
            messagebox.showerror("No File", "Please select a valid health check file.")
            return
        self._cancel_event.clear()
        self._run_btn.config(state="disabled")
        self._cancel_btn.config(state="normal")
        self._open_btn.config(state="disabled")
        self._pptx_btn.config(state="disabled")
        self._log_clear()
        self._set_progress(0, "Starting...")
        self._set_status("Generating report via streaming API...")
        threading.Thread(target=self._run_generation, daemon=True).start()

    def _on_cancel(self):
        self._cancel_event.set()
        self._cancel_btn.config(state="disabled")
        self._set_status("Cancelling...")

    def _on_open_last(self):
        if self._last_output_file and self._last_output_file.exists():
            self._open_in_edge(self._last_output_file)
        else:
            messagebox.showinfo("No Report", "No report has been generated yet.")

    def _run_generation(self):
        try:
            health_check_text = load_text(self._selected_file)
            template          = load_text(TEMPLATE_FILE)
            release_info      = load_text(RELEASE_FILE) if RELEASE_FILE.exists() else ""

            sid = extract_sid(health_check_text)
            if sid == "UNKNOWN":
                sid = self._selected_file.stem.split("_")[0].upper()

            ts          = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = OUTPUT_DIR / (sid + "_HANA_Analysis_Report_" + ts + ".html")
            OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

            prompt = build_prompt(template, release_info, health_check_text)
            discovered_mc = _discover_minichecks(health_check_text)

            self.after(0, self._log, "File    : " + self._selected_file.name + "\n")
            self.after(0, self._log, "Prompt  : {:,} chars (~{:,} tokens)\n".format(
                len(prompt), len(prompt) // 4))
            self.after(0, self._log, "Minichecks detected: {}\n".format(len(discovered_mc)))
            if discovered_mc:
                for mc in discovered_mc:
                    self.after(0, self._log, "  • " + mc + "\n")
            self.after(0, self._log, "Model   : " + self._model + "\n")
            self.after(0, self._log, "Output  : " + output_file.name + "\n")
            self.after(0, self._log, "Logo    : " + (
                'embedded' if self._logo_b64 else 'not found') + "\n")
            self.after(0, self._set_progress, 5, "Connecting to proxy...")

            client = anthropic.Anthropic(base_url=self._base_url, api_key=self._token)

            self.after(0, self._log, "Streaming response...\n")
            self.after(0, self._set_progress, 10, "Claude is generating the HTML report...")

            html_chunks   = []
            output_tokens = 0
            start_time    = datetime.now()

            with client.messages.stream(
                model=self._model,
                max_tokens=MAX_OUTPUT_TOKENS,
                messages=[{"role": "user", "content": prompt}],
            ) as stream:
                for text_chunk in stream.text_stream:
                    if self._cancel_event.is_set():
                        break
                    html_chunks.append(text_chunk)
                    output_tokens += len(text_chunk) // 4
                    pct     = min(10 + int(output_tokens / EXPECTED_OUT_TOKENS * 85), 95)
                    elapsed = int((datetime.now() - start_time).total_seconds())
                    self.after(0, self._set_progress, pct,
                               "Generating...  ~{:,} tokens  [{}m {:02d}s]".format(
                                   output_tokens, elapsed // 60, elapsed % 60))

            if self._cancel_event.is_set():
                self.after(0, self._on_generation_error, "Cancelled by user.")
                return

            elapsed = int((datetime.now() - start_time).total_seconds())
            self.after(0, self._log,
                       "Stream complete in {}m {:02d}s\n".format(elapsed // 60, elapsed % 60))
            self.after(0, self._set_progress, 96, "Post-processing HTML...")

            raw_html = "".join(html_chunks)
            raw_html = _strip_code_fences(raw_html)

            if not raw_html.startswith("<"):
                raise ValueError(
                    "Response is not valid HTML.\n\nFirst 300 chars:\n" + raw_html[:300])

            self.after(0, self._log, "Applying enhancements: CSS / logo / footer...\n")
            full_html = post_process_html(raw_html, self._logo_b64)

            self.after(0, self._set_progress, 98, "Saving...")
            output_file.write_text(full_html, encoding="utf-8")
            size_kb = output_file.stat().st_size // 1024
            self.after(0, self._log, "Saved: {}  ({} KB)\n".format(output_file.name, size_kb))

            self._last_output_file = output_file
            self.after(0, self._on_generation_complete, output_file)

        except anthropic.APIStatusError as e:
            self.after(0, self._on_generation_error,
                       "API error {}: {}".format(e.status_code, e.message))
        except Exception as e:
            self.after(0, self._on_generation_error, str(e))

    def _on_generation_complete(self, output_file):
        self._set_progress(100, "Complete!")
        self._run_btn.config(state="normal")
        self._cancel_btn.config(state="disabled")
        self._open_btn.config(state="normal")
        if PPTX_OK:
            self._pptx_btn.config(state="normal")
        self._log("\n  Report ready: {}\n".format(output_file))
        self._set_status("Done \u2014 " + output_file.name)
        self._open_in_edge(output_file)

    def _on_generation_error(self, msg):
        self._set_progress(0, "")
        self._run_btn.config(state="normal")
        self._cancel_btn.config(state="disabled")
        self._log("\n  Error: {}\n".format(msg))
        self._set_status("Error \u2014 see log.")
        messagebox.showerror("Generation Failed", msg)

    # ── PowerPoint export ──────────────────────────────────────────────────────

    def _on_pptx_generate(self):
        if not PPTX_OK:
            messagebox.showerror("python-pptx not installed",
                "Install it with:\n\n    pip install python-pptx\n\nthen restart the app.")
            return
        if not self._last_output_file or not self._last_output_file.exists():
            messagebox.showinfo("No Report", "Generate an HTML report first.")
            return
        self._pptx_btn.config(state="disabled")
        self._set_status("Building PowerPoint presentation...")
        self._log("\nBuilding PowerPoint...\n")
        threading.Thread(
            target=self._run_pptx_generation,
            args=(self._last_output_file,),
            daemon=True,
        ).start()

    def _run_pptx_generation(self, html_path):
        try:
            self.after(0, self._log, "Parsing: {}\n".format(html_path.name))
            pptx_path = build_pptx_from_html(html_path, LOGO_PATH)
            size_kb   = pptx_path.stat().st_size // 1024
            self.after(0, self._log, "Saved: {}  ({} KB)\n".format(pptx_path.name, size_kb))
            self.after(0, self._on_pptx_complete, pptx_path)
        except Exception as e:
            self.after(0, self._on_pptx_error, str(e))

    def _on_pptx_complete(self, pptx_path):
        self._pptx_btn.config(state="normal")
        self._set_status("PowerPoint ready \u2014 " + pptx_path.name)
        self._log("  PowerPoint ready: {}\n".format(pptx_path))
        import subprocess
        try:
            os.startfile(str(pptx_path))
        except Exception:
            try:
                subprocess.Popen(["explorer", "/select,", str(pptx_path)])
            except Exception:
                pass
        messagebox.showinfo("PowerPoint Ready",
            "Saved to:\n{}\n\nOpening file...".format(pptx_path))

    def _on_pptx_error(self, msg):
        self._pptx_btn.config(state="normal")
        self._set_status("PowerPoint generation failed \u2014 see log.")
        self._log("\n  PPTX Error: {}\n".format(msg))
        messagebox.showerror("PowerPoint Failed", msg)

    # ── Utilities ──────────────────────────────────────────────────────────────

    def _open_in_edge(self, path):
        import subprocess
        try:
            subprocess.Popen([EDGE_PATH, str(path)])
            self._set_status("Opened in Edge: " + path.name)
        except Exception as e:
            messagebox.showerror("Could not open Edge", str(e))

    def _set_progress(self, pct, phase):
        self._progress_var.set(pct)
        self._pct_var.set("{}%".format(pct) if pct > 0 else "")
        self._phase_label.config(text=phase)

    def _log(self, msg):
        self._log_text.config(state="normal")
        self._log_text.insert("end", msg)
        self._log_text.see("end")
        self._log_text.config(state="disabled")

    def _log_clear(self):
        self._log_text.config(state="normal")
        self._log_text.delete("1.0", "end")
        self._log_text.config(state="disabled")

    def _set_status(self, msg):
        self._status_var.set(msg)


if __name__ == "__main__":
    app = AnalysisApp()
    app.mainloop()
