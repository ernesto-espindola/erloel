#!/usr/bin/env python3
"""Generate PowerPoint from BHPXXX HANA Analysis Report."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────
BASE = Path(r"C:\Users\I522148\OneDrive - SAP SE\SWAT\AI\Claude\working directory")
OUTPUT = BASE / "Results" / "BHPXXX_HANA_Analysis_Report_20260813.pptx"
LOGO  = BASE / "SAP_LOGO.png"

# ── Colors ─────────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x00, 0x33, 0x66)
C_BLUE   = RGBColor(0x00, 0x57, 0xa8)
C_CYAN   = RGBColor(0x00, 0xb4, 0xe6)
C_RED    = RGBColor(0xd9, 0x53, 0x4f)
C_ORANGE = RGBColor(0xf0, 0xad, 0x4e)
C_GREEN  = RGBColor(0x5c, 0xb8, 0x5c)
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_BG     = RGBColor(0xe8, 0xf4, 0xfd)
C_CARD   = RGBColor(0xf0, 0xf7, 0xff)
C_BODY   = RGBColor(0x21, 0x25, 0x29)
C_GREY   = RGBColor(0x66, 0x66, 0x66)
C_EDGE   = RGBColor(0xe0, 0xec, 0xf8)
C_RED_BG   = RGBColor(0xff, 0xe8, 0xe8)
C_ORG_BG   = RGBColor(0xff, 0xfd, 0xf0)
C_GRN_BG   = RGBColor(0xf0, 0xff, 0xf4)

# ── Slide dimensions (16:9) ────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────────────
def R(slide, x, y, w, h, bg, border=None):
    """Add a solid rectangle."""
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = bg
    if border:
        s.line.color.rgb = border; s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s

def T(slide, text, x, y, w, h, sz=10, bold=False, color=None,
      align=PP_ALIGN.LEFT, italic=False):
    """Add a text box."""
    color = color or C_BODY
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return tb

def hdr(slide, icon, title, badge, badge_col):
    """Dark-blue header bar with title and status badge."""
    R(slide, 0, 0, SW, Inches(0.88), C_DARK)
    T(slide, f"{icon}  {title}", Inches(0.3), Inches(0.1),
      Inches(10.5), Inches(0.7), sz=18, bold=True, color=C_WHITE)
    bw = Inches(2.0)
    bx = SW - bw - Inches(0.25)
    R(slide, bx, Inches(0.2), bw, Inches(0.5), badge_col)
    T(slide, badge, bx, Inches(0.2), bw, Inches(0.5),
      sz=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def card(slide, title, lines, x, y, w, h, left_col=C_BLUE):
    """Info card with colored left accent bar."""
    bg = {C_RED: C_RED_BG, C_ORANGE: C_ORG_BG, C_GREEN: C_GRN_BG}.get(left_col, C_CARD)
    R(slide, x, y, w, h, bg, left_col)
    R(slide, x, y, Inches(0.08), h, left_col)
    T(slide, title, x + Inches(0.16), y + Inches(0.08),
      w - Inches(0.24), Inches(0.28), sz=9.5, bold=True, color=C_DARK)
    body = "\n".join(f"▶  {l}" for l in lines)
    T(slide, body, x + Inches(0.16), y + Inches(0.38),
      w - Inches(0.24), h - Inches(0.44), sz=8.5, color=C_BODY)

def kpi(slide, val, lbl, x, y, w, h, vc=C_BLUE):
    """KPI metric box with top color bar."""
    R(slide, x, y, w, h, C_WHITE, C_EDGE)
    R(slide, x, y, w, Inches(0.06), vc)
    T(slide, val, x, y + Inches(0.1), w, Inches(0.45),
      sz=22, bold=True, color=vc, align=PP_ALIGN.CENTER)
    T(slide, lbl, x, y + Inches(0.54), w, Inches(0.26),
      sz=7.5, color=C_GREY, align=PP_ALIGN.CENTER)

def tbl_hdr(slide, cols, y):
    """Table header row: list of (text, x_in, w_in)."""
    R(slide, Inches(0.28), y, SW - Inches(0.56), Inches(0.36), C_DARK)
    for txt, xi, wi in cols:
        T(slide, txt, Inches(xi), y + Inches(0.04),
          Inches(wi), Inches(0.28), sz=8.5, bold=True, color=C_WHITE)

def tbl_row(slide, cells, y, h, even, left_col=None):
    """Table data row."""
    bg = C_CARD if even else C_WHITE
    R(slide, Inches(0.28), y, SW - Inches(0.56), h, bg, C_EDGE)
    if left_col:
        R(slide, Inches(0.28), y, Inches(0.07), h, left_col)
    for txt, xi, wi, col, bld in cells:
        T(slide, txt, Inches(xi), y + Inches(0.06), Inches(wi),
          h - Inches(0.1), sz=8.5, bold=bld, color=col or C_BODY)

def pill(slide, text, x, y, w, h, bg, fg=C_WHITE):
    R(slide, x, y, w, h, bg)
    T(slide, text, x, y, w, h, sz=8, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_DARK)
R(sl, 0, SH - Inches(0.1), SW, Inches(0.1), C_CYAN)

# SAP logo image (if available)
logo_y = Inches(0.35)
if LOGO.exists():
    sl.shapes.add_picture(str(LOGO), Inches(0.45), logo_y, height=Inches(0.55))

T(sl, "SAP HANA Database",
  Inches(0.45), Inches(1.1), Inches(12), Inches(0.7),
  sz=20, color=C_CYAN)
T(sl, "Health Check Report",
  Inches(0.45), Inches(1.82), Inches(12), Inches(1.1),
  sz=42, bold=True, color=C_WHITE)
T(sl, "BHP / HP4  —  Tenant PS4",
  Inches(0.45), Inches(2.95), Inches(12), Inches(0.65),
  sz=24, bold=False, color=RGBColor(0xa8, 0xd0, 0xf0))

R(sl, Inches(0.45), Inches(3.72), Inches(8), Inches(0.04), C_CYAN)

meta = [
    ("Host",         "vhgcbhp4db01"),
    ("Instance",     "02"),
    ("HANA Version", "2.00.089.00  (SPS08)"),
    ("System Type",  "MDC  ·  VM  ·  SAP Data Center  ·  SLES 15 SP5"),
    ("Report Date",  "2026-08-12"),
]
for i, (lbl, val) in enumerate(meta):
    ry = Inches(3.9 + i * 0.5)
    T(sl, lbl, Inches(0.45), ry, Inches(2.5), Inches(0.42),
      sz=10, bold=True, color=RGBColor(0xa8, 0xd0, 0xf0))
    T(sl, val, Inches(3.1), ry, Inches(9.5), Inches(0.42),
      sz=10, color=C_WHITE)

T(sl, "Prepared by SAP ECS SWAT  |  Health Check Tool v1.00  |  Enhanced v4",
  Inches(0.45), SH - Inches(0.55), Inches(12), Inches(0.38),
  sz=8.5, color=RGBColor(0x60, 0x88, 0xb8))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_BG)
hdr(sl, "📋", "Executive Overview", "⚠ 1 Critical  2 Warnings", C_ORANGE)

health = [
    ("💰 License",     "🔴 CRITICAL", "License type undefined\n40 P1 alerts firing non-stop",           C_RED),
    ("💾 Memory",      "🟡 WARNING",  "85.7% RAM after restart\nAllocation limit near physical RAM",    C_ORANGE),
    ("🔐 Security",    "🟡 WARNING",  "CVE-2026-0492 active\nSelf-signed certificates in use",          C_ORANGE),
    ("💿 Backups",     "🟢 OK",       "0 failures in 40 days\n45 successful runs, healthy throughput",  C_GREEN),
    ("⚡ Performance", "🟡 WARNING",  "BATCH_USER: 28–30 GB/query\nBALDAT lock timeouts up to 1800s",   C_ORANGE),
    ("🛡 Stability",   "🟢 OK",       "0 OOM events, 0 CPU spikes\n0.42% average CPU utilization",      C_GREEN),
]

cw = Inches(4.1); ch = Inches(1.28)
x0 = Inches(0.28); gap = Inches(0.24)
for i, (dim, status, detail, col) in enumerate(health):
    cx = x0 + (i % 3) * (cw + gap)
    cy = Inches(1.0) + (i // 3) * (ch + Inches(0.14))
    R(sl, cx, cy, cw, ch, C_WHITE, col)
    R(sl, cx, cy, cw, Inches(0.07), col)
    T(sl, dim,    cx + Inches(0.12), cy + Inches(0.1),  cw - Inches(0.2), Inches(0.28), sz=10, bold=True,  color=C_DARK)
    T(sl, status, cx + Inches(0.12), cy + Inches(0.4),  cw - Inches(0.2), Inches(0.28), sz=9,  bold=True,  color=col)
    T(sl, detail, cx + Inches(0.12), cy + Inches(0.7),  cw - Inches(0.2), Inches(0.52), sz=8.5, color=C_BODY)

# Critical finding bar
R(sl, Inches(0.28), Inches(3.85), SW - Inches(0.56), Inches(1.0), C_RED_BG, C_RED)
R(sl, Inches(0.28), Inches(3.85), Inches(0.08), Inches(1.0), C_RED)
T(sl, "⚠  Critical Finding",
  Inches(0.46), Inches(3.9), Inches(12.5), Inches(0.28), sz=10, bold=True, color=C_RED)
T(sl, "40 continuous P1 compliance alerts caused by a single misconfiguration: "
      "license_usage = 'undefined'. A single ALTER SYSTEM command clears all 40 alerts immediately. "
      "Additionally, gc_unused_memory_threshold_abs is set 2× the ECS standard, "
      "delaying memory garbage collection. Both are Priority 1 — resolvable in under 1 hour.",
  Inches(0.46), Inches(4.2), SW - Inches(0.7), Inches(0.58), sz=9, color=C_BODY)

T(sl, "The BHP/PS4 system is operationally stable — zero OOM events, CPU spikes, or backup failures over 40 days. "
      "Primary focus: license compliance and technical table cleanup to reclaim disk space and reduce memory footprint.",
  Inches(0.28), Inches(5.05), SW - Inches(0.56), Inches(0.55),
  sz=9, italic=True, color=C_GREY)

card(sl, "Key Numbers at a Glance",
     ["40 P1 Alerts (all = Alert 140 license)  ·  1 ECS Parameter Error",
      "75.5 GB in 11 technical tables  ·  82.7% disk utilization",
      "67.7% HANA memory used  ·  12 action items across 3 priority tiers"],
     Inches(0.28), Inches(5.72), SW - Inches(0.56), Inches(1.0))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ENVIRONMENT SUMMARY
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_BG)
hdr(sl, "🖥", "Environment Summary", "System Info", C_BLUE)

kpis = [
    ("67%",    "Memory Used",    C_ORANGE),
    ("83%",    "Disk Used",      C_ORANGE),
    ("0.42%",  "Avg CPU",        C_GREEN),
    ("40",     "P1 Alerts",      C_RED),
    ("0",      "OOM Events",     C_GREEN),
    ("0",      "Failed Backups", C_GREEN),
]
kw = Inches(2.05); kh = Inches(0.88)
for i, (v, l, c) in enumerate(kpis):
    kpi(sl, v, l, Inches(0.28) + i * (kw + Inches(0.1)), Inches(1.0), kw, kh, c)

# Two-column config grid
left_data = [
    ("SID / Tenant",   "HP4 (BHP)  /  PS4"),
    ("Host",           "vhgcbhp4db01"),
    ("Instance",       "02"),
    ("HANA Version",   "2.00.089.00  (SPS08)"),
    ("ECS Standard",   "2.00.089.03  — 3 patches behind"),
    ("System Type",    "MDC  ·  VM"),
    ("Cloud",          "SAP Data Center"),
    ("OS",             "SLES 15 SP5"),
]
right_data = [
    ("Physical RAM",   "503.83 GB"),
    ("HANA Alloc",     "484 GB  (alloc limit)"),
    ("HANA Mem Used",  "144.9 / 214 GB  (67%)"),
    ("Data Volume",    "988 / 1195 GB  (83%)"),
    ("Log Volume",     "74 / 256 GB  (29%)"),
    ("Logical CPUs",   "32"),
    ("Replication",    "Not configured"),
    ("NSE Tables",     "297  |  LOBs: 31.1M packed"),
]
for col_data, x0i in [(left_data, 0.28), (right_data, 6.88)]:
    cw2 = Inches(6.2) if x0i == 0.28 else Inches(6.17)
    R(sl, Inches(x0i), Inches(2.02), cw2, Inches(0.34), C_DARK)
    lbl = "System Configuration" if x0i == 0.28 else "Resource Utilization"
    T(sl, lbl, Inches(x0i + 0.12), Inches(2.05), cw2 - Inches(0.2), Inches(0.26),
      sz=9.5, bold=True, color=C_WHITE)
    for i, (k, v) in enumerate(col_data):
        bg = C_CARD if i % 2 == 0 else C_WHITE
        ry = Inches(2.36 + i * 0.46)
        R(sl, Inches(x0i), ry, cw2, Inches(0.44), bg, C_EDGE)
        T(sl, k, Inches(x0i + 0.12), ry + Inches(0.06), Inches(2.5), Inches(0.34), sz=8.5, bold=True, color=C_DARK)
        T(sl, v, Inches(x0i + 2.72), ry + Inches(0.06), cw2 - Inches(2.82), Inches(0.34), sz=8.5, color=C_BODY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HANA VERSION
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_BG)
hdr(sl, "🔢", "HANA Database Version", "⚠ Minor Patch Behind", C_ORANGE)

# Version banner
R(sl, Inches(0.28), Inches(1.02), SW - Inches(0.56), Inches(1.75), C_BLUE)
T(sl, "INSTALLED", Inches(0.55), Inches(1.1), Inches(5.5), Inches(0.28), sz=8.5, bold=True, color=C_CYAN)
T(sl, "2.00.089.00", Inches(0.55), Inches(1.38), Inches(5.5), Inches(0.7), sz=32, bold=True, color=C_WHITE)
T(sl, "SPS08  ·  Revision 89  ·  Initial Release", Inches(0.55), Inches(2.1), Inches(5.5), Inches(0.3),
  sz=9.5, color=RGBColor(0xd0, 0xea, 0xff))

T(sl, "vs", Inches(6.4), Inches(1.65), Inches(0.6), Inches(0.4), sz=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

T(sl, "ECS STANDARD", Inches(7.2), Inches(1.1), Inches(5.8), Inches(0.28), sz=8.5, bold=True, color=C_CYAN)
T(sl, "2.00.089.03", Inches(7.2), Inches(1.38), Inches(5.8), Inches(0.7), sz=32, bold=True, color=RGBColor(0xa8, 0xff, 0xd0))
T(sl, "SPS08  ·  Revision 89  ·  Patch 03", Inches(7.2), Inches(2.1), Inches(5.8), Inches(0.3),
  sz=9.5, color=RGBColor(0xd0, 0xea, 0xff))

# Match badge
R(sl, Inches(11.0), Inches(1.42), Inches(2.1), Inches(0.38), C_ORANGE)
T(sl, "⚠ MINOR UPDATE", Inches(11.0), Inches(1.42), Inches(2.1), Inches(0.38),
  sz=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

card(sl, "⚠  Upgrade Recommendation",
     ["Revision 89.00 is 3 patch levels behind ECS standard 89.03 (same SPS08 track — low-risk upgrade)",
      "CVE-2026-0492 is active for SAP HANA; revision 89.03 may include the security fix",
      "Plan the upgrade during the next maintenance window (Priority 2 — This Week)",
      "Owner: Security + Basis  ·  Effort: Maintenance window  ·  Prevents: CVE exposure"],
     Inches(0.28), Inches(3.0), SW - Inches(0.56), Inches(1.25), C_ORANGE)

card(sl, "SPS08 Context",
     ["SPS08 is the current and latest Support Package Stack for HANA 2.0",
      "No SPS upgrade required — revision patch only (89.00 → 89.03)",
      "Always verify SAP Note compatibility list before applying patches"],
     Inches(0.28), Inches(4.38), SW - Inches(0.56), Inches(1.1))

card(sl, "What's New in 89.03 (vs 89.00)",
     ["Security patches including potential CVE-2026-0492 remediation",
      "Stability fixes for column store and query processing",
      "Review patch release notes (SAP Launchpad) for full change list before upgrade"],
     Inches(0.28), Inches(5.6), SW - Inches(0.56), Inches(1.1), C_BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ECS PARAMETERS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_BG)
hdr(sl, "⚙", "ECS Standard Parameters", "1 ERROR  |  Parameters Checked", C_RED)

for v, l, c, xi in [("1", "ERROR", C_RED, 0.28), ("87", "OK", C_GREEN, 3.2), ("12", "OK-Strict", C_ORANGE, 6.12)]:
    kpi(sl, v, l, Inches(xi), Inches(1.0), Inches(2.6), Inches(0.86), c)

# Error param table
tcols = [("INI File", 0.3, 2.0), ("Section", 2.35, 1.9), ("Parameter", 4.3, 3.1),
         ("ECS Expected", 7.45, 2.4), ("Actual Value", 9.9, 1.55), ("Status", 11.5, 1.3)]
tbl_hdr(sl, tcols, Inches(2.0))

R(sl, Inches(0.28), Inches(2.36), SW - Inches(0.56), Inches(0.46), C_RED_BG, C_RED)
R(sl, Inches(0.28), Inches(2.36), Inches(0.07), Inches(0.46), C_RED)
row_vals = [("global.ini", 0.3, 2.0), ("memorymanager", 2.35, 1.9),
            ("gc_unused_memory_threshold_abs", 4.3, 3.1),
            ("36864 MB", 7.45, 2.4), ("78643 MB", 9.9, 1.55)]
for txt, xi, wi in row_vals:
    T(sl, txt, Inches(xi), Inches(2.4), Inches(wi), Inches(0.36), sz=8.5, color=C_BODY)
pill(sl, "ERROR", Inches(11.55), Inches(2.44), Inches(1.15), Inches(0.3), C_RED)

# Consequence block
R(sl, Inches(0.28), Inches(2.9), SW - Inches(0.56), Inches(1.2), C_RED_BG, C_RED)
R(sl, Inches(0.28), Inches(2.9), Inches(0.08), Inches(1.2), C_RED)
T(sl, "⚠  Security / Stability Consequence",
  Inches(0.46), Inches(2.95), Inches(12.5), Inches(0.28), sz=9.5, bold=True, color=C_RED)
T(sl, "gc_unused_memory_threshold_abs controls when HANA garbage-collects unused memory. "
      "At 78,643 MB (2× the ECS standard of 36,864 MB), HANA retains freed memory far longer than necessary, "
      "increasing the overall memory footprint unnecessarily.\n"
      "Fix:  ALTER SYSTEM ALTER CONFIGURATION ('global.ini','SYSTEM') "
      "SET ('memorymanager','gc_unused_memory_threshold_abs') = '36864' WITH RECONFIGURE;",
  Inches(0.46), Inches(3.26), SW - Inches(0.7), Inches(0.77), sz=8.5, color=C_BODY)

card(sl, "Notable OK-Strictness Deviations (no immediate action required)",
     ["sslminprotocolversion — TLS config reviewed; within acceptable range but should be monitored",
      "statement_memory_limit — not set; consider adding a limit to cap BATCH_USER query memory",
      "12 parameters total at OK-strictness high/low — review during next configuration audit"],
     Inches(0.28), Inches(4.26), SW - Inches(0.56), Inches(1.1))

card(sl, "Priority 1 Action — Immediate (≤ 48h)",
     ["Run the ALTER SYSTEM command above to reset gc_unused_memory_threshold_abs to 36,864 MB",
      "Owner: DBA  ·  Effort: 15 minutes  ·  Prevents: Elevated memory footprint, delayed GC cycles"],
     Inches(0.28), Inches(5.5), SW - Inches(0.56), Inches(0.88), C_RED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — P1 ALERTS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_BG)
hdr(sl, "🚨", "P1 Alerts – Last 40 Days", "40 Active Alerts", C_RED)

# Big count
R(sl, Inches(0.28), Inches(1.0), Inches(3.0), Inches(2.0), C_RED_BG, C_RED)
T(sl, "40", Inches(0.28), Inches(1.06), Inches(3.0), Inches(1.3), sz=72, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
T(sl, "P1 ALERT EVENTS",  Inches(0.28), Inches(2.38), Inches(3.0), Inches(0.32), sz=9, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
T(sl, "Continuous — 40 days", Inches(0.28), Inches(2.72), Inches(3.0), Inches(0.24), sz=8, color=C_GREY, align=PP_ALIGN.CENTER)

# Detail block
R(sl, Inches(3.55), Inches(1.0), SW - Inches(3.83), Inches(2.0), C_RED_BG, C_RED)
R(sl, Inches(3.55), Inches(1.0), Inches(0.08), Inches(2.0), C_RED)
T(sl, "Alert 140 — License Usage Type Not Defined",
  Inches(3.74), Inches(1.06), Inches(9.3), Inches(0.32), sz=12, bold=True, color=C_DARK)
T(sl, "HANA Alert 140 fires when license_usage_type is not configured for a production system.\n"
      "On BHP/PS4, the global.ini parameter  [system]  license_usage  is set to 'undefined'.\n\n"
      "Root Cause:  Parameter was never set after system provisioning.\n"
      "Impact:  Ongoing P1 alert storm → alert fatigue → risk of missing real incidents.\n"
      "Fix:  One command clears all 40 alerts immediately (SAP Note 2779499).",
  Inches(3.74), Inches(1.42), Inches(9.3), Inches(1.46), sz=9, color=C_BODY)

# DB breakdown
for txt, detail, col, xi in [
    ("SYSTEMDB:  0 P1 Alerts", "No alerts on the system database.", C_GREEN, 0.28),
    ("Tenant PS4:  40 P1 Alerts  (Alert 140)", "All events = same root cause — single fix.", C_RED, 6.88),
]:
    bw = Inches(6.2) if xi == 0.28 else Inches(6.17)
    R(sl, Inches(xi), Inches(3.2), bw, Inches(0.82), C_CARD if col == C_GREEN else C_RED_BG, col)
    R(sl, Inches(xi), Inches(3.2), Inches(0.08), Inches(0.82), col)
    T(sl, txt, Inches(xi + 0.18), Inches(3.25), bw - Inches(0.28), Inches(0.3), sz=9.5, bold=True, color=col)
    T(sl, detail, Inches(xi + 0.18), Inches(3.57), bw - Inches(0.28), Inches(0.38), sz=8.5, color=C_BODY)

card(sl, "Priority 1 Action — Immediate (≤ 48h)",
     ["ALTER SYSTEM ALTER CONFIGURATION ('global.ini','SYSTEM') "
      "SET ('system','license_usage') = 'production' WITH RECONFIGURE;",
      "Owner: DBA  ·  Effort: 30 minutes  ·  Clears all 40 P1 alerts  ·  SAP Note 2779499"],
     Inches(0.28), Inches(4.18), SW - Inches(0.56), Inches(0.88), C_RED)

card(sl, "Alert Summary",
     ["100% of P1 events share one root cause — single parameter misconfiguration",
      "Resolving license_usage immediately clears the entire alert backlog",
      "No other alert types were detected in the 40-day observation window"],
     Inches(0.28), Inches(5.2), SW - Inches(0.56), Inches(1.0))

T(sl, "Tip: After fixing, run  HANA_Alerts_Notifications  minicheck to confirm alerts cleared.",
  Inches(0.28), Inches(6.32), SW - Inches(0.56), Inches(0.36),
  sz=8.5, italic=True, color=C_GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BIG TECHNICAL TABLES
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_BG)
hdr(sl, "🗄", "Big Technical Tables", "11 Tables  |  75.5 GB Total", C_ORANGE)

for v, l, c, xi in [("11", "Tables Found", C_ORANGE, 0.28), ("75.5 GB", "Total Disk", C_RED, 3.2),
                    ("18.3 GB", "Total Memory", C_ORANGE, 6.12), ("7", "LOB Tables", C_BLUE, 9.04)]:
    kpi(sl, v, l, Inches(xi), Inches(1.0), Inches(2.6), Inches(0.82), c)

# ESC banner
R(sl, Inches(0.28), Inches(1.98), SW - Inches(0.56), Inches(0.38), C_RED_BG, C_RED)
T(sl, "🔴  ESC ALERT: SOFFCONT1 = 405 GB disk on largest-tables minicheck  (exceeds 100 GB threshold)",
  Inches(0.44), Inches(2.02), SW - Inches(0.7), Inches(0.28), sz=9.5, bold=True, color=C_RED)

# Table header
tcols2 = [("Schema", 0.3, 2.1), ("Table Name", 2.45, 2.55), ("Type", 5.05, 1.05),
          ("Disk GB", 6.15, 1.2), ("Mem GB", 7.4, 1.2), ("LOB", 8.65, 1.5), ("Recommended Action", 10.2, 2.9)]
tbl_hdr(sl, tcols2, Inches(2.5))

rows7 = [
    ("SAPHANADB", "REPOLOAD",         "Column", "36.14", " 9.28", "H1/H2", "Archive load objects via SE38 / SAP Note",      C_RED),
    ("SAPHANADB", "REPOSRC",          "Column", "11.19", " 4.62", "H2",    "Delete inactive ABAP sources (SCTS_CLEANUP)",    C_ORANGE),
    ("SAPHANADB", "GRACSODREPDATA",   "Column", "10.09", " 0.62", "—",     "Review GRC SoD retention; archive old reports",  C_ORANGE),
    ("SAPHANADB", "/AIF/TRFCQDATA",   "Column", " 5.07", " 0.31", "H2",    "Archive processed AIF queue entries",            C_ORANGE),
    ("_SYS_REPO", "OBJECT_HISTORY",   "Column", " 4.10", " 0.01", "H2",    "Run _SYS_REPO cleanup procedure",               C_ORANGE),
    ("SAPHANADB", "DYNPLOAD",         "Column", " 2.34", " 0.20", "H1",    "Review and remove inactive dynpros",             C_BLUE),
    ("SAPHANADB", "D010TAB",          "Column", " 1.52", " 1.54", "—",     "ABAP cross-ref cleanup via SAAB",               C_BLUE),
]
for i, (sch, tbl, typ, disk, mem, lob, act, col) in enumerate(rows7):
    rh = Inches(0.44); ry = Inches(2.86 + i * rh)
    bg = C_CARD if i % 2 == 0 else C_WHITE
    R(sl, Inches(0.28), ry, SW - Inches(0.56), rh, bg, C_EDGE)
    R(sl, Inches(0.28), ry, Inches(0.07), rh, col)
    for val, xi, wi in [(sch, 0.3, 2.1), (tbl, 2.45, 2.55), (typ, 5.05, 1.05),
                        (disk, 6.15, 1.2), (mem, 7.4, 1.2), (lob, 8.65, 1.5), (act, 10.2, 2.9)]:
        c2 = col if xi in [6.15, 7.4] else (C_DARK if xi == 2.45 else C_BODY)
        T(sl, val, Inches(xi), ry + Inches(0.06), Inches(wi), rh - Inches(0.1),
          sz=8, bold=(c2 != C_BODY and c2 != C_DARK), color=c2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — LARGE RECORDS / OOM / CPU  (all green)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_BG)
hdr(sl, "✅", "Tables / OOM / CPU — All Clear", "3 Sections  |  All OK", C_GREEN)

panels = [
    ("📊 Tables > 1.5B Records", "✅  None Found",
     "No tables exceed 1.5 billion records in\nSYSTEMDB or tenant PS4.\n\n"
     "⚠ Growth Watch:\n"
     "• FQM_FLOW: +319% (9.9M → 41.6M records)\n"
     "• WBCROSSGT: +35% (163M → 221M)\n"
     "• RSAU_BUF_DATA: 727M records (audit log)\n"
     "• D010TAB: +35% (117M → 158M)\n\n"
     "FQM_FLOW rapid growth requires investigation."),
    ("💾 Out-of-Memory Events", "✅  0 Events",
     "No OOM events in 40-day window.\n\n"
     "Memory Status:\n"
     "• Allocated: 214 GB  |  Used: 144.9 GB (67.7%)\n"
     "• Available: 69.1 GB\n\n"
     "⚠ Post-restart (Aug 9): Physical RAM hit\n"
     "85.7% (431/504 GB) — monitor closely.\n\n"
     "Top allocator: CS/Main/Dict at 9.36% (19.92 GB)"),
    ("⚡ CPU Spikes ≥ 95%", "✅  0 Spikes",
     "No CPU spikes in 40-day window.\n\n"
     "CPU Status:\n"
     "• Average CPU: 0.42%\n"
     "• Peak events (≥95%): 0\n"
     "• Logical CPUs available: 32\n\n"
     "BATCH_USER analytical queries consume up\n"
     "to 1005 CPU-seconds per execution but do\n"
     "not cause system-wide CPU saturation."),
]
pw = Inches(4.1); ph = Inches(4.1)
px0 = Inches(0.28); gap = Inches(0.24)
for i, (ttl, ok, detail) in enumerate(panels):
    cx = px0 + i * (pw + gap)
    R(sl, cx, Inches(1.0), pw, ph, C_GRN_BG, C_GREEN)
    R(sl, cx, Inches(1.0), Inches(0.08), ph, C_GREEN)
    T(sl, ttl, cx + Inches(0.16), Inches(1.06), pw - Inches(0.24), Inches(0.3), sz=9.5, bold=True, color=C_DARK)
    T(sl, ok,  cx + Inches(0.16), Inches(1.42), pw - Inches(0.24), Inches(0.4), sz=16, bold=True, color=C_GREEN)
    T(sl, detail, cx + Inches(0.16), Inches(1.9), pw - Inches(0.24), Inches(3.1), sz=8.5, color=C_BODY)

card(sl, "Preventive Recommendations",
     ["Monitor physical RAM trend — alert threshold: 85%; post-restart at 85.7% needs attention",
      "LOBPage allocator at 13.99 GB correlates with REPOLOAD — archiving will reduce memory pressure",
      "Investigate FQM_FLOW +319% growth — verify FQM background job clears processed flow entries",
      "Apply WorkloadClass memory limits to BATCH_USER to cap queries at 20 GB (Priority 2 — This Week)"],
     Inches(0.28), Inches(5.22), SW - Inches(0.56), Inches(1.26))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FAILED BACKUPS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_BG)
hdr(sl, "💿", "Failed Backups", "✅  0 Failures", C_GREEN)

# Big 0
R(sl, Inches(0.28), Inches(1.0), Inches(2.8), Inches(1.9), C_GRN_BG, C_GREEN)
T(sl, "0", Inches(0.28), Inches(1.04), Inches(2.8), Inches(1.26), sz=72, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
T(sl, "FAILED BACKUPS", Inches(0.28), Inches(2.3), Inches(2.8), Inches(0.3), sz=9, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
T(sl, "in 40-day window", Inches(0.28), Inches(2.62), Inches(2.8), Inches(0.24), sz=8, color=C_GREY, align=PP_ALIGN.CENTER)

# Stats
stats = [("45", "Successful Runs"), ("~987 GB", "Last Backup (compressed)"), ("40–65 min", "Recent Duration"), ("182–355 MB/s", "Throughput")]
sw2 = Inches(2.3); sh2 = Inches(0.84)
for i, (v, l) in enumerate(stats):
    kpi(sl, v, l, Inches(3.45) + i * (sw2 + Inches(0.16)), Inches(1.0), sw2, sh2, C_BLUE)

card(sl, "Backup Run Analysis",
     ["Daily full data backups running consistently — 45 successful runs observed in 40 days",
      "Backup sizes: ~987 GB compressed / ~1,010 GB uncompressed (full tenant PS4)",
      "Duration improved: July backups took 73–134 min; recent runs complete in 40–65 minutes",
      "Improvement noted after data volume optimization on July 15",
      "Log backups occurring regularly — automatic log backup enabled and active"],
     Inches(0.28), Inches(2.06), SW - Inches(0.56), Inches(1.52))

# Warning
R(sl, Inches(0.28), Inches(3.74), SW - Inches(0.56), Inches(0.94), C_ORG_BG, C_ORANGE)
R(sl, Inches(0.28), Inches(3.74), Inches(0.08), Inches(0.94), C_ORANGE)
T(sl, "⚠  Backup Storage Warning: /hana_backup/HP4 — only 3.98 GB total capacity",
  Inches(0.46), Inches(3.8), SW - Inches(0.7), Inches(0.28), sz=9.5, bold=True, color=RGBColor(0x85, 0x64, 0x04))
T(sl, "Local backup target has extremely limited space (0.03/3.98 GB used). "
      "Verify that backups are streaming to an off-site Backint/network destination. "
      "A single failed Backint run would fill the local path and cause the next backup to fail.",
  Inches(0.46), Inches(4.1), SW - Inches(0.7), Inches(0.5), sz=9, color=C_BODY)

card(sl, "Recommendations",
     ["Verify off-site backup retention meets RPO/RTO requirements (30 days minimum recommended)",
      "Monitor /hana_backup/HP4 — only 3.98 GB available; confirm Backint routing is active",
      "Review daily backup window (01:00–07:00) to avoid overlap with future maintenance windows"],
     Inches(0.28), Inches(4.84), SW - Inches(0.56), Inches(1.1))

T(sl, "Backup path: /hana_backup/HP4  |  Catalog system: M_BACKUP_CATALOG  |  Type: Backint (off-site assumed)",
  Inches(0.28), Inches(6.1), SW - Inches(0.56), Inches(0.34), sz=8, italic=True, color=C_GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ADDITIONAL RISKS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_BG)
hdr(sl, "⚠", "Additional Risks", "6 Items Identified", C_ORANGE)

risks = [
    ("License Compliance",  "License type 'undefined' on production — 40 continuous P1 alerts",        "HIGH",   C_RED,    "Set license_usage parameter (SAP Note 2779499)"),
    ("Lock Contention",     "BALDAT lock waits up to 1800s — ARFC_DEST_SHIP blocking multiple sessions","MEDIUM", C_ORANGE, "Archive app logs; review SBAL_DELETE schedule"),
    ("High Memory Queries", "BATCH_USER analytical queries: 16–30 GB memory per execution",             "MEDIUM", C_ORANGE, "Apply WorkloadClass memory limit for BATCH_USER"),
    ("Timezone Config",     "Brazil DST rule 'BRZLEA' missing — incorrect date/time calculations",      "MEDIUM", C_ORANGE, "Apply TZ update via STZAC / SAP Timezone Note"),
    ("FQM_FLOW Growth",     "Table grew +319% in 40 days (9.9M → 41.6M records)",                       "MEDIUM", C_ORANGE, "Investigate FQM background job and clearing logic"),
    ("Active CVE",          "CVE-2026-0492 active for SAP HANA — revision 89.03 may contain the fix",  "MEDIUM", C_ORANGE, "Assess CVE; plan patch upgrade to revision 89.03"),
]
tcols3 = [("Risk Area", 0.3, 2.2), ("Observation", 2.55, 4.3), ("Sev", 6.9, 1.1), ("Recommended Action", 8.05, 5.0)]
tbl_hdr(sl, tcols3, Inches(1.0))

for i, (area, obs, sev, col, action) in enumerate(risks):
    bg = C_CARD if i % 2 == 0 else C_WHITE
    ry = Inches(1.36 + i * 0.82); rh = Inches(0.8)
    R(sl, Inches(0.28), ry, SW - Inches(0.56), rh, bg, C_EDGE)
    R(sl, Inches(0.28), ry, Inches(0.07), rh, col)
    T(sl, area,   Inches(0.4),  ry + Inches(0.06), Inches(2.1),  rh - Inches(0.1), sz=9,   bold=True,  color=C_DARK)
    T(sl, obs,    Inches(2.55), ry + Inches(0.06), Inches(4.25), rh - Inches(0.1), sz=8.5, bold=False, color=C_BODY)
    pill(sl, sev, Inches(6.95), ry + Inches(0.24), Inches(1.0), Inches(0.3), col)
    T(sl, action, Inches(8.1),  ry + Inches(0.06), Inches(4.9),  rh - Inches(0.1), sz=8.5, bold=False, color=C_BODY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — MINICHECKS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_BG)
hdr(sl, "🔬", "Minichecks Analysis", "7 Scripts  |  Key Findings", C_BLUE)

minichecks = [
    ("📈 CPU & Memory",       "Post-restart (Aug 9): physical RAM at 85.7% (431/504 GB). CPU very healthy at <1%. Monitor RAM trend closely.", C_ORANGE),
    ("🔍 Expensive SQL",      "BATCH_USER: 28–30 GB memory/exec. GCBADM: BALDAT lock timeout 1800s. REPOLOAD slow queries from FABRICIO.N.",   C_RED),
    ("📋 Largest Tables",     "SOFFCONT1: 405 GB LOB (dominant). DBTABLOG: 44 GB growing. RSAU_BUF_DATA: 40 GB. REPOLOAD: 39 GB.",             C_ORANGE),
    ("⚙ Configuration",       "License = undefined. Disk at 82.7%. 297 NSE tables. 31.1M packed LOBs. Self-signed certificates active.",       C_ORANGE),
    ("💽 Disk Partitions",    "Data partition 988/1195 GB (82.7%). /hana/data/HP4 at 71.1% FS. Backup path only 3.98 GB total available.",      C_ORANGE),
    ("📡 IO Disk Details",    "Device sdc: high latency 12.5ms at very low IO (0.07 IO/s). All other disks normal (3–5ms latency).",             C_ORANGE),
    ("💾 Savepoints",         "All healthy. Critical phase consistently < 0.1s. Throughput 400–1000 MB/s. Zero blocked savepoints.",            C_GREEN),
]
cw3 = Inches(5.9); ch3 = Inches(0.84)
for i, (name, findings, col) in enumerate(minichecks):
    cx = Inches(0.28) if i % 2 == 0 else Inches(6.88)
    if i == 6: cx = Inches(3.58)
    cy = Inches(1.0) + (i // 2) * Inches(0.94)
    R(sl, cx, cy, cw3, ch3, C_WHITE, col)
    R(sl, cx, cy, Inches(0.08), ch3, col)
    T(sl, name, cx + Inches(0.16), cy + Inches(0.07), cw3 - Inches(0.26), Inches(0.28), sz=9.5, bold=True, color=C_DARK)
    T(sl, findings, cx + Inches(0.16), cy + Inches(0.38), cw3 - Inches(0.26), Inches(0.42), sz=8, color=C_BODY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ACTION PLAN
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, C_BG)
hdr(sl, "✅", "Consolidated Action Plan", "12 Actions  |  3 Priority Tiers", C_BLUE)

# P1 section
R(sl, Inches(0.28), Inches(1.0), SW - Inches(0.56), Inches(0.36), C_RED)
T(sl, "🔴  Priority 1 — Immediate (≤ 48 Hours)",
  Inches(0.42), Inches(1.03), Inches(12.5), Inches(0.27), sz=10, bold=True, color=C_WHITE)

p1 = [
    ("1. Set Production License Type",       "Owner: DBA  ·  30 min  ·  Clears all 40 P1 alerts"),
    ("2. Fix gc_unused_memory_threshold_abs", "Owner: DBA  ·  15 min  ·  Prevents elevated memory footprint"),
]
for i, (t, d) in enumerate(p1):
    rx = Inches(0.28) + i * Inches(6.55)
    R(sl, rx, Inches(1.36), Inches(6.27), Inches(0.66), C_RED_BG, C_RED)
    R(sl, rx, Inches(1.36), Inches(0.07), Inches(0.66), C_RED)
    T(sl, t, rx + Inches(0.16), Inches(1.4),  Inches(6.0), Inches(0.28), sz=9, bold=True, color=C_DARK)
    T(sl, d, rx + Inches(0.16), Inches(1.68), Inches(6.0), Inches(0.28), sz=8.5, color=C_GREY)

# P2 section
R(sl, Inches(0.28), Inches(2.14), SW - Inches(0.56), Inches(0.34), C_ORANGE)
T(sl, "🟡  Priority 2 — This Week",
  Inches(0.42), Inches(2.17), Inches(12.5), Inches(0.25), sz=10, bold=True, color=C_WHITE)

p2 = [
    ("3. Resolve BALDAT Lock Contention",   "AppTeam+DBA  ·  2–4h  ·  Prevents lock timeouts"),
    ("4. Archive REPOLOAD Table (36 GB)",   "Basis  ·  2–4h  ·  Frees disk, fixes slow queries"),
    ("5. WorkloadClass for BATCH_USER",     "DBA  ·  1–2h  ·  Caps query memory at 20 GB"),
    ("6. Fix Brazil Timezone (BRZLEA)",     "Basis  ·  1h  ·  Prevents date calc errors"),
    ("7. Patch to Rev 89.03 / CVE-0492",   "Security+Basis  ·  MW  ·  CVE remediation"),
]
aw = Inches(4.1); ag = Inches(0.19)
for i, (t, d) in enumerate(p2):
    rx = Inches(0.28) + (i % 3) * (aw + ag)
    ry = Inches(2.56) + (i // 3) * Inches(0.72)
    R(sl, rx, ry, aw, Inches(0.66), C_ORG_BG, C_ORANGE)
    R(sl, rx, ry, Inches(0.07), Inches(0.66), C_ORANGE)
    T(sl, t, rx + Inches(0.16), ry + Inches(0.05), aw - Inches(0.24), Inches(0.28), sz=8.5, bold=True, color=C_DARK)
    T(sl, d, rx + Inches(0.16), ry + Inches(0.34), aw - Inches(0.24), Inches(0.26), sz=7.5, color=C_GREY)

# P3 section
R(sl, Inches(0.28), Inches(4.14), SW - Inches(0.56), Inches(0.34), C_GREEN)
T(sl, "🟢  Priority 3 — This Month",
  Inches(0.42), Inches(4.17), Inches(12.5), Inches(0.25), sz=10, bold=True, color=C_WHITE)

p3 = [
    ("8. Archive REPOSRC, SOFFCONT1, DBTABLOG", "AppTeam+Basis  ·  1–2d  ·  Prevents disk exhaustion"),
    ("9. Investigate FQM_FLOW Growth (+319%)",  "AppTeam  ·  2–4h  ·  Prevents uncontrolled growth"),
    ("10. Optimize BATCH_USER SQL Queries",     "AppTeam+DBA  ·  1–2w  ·  Reduces memory pressure"),
    ("11. Plan Data Volume Expansion (82.7%)",  "Basis+Infra  ·  MW  ·  Prevents disk-full outage"),
    ("12. Review Physical RAM Allocation",      "DBA+Infra  ·  2–4h  ·  Prevents OS memory pressure"),
]
for i, (t, d) in enumerate(p3):
    rx = Inches(0.28) + (i % 3) * (aw + ag)
    ry = Inches(4.55) + (i // 3) * Inches(0.72)
    R(sl, rx, ry, aw, Inches(0.66), C_GRN_BG, C_GREEN)
    R(sl, rx, ry, Inches(0.07), Inches(0.66), C_GREEN)
    T(sl, t, rx + Inches(0.16), ry + Inches(0.05), aw - Inches(0.24), Inches(0.28), sz=8.5, bold=True, color=C_DARK)
    T(sl, d, rx + Inches(0.16), ry + Inches(0.34), aw - Inches(0.24), Inches(0.26), sz=7.5, color=C_GREY)


# ── Save ───────────────────────────────────────────────────────────────────
OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT))
print(f"Done - {len(prs.slides)} slides - {OUTPUT}")
