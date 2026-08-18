"""
360 Report Generator
Standalone GUI app — select a PLA Excel file and generate the 360 PPTX.
Self-installs any missing dependencies on first run.
Requires: 360_dummy.pptx template in the same folder as this script.
"""
import sys
import subprocess

# ── Auto-install dependencies ─────────────────────────────────────────────────
_PACKAGES = ["openpyxl", "python-pptx", "Pillow", "pywin32", "lxml"]
_MOD_MAP   = {"python-pptx": "pptx", "Pillow": "PIL", "pywin32": "win32com"}

def _ensure_deps():
    import importlib
    import importlib.util
    missing = []
    for p in _PACKAGES:
        mod = _MOD_MAP.get(p, p)
        try:
            importlib.import_module(mod)
        except ImportError:
            missing.append(p)
    if missing:
        print(f"Installing missing packages: {missing}")
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", "--quiet"] + missing)

_ensure_deps()

# ── Standard imports ──────────────────────────────────────────────────────────
import os, io, zipfile, datetime, shutil, threading, traceback
import xml.etree.ElementTree as ET
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

import openpyxl
from lxml import etree
from PIL import Image, ImageGrab
from pptx import Presentation
from pptx.util import Inches
import win32com.client
import time

# ── Constants ─────────────────────────────────────────────────────────────────
SCRIPT_DIR    = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PPTX = os.path.join(SCRIPT_DIR, "360_dummy.pptx")

NS_DML = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_R   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS_X   = 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'


# ═══════════════════════════════════════════════════════════════════════════════
#  XLSX helpers
# ═══════════════════════════════════════════════════════════════════════════════

def _find_label_value(ws, label):
    """Return the cell value immediately to the right of the first cell whose
    text contains *label* (case-insensitive)."""
    needle = label.lower().rstrip(':')
    for row in ws.iter_rows(values_only=True):
        for idx, val in enumerate(row):
            if val and needle in str(val).lower():
                nxt = idx + 1
                if nxt < len(row) and row[nxt] is not None:
                    return str(row[nxt]).strip()
    return ""


def _find_label_value_rows(rows, label):
    """Same as _find_label_value but operates on a pre-fetched list of rows."""
    needle = label.lower().rstrip(':')
    for row in rows:
        for idx, val in enumerate(row):
            if val and needle in str(val).lower():
                nxt = idx + 1
                if nxt < len(row) and row[nxt]:
                    return str(row[nxt]).strip()
    return ""


def _sheet_rels_path(xlsx_path, sheet_name):
    """Dynamically resolve xl/worksheets/_rels/sheetN.xml.rels for a named sheet,
    regardless of how the workbook numbers its sheets."""
    with zipfile.ZipFile(xlsx_path, 'r') as z:
        wb_root  = ET.fromstring(z.read('xl/workbook.xml').decode('utf-8'))
        rel_root = ET.fromstring(z.read('xl/_rels/workbook.xml.rels').decode('utf-8'))
        rid_map  = {r.get('Id'): r.get('Target') for r in rel_root}
        for s in wb_root.iter(f'{{{NS_X}}}sheet'):
            if s.get('name') == sheet_name:
                target  = rid_map.get(s.get(f'{{{NS_R}}}id'), '')
                sh_stem = target.split('/')[-1].replace('.xml', '')
                return f'xl/worksheets/_rels/{sh_stem}.xml.rels'
    return None


def _extract_sheet_image(xlsx_path, sheet_name):
    """Extract the first embedded PNG image from a named sheet's drawing."""
    rels_path = _sheet_rels_path(xlsx_path, sheet_name)
    if not rels_path:
        return None
    with zipfile.ZipFile(xlsx_path, 'r') as z:
        try:
            rels_xml = z.read(rels_path).decode('utf-8')
        except KeyError:
            return None
        rels_root    = etree.fromstring(rels_xml.encode())
        drawing_path = None
        for rel in rels_root:
            if 'drawing' in rel.get('Target', ''):
                drawing_path = rel.get('Target').replace('..', 'xl')
                break
        if not drawing_path:
            return None
        draw_file = drawing_path.split('/')[-1]
        try:
            dr_rels_xml = z.read(f'xl/drawings/_rels/{draw_file}.rels').decode('utf-8')
        except KeyError:
            return None
        dr_root   = etree.fromstring(dr_rels_xml.encode())
        img_path  = None
        for rel in dr_root:
            t = rel.get('Target', '')
            if 'media/image' in t:
                img_path = t.replace('../', 'xl/')
                break
        if not img_path:
            return None
        return Image.open(io.BytesIO(z.read(img_path))).copy()


# ═══════════════════════════════════════════════════════════════════════════════
#  PPTX helpers
# ═══════════════════════════════════════════════════════════════════════════════

def _set_cell_text(cell_obj, txt, sz=None):
    """Replace all text runs in a table cell's first paragraph with *txt*,
    preserving the original run's character formatting (rPr).
    Optionally override font size with *sz* (in hundredths of a point, e.g. 1000 = 10pt)."""
    tf   = cell_obj.text_frame
    para = tf.paragraphs[0]
    p_el = para._p
    runs = p_el.findall(f'{{{NS_DML}}}r')
    if runs:
        rPr_src  = runs[0].find(f'{{{NS_DML}}}rPr')
        rPr_copy = etree.fromstring(etree.tostring(rPr_src)) if rPr_src is not None else None
        for r in runs:
            p_el.remove(r)
        run = para.add_run()
        run.text = txt
        if rPr_copy is not None:
            r_el     = run._r
            existing = r_el.find(f'{{{NS_DML}}}rPr')
            if existing is not None:
                r_el.remove(existing)
            if sz is not None:
                rPr_copy.set('sz', str(sz))
            r_el.insert(0, rPr_copy)
        elif sz is not None:
            r_el = run._r
            rPr  = etree.SubElement(r_el, f'{{{NS_DML}}}rPr')
            rPr.set('sz', str(sz))
            r_el.insert(0, rPr)
    else:
        run = para.add_run()
        run.text = txt
        if sz is not None:
            r_el = run._r
            rPr  = etree.SubElement(r_el, f'{{{NS_DML}}}rPr')
            rPr.set('sz', str(sz))
            r_el.insert(0, rPr)


def _update_label_value_para(para, new_value):
    """In a paragraph structured as [label_run, value_run], replace the value run.
    Falls back gracefully for single-run or no-run paragraphs."""
    runs = para.runs
    if len(runs) >= 2:
        runs[1].text = '\xa0' + new_value
    elif len(runs) == 1:
        text  = runs[0].text
        label = text.split('\xa0')[0] if '\xa0' in text else text
        runs[0].text = label + '\xa0' + new_value


def _remove_hyperlinks(shape):
    """Strip all hlinkClick hyperlink elements from every run in a shape."""
    if not hasattr(shape, 'text_frame'):
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(f'{{{NS_DML}}}rPr')
            if rPr is not None:
                for hl in rPr.findall(f'{{{NS_DML}}}hlinkClick'):
                    rPr.remove(hl)
        for endRPr in para._p.findall(f'{{{NS_DML}}}endParaRPr'):
            for hl in endRPr.findall(f'{{{NS_DML}}}hlinkClick'):
                endRPr.remove(hl)


def _rename_shape_title(shape, new_title):
    """Replace a (possibly multi-paragraph) title shape with a single clean line."""
    tf    = shape.text_frame
    txBody = tf._txBody
    # Clear all runs in all paragraphs
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    # Write new title into first paragraph
    if tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = new_title
    else:
        tf.paragraphs[0].add_run().text = new_title
    # Remove any extra paragraphs (e.g. old line-break continuation)
    all_paras = txBody.findall(f'{{{NS_DML}}}p')
    for extra in all_paras[1:]:
        txBody.remove(extra)


def _delete_slide(prs, idx):
    """Remove slide at 0-based *idx* from the presentation."""
    sldIdLst = prs.slides._sldIdLst
    sldId    = sldIdLst[idx]
    rId      = sldId.get(f'{{{NS_R}}}id')
    sldIdLst.remove(sldId)
    prs.part._rels.pop(rId)


def _remove_shapes(slide, types=(19, 13)):
    """Remove all shapes of given shape_type codes from *slide*."""
    for el in [s._element for s in slide.shapes if s.shape_type in types]:
        el.getparent().remove(el)


def _fit_image(img_path, avail_w, avail_h):
    """Return (w, h) in EMU that fits the image within the available area."""
    img = Image.open(img_path)
    iw, ih = img.size
    asp = iw / ih
    if avail_w / asp <= avail_h:
        return avail_w, int(avail_w / asp)
    return int(avail_h * asp), avail_h


def _place_image(slide, img_path, left, top, avail_w, avail_h):
    """Insert *img_path* into *slide* centred within the given area."""
    if not os.path.exists(img_path):
        return
    w, h  = _fit_image(img_path, avail_w, avail_h)
    slide.shapes.add_picture(img_path,
                             left + (avail_w - w) // 2,
                             top  + (avail_h - h) // 2,
                             w, h)


def _insert_body_screenshot(slide, img_path, slide_w, slide_h,
                             top_in=1.6, pad_in=0.15, bottom_margin_in=0.3):
    """Insert *img_path* horizontally centred below *top_in* inches."""
    if not os.path.exists(img_path):
        return
    pad      = Inches(pad_in)
    avail_w  = slide_w - 2 * pad
    avail_h  = slide_h - Inches(top_in) - Inches(bottom_margin_in)
    w, h     = _fit_image(img_path, avail_w, avail_h)
    slide.shapes.add_picture(img_path,
                             (slide_w - w) // 2,
                             Inches(top_in),
                             w, h)


# ═══════════════════════════════════════════════════════════════════════════════
#  Main generation routine
# ═══════════════════════════════════════════════════════════════════════════════

def generate_360(xlsx_path, log_cb):
    """Generate the 360 PPTX from *xlsx_path*.
    *log_cb(msg)* is called with progress strings throughout."""

    if not os.path.exists(TEMPLATE_PPTX):
        raise FileNotFoundError(
            f"Template not found: {TEMPLATE_PPTX}\n"
            "Place 360_dummy.pptx in the same folder as this script.")

    work_dir  = os.path.dirname(xlsx_path)
    xlsx_copy = os.path.join(work_dir, "_pla_work_copy.xlsx")
    img_dir   = os.path.join(work_dir, "_imgs_tmp")
    os.makedirs(img_dir, exist_ok=True)

    # ── 1. Copy xlsx (avoids lock conflicts) ──────────────────────────────────
    log_cb("Copying Excel file...")
    if os.path.exists(xlsx_copy):
        os.remove(xlsx_copy)
    shutil.copy2(xlsx_path, xlsx_copy)

    # ── 2. Read xlsx data ─────────────────────────────────────────────────────
    log_cb("Reading data from Excel...")
    wb = openpyxl.load_workbook(xlsx_copy, data_only=True)

    # Customer Info
    ci            = wb["Customer Info."]
    customer_name = _find_label_value(ci, "Customer Name")
    crm_id        = _find_label_value(ci, "CRM ID")
    erp_id        = _find_label_value(ci, "ERP ID")
    cdm_name      = _find_label_value(ci, "CDM Name")
    tsm_name      = _find_label_value(ci, "TSM Name")
    subcontractor = _find_label_value(ci, "Subcontractor")

    systems_in_scope = []
    in_scope = False
    for row in ci.iter_rows(values_only=True):
        if row[3] and "Ext. SID" in str(row[3]):
            in_scope = True
            continue
        if in_scope and row[3] and str(row[3]).strip() and '>' not in str(row[3]):
            systems_in_scope.append({
                "ext_sid":    str(row[3]).strip(),
                "sys_num":    str(row[4]).strip() if row[4] else "",
                "sys_type":   str(row[5]).strip() if row[5] else "",
                "sys_role":   str(row[6]).strip() if row[6] else "",
                "datacenter": str(row[7]).strip() if row[7] else "",
                "eu_dp":      str(row[8]).strip() if row[8] else "",
                "gxp":        str(row[9]).strip() if row[9] else "",
            })

    # System info
    si_rows   = list(wb["System info."].iter_rows(values_only=True))
    sys_type  = _find_label_value_rows(si_rows, "System Type")
    sys_role  = _find_label_value_rows(si_rows, "System Role")
    db_type   = _find_label_value_rows(si_rows, "Database Type")
    db_name   = _find_label_value_rows(si_rows, "Database Name")
    net_seg   = _find_label_value_rows(si_rows, "Network Segment ID")
    main_fqdn = _find_label_value_rows(si_rows, "Main Appl Host FQDN")

    # BASIS Checks
    bc_rows       = list(wb["BASIS Checks"].iter_rows(values_only=True))
    basis_sid     = str(bc_rows[0][0]).strip() if bc_rows[0][0] else ""
    basis_product = basis_version = basis_db_type = basis_os_type = basis_ewa = ""
    for row in bc_rows[:8]:
        if not row[1]:
            continue
        lbl = str(row[1]).strip().rstrip(':').lower()
        val = str(row[2]).strip() if row[2] else ""
        if   "lead product" in lbl: basis_product = val
        elif "version"      in lbl: basis_version = val
        elif "db type"      in lbl: basis_db_type = val
        elif "os type"      in lbl: basis_os_type = val
        elif "ewa"          in lbl: basis_ewa      = val

    basis_checks = []
    for row in bc_rows:
        if not row[3]:
            continue
        lbl = str(row[3]).strip()
        if lbl.lower() == "checks:":
            continue
        basis_checks.append((
            lbl,
            str(row[4]).strip() if row[4] else "",
            str(row[5]).strip() if row[5] else "",
        ))

    # Recommendations
    rec_rows         = list(wb["Recommendations"].iter_rows(values_only=True))
    rec_tbl_starts   = [i for i, r in enumerate(rec_rows)
                        if r[1] and "Item No" in str(r[1])]

    def _parse_rec(start):
        items = []
        for row in rec_rows[start + 1:]:
            if not any(c for c in row if c is not None):
                break
            items.append(tuple(str(row[j]).strip() if row[j] else ""
                               for j in range(1, 9)))
        return items

    basis_recs = _parse_rec(rec_tbl_starts[0]) if len(rec_tbl_starts) > 0 else []
    db_recs    = _parse_rec(rec_tbl_starts[1]) if len(rec_tbl_starts) > 1 else []
    sm_recs    = _parse_rec(rec_tbl_starts[2]) if len(rec_tbl_starts) > 2 else []

    log_cb(f"Customer : {customer_name}")
    log_cb(f"BASIS    : {len(basis_checks)} checks | "
           f"Recs BASIS/DB/SM: {len(basis_recs)}/{len(db_recs)}/{len(sm_recs)}")

    # ── 3. Extract embedded sheet images (dynamic sheet-name lookup) ──────────
    log_cb("Extracting sheet images...")
    db_img_path = os.path.join(img_dir, "db_checks.png")
    sm_img_path = os.path.join(img_dir, "sm_checks.png")

    for sheet_name, dest in [("DB Checks", db_img_path), ("SM Checks", sm_img_path)]:
        img = _extract_sheet_image(xlsx_copy, sheet_name)
        if img:
            img.save(dest)
            log_cb(f"  {sheet_name} image: {img.size}")
        else:
            log_cb(f"  WARNING: no image found in '{sheet_name}' sheet")

    # ── 4. Screenshot Recommendations tables via Excel COM ────────────────────
    def _xr(s, e):
        return f"B{s+1}:I{e+1}"

    basis_range = _xr(rec_tbl_starts[0], rec_tbl_starts[0] + len(basis_recs)) if rec_tbl_starts       else "B3:I11"
    db_range    = _xr(rec_tbl_starts[1], rec_tbl_starts[1] + len(db_recs))    if len(rec_tbl_starts)>1 else "B15:I20"
    sm_range    = _xr(rec_tbl_starts[2], rec_tbl_starts[2] + len(sm_recs))    if len(rec_tbl_starts)>2 else "B23:I26"

    basis_rec_img = os.path.join(img_dir, "rec_basis.png")
    db_rec_img    = os.path.join(img_dir, "rec_db.png")
    sm_rec_img    = os.path.join(img_dir, "rec_sm.png")

    log_cb("Opening Excel for Recommendations screenshots...")
    xl     = win32com.client.Dispatch("Excel.Application")
    xl.Visible = True
    xl.DisplayAlerts = False
    time.sleep(1.5)
    wb_xl  = xl.Workbooks.Open(os.path.abspath(xlsx_copy))
    ws_rec = wb_xl.Sheets("Recommendations")
    ws_rec.Activate()

    for rng, dest, label in [
        (basis_range, basis_rec_img, "BASIS"),
        (db_range,    db_rec_img,    "Database"),
        (sm_range,    sm_rec_img,    "SM"),
    ]:
        ws_rec.Range(rng).CopyPicture(Appearance=1, Format=2)
        time.sleep(1.5)
        img = ImageGrab.grabclipboard()
        if img:
            img.save(dest)
            log_cb(f"  {label} recommendations screenshot: {img.size}")
        else:
            log_cb(f"  WARNING: clipboard empty for {label} ({rng})")

    wb_xl.Close(False)
    xl.Quit()
    time.sleep(1.0)
    log_cb("Excel closed.")

    # ── 5. Build PPTX ─────────────────────────────────────────────────────────
    log_cb("Building PPTX from template...")
    now        = datetime.datetime.now()
    month_year = now.strftime("%B %Y")
    prs        = Presentation(TEMPLATE_PPTX)
    SW         = prs.slide_width
    SH         = prs.slide_height

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    for shape in prs.slides[0].shapes:
        if "Title" in shape.name and hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Innova Sport" in run.text:
                        run.text = run.text.replace("Innova Sport S.A De C.V.", customer_name)
        if "Date" in shape.name and hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = month_year

    # ── Slide 3: Customer Card Overview ──────────────────────────────────────
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.name == "TextBox 9" and hasattr(shape, 'text_frame'):
            # Rebuild TextBox 9 with proper inline label+value runs at size 16
            # Structure: para0 = CRM+CDM+TSM, para1 = SUBCONTRACTOR, para2 = ERP ID
            tf  = shape.text_frame
            ns  = NS_DML

            def _make_run(p_el, text, bold):
                """Append a new run to p_el, copying rPr from first existing run,
                then override sz and bold."""
                existing = p_el.findall(f'{{{ns}}}r')
                base_rPr = None
                if existing:
                    base_rPr = existing[0].find(f'{{{ns}}}rPr')

                r_el  = etree.SubElement(p_el, f'{{{ns}}}r')
                rPr   = etree.SubElement(r_el, f'{{{ns}}}rPr')
                # Copy all attribs from base if available
                if base_rPr is not None:
                    for k, v in base_rPr.attrib.items():
                        rPr.set(k, v)
                    for child in base_rPr:
                        rPr.append(etree.fromstring(etree.tostring(child)))
                rPr.set('sz', '1600')         # 16pt  (hundredths of a point)
                rPr.set('b',  '1' if bold else '0')
                t_el      = etree.SubElement(r_el, f'{{{ns}}}t')
                t_el.text = text
                return r_el

            # Wipe all existing runs from all 3 paragraphs
            for para in tf.paragraphs:
                p_el = para._p
                for r in p_el.findall(f'{{{ns}}}r'):
                    p_el.remove(r)

            # Para 0: CRM ID  CDM  TSM
            p0 = tf.paragraphs[0]._p
            _make_run(p0, 'CRM ID:',        bold=True)
            _make_run(p0, f'\xa0{crm_id}\xa0\t', bold=False)
            _make_run(p0, 'CDM:',           bold=True)
            _make_run(p0, f'\xa0{cdm_name}\xa0', bold=False)
            _make_run(p0, 'TSM:',           bold=True)
            _make_run(p0, f'\xa0{tsm_name}\xa0', bold=False)

            # Para 1: SUBCONTRACTOR
            p1 = tf.paragraphs[1]._p
            _make_run(p1, 'SUBCONTRACTOR:', bold=True)
            _make_run(p1, f'\xa0{subcontractor or "N/A"}\xa0', bold=False)

            # Para 2: ERP ID
            p2 = tf.paragraphs[2]._p
            _make_run(p2, 'ERP ID:',        bold=True)
            _make_run(p2, f'\xa0{erp_id}',  bold=False)

        if shape.name == "TextBox 20" and hasattr(shape, 'text_frame'):
            # index → value (row 5 = MAIN APPL HOST FQDN)
            val_map = {0: sys_type, 1: db_type, 2: sys_role,
                       3: db_name,  4: net_seg,  5: main_fqdn}
            for i, para in enumerate(shape.text_frame.paragraphs):
                if i in val_map:
                    _update_label_value_para(para, val_map[i])

    # ── Slide 6: BASIS Checks ─────────────────────────────────────────────────
    for shape in prs.slides[5].shapes:
        if shape.shape_type != 19:
            continue
        tbl   = shape.table
        nrows = len(tbl.rows)
        if nrows == 6:
            # System header block — 10pt
            _set_cell_text(tbl.rows[0].cells[0], basis_sid, sz=1000)
            for r, val in enumerate(
                [basis_product, basis_version, basis_db_type, basis_os_type, basis_ewa], 1
            ):
                if r < nrows:
                    _set_cell_text(tbl.rows[r].cells[2], val, sz=1000)
        elif nrows >= 4:
            # Checks table — 10pt
            for idx, (lbl, result, details) in enumerate(basis_checks):
                ri = idx + 1
                if ri >= nrows:
                    break
                _set_cell_text(tbl.rows[ri].cells[0], lbl,     sz=1000)
                _set_cell_text(tbl.rows[ri].cells[1], result,  sz=1000)
                _set_cell_text(tbl.rows[ri].cells[2], details, sz=1000)

    # ── Slide 7: DB Checks ────────────────────────────────────────────────────
    slide7 = prs.slides[6]

    # Update system-ID text box (TextBox 5)
    for shape in slide7.shapes:
        if shape.name == "TextBox 5" and hasattr(shape, 'text_frame'):
            paras = shape.text_frame.paragraphs
            for para in paras:
                for run in para.runs:
                    run.text = ""
            for para in reversed(paras):
                if para.runs:
                    para.runs[0].text = basis_sid
                    break

    # Update info text box (TextBox 3)
    for shape in slide7.shapes:
        if shape.name == "TextBox 3" and hasattr(shape, 'text_frame'):
            subs = [
                ("SAP HANA, platform edition",              basis_product),
                ("2.0 / 079 / 08 (2.007908)",               basis_version),
                ("Linux",                                    basis_os_type),
                ("HANA MDC (Multi-Tenant Database Container)", db_type),
                ("HANA MDC",                                 db_type),
            ]
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in subs:
                        if old in run.text:
                            run.text = run.text.replace(old, new)

    # Replace DB Checks screenshot — fill the area below the divider line
    if os.path.exists(db_img_path):
        pic = next((s for s in slide7.shapes if s.shape_type == 13), None)
        if pic:
            pic._element.getparent().remove(pic._element)
        area_top  = Inches(2.68)
        area_left = Inches(0.3)
        avail_w   = SW - 2 * Inches(0.3)
        avail_h   = Inches(7.0) - area_top
        _place_image(slide7, db_img_path, area_left, area_top, avail_w, avail_h)

    # ── Slide 8: SM Checks ────────────────────────────────────────────────────
    # Template has:
    #   Table 7  (6x3, top=1.26") = system header  → fill from BASIS data
    #   Table 2  (16x3, top=2.85") = checks body   → replace with SM image
    slide8 = prs.slides[7]

    # Fill system header table (same data as BASIS header)
    for shape in slide8.shapes:
        if shape.shape_type == 19 and len(shape.table.rows) == 6:
            tbl = shape.table
            _set_cell_text(tbl.rows[0].cells[0], basis_sid)
            for r, val in enumerate(
                [basis_product, basis_version, basis_db_type, basis_os_type, basis_ewa], 1
            ):
                if r < len(tbl.rows):
                    _set_cell_text(tbl.rows[r].cells[2], val)

    # Remove only the body checks table and any existing pictures, then insert SM image
    for shape in slide8.shapes:
        if shape.shape_type == 19 and len(shape.table.rows) != 6:
            shape._element.getparent().remove(shape._element)
            break
    for shape in list(slide8.shapes):
        if shape.shape_type == 13:
            shape._element.getparent().remove(shape._element)

    if os.path.exists(sm_img_path):
        # Place image in the body area (below divider line at ~3.01")
        area_top  = Inches(3.1)
        area_left = Inches(0.3)
        avail_w   = SW - 2 * Inches(0.3)
        avail_h   = Inches(7.0) - area_top
        _place_image(slide8, sm_img_path, area_left, area_top, avail_w, avail_h)

    # Rename title shape from "Server Management\n Checks" → "SM Checks"
    for shape in slide8.shapes:
        if hasattr(shape, 'text_frame') and 'management' in shape.text.lower():
            _rename_shape_title(shape, "SM Checks")
            break

    # ── Slide 11: Recommendations – BASIS ────────────────────────────────────
    _remove_shapes(prs.slides[10])
    _insert_body_screenshot(prs.slides[10], basis_rec_img, SW, SH, top_in=1.6)

    # ── Slide 12: Recommendations – Database ─────────────────────────────────
    _remove_shapes(prs.slides[11])
    _insert_body_screenshot(prs.slides[11], db_rec_img, SW, SH, top_in=1.4)

    # ── Slide 14: Recommendations – SM ───────────────────────────────────────
    _remove_shapes(prs.slides[13])
    _insert_body_screenshot(prs.slides[13], sm_rec_img, SW, SH, top_in=1.6)

    # ── Slide 15: Attachments — update filename, strip hyperlinks ────────────
    for shape in prs.slides[14].shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Innova Sport" in run.text:
                        run.text = run.text.replace(
                            "PLA_Innova Sport S.A De C.V.xlsx",
                            os.path.basename(xlsx_path))
            _remove_hyperlinks(shape)

    # ── Remove duplicate slides (highest index first to avoid index shift) ────
    log_cb("Removing duplicate slides...")
    _delete_slide(prs, 12)   # Recommendations – Database duplicate (slide 13)
    _delete_slide(prs,  8)   # SM Checks duplicate (slide 9)
    log_cb(f"  Slides remaining: {len(prs.slides)}")

    # ── Save ──────────────────────────────────────────────────────────────────
    safe_name = customer_name.replace("/", "-").replace("\\", "-").strip()
    out_path  = os.path.join(work_dir, f"360_{safe_name}.pptx")
    prs.save(out_path)

    # Cleanup temp artefacts
    try:
        os.remove(xlsx_copy)
        shutil.rmtree(img_dir, ignore_errors=True)
    except Exception:
        pass

    log_cb(f"\nDone!  Saved: {out_path}")
    return out_path


# ═══════════════════════════════════════════════════════════════════════════════
#  GUI
# ═══════════════════════════════════════════════════════════════════════════════

class App(tk.Tk):

    # ── Colours ───────────────────────────────────────────────────────────────
    BG_DARK  = "#0f1e2d"
    BG_MID   = "#1a2b3c"
    BG_ENTRY = "#243447"
    FG_WHITE = "#ffffff"
    FG_MUTED = "#8faacc"
    FG_LABEL = "#ccd9e8"
    BTN_BLUE = "#2e6da4"
    BTN_GRN  = "#1e7a3e"
    BTN_GREY = "#5a5a5a"

    def __init__(self):
        super().__init__()
        self.title("360 Report Generator")
        self.resizable(False, False)
        self.configure(bg=self.BG_MID)
        self._build_ui()

    def _build_ui(self):
        # Header
        hdr = tk.Frame(self, bg=self.BG_DARK, pady=14)
        hdr.pack(fill="x")
        tk.Label(hdr, text="360° Report Generator",
                 font=("Segoe UI", 18, "bold"),
                 fg=self.FG_WHITE, bg=self.BG_DARK).pack()
        tk.Label(hdr, text="Select a PLA Excel file to generate the PowerPoint report",
                 font=("Segoe UI", 10),
                 fg=self.FG_MUTED, bg=self.BG_DARK).pack(pady=(2, 0))

        # File picker row
        frm = tk.Frame(self, bg=self.BG_MID, padx=20, pady=16)
        frm.pack(fill="x")
        tk.Label(frm, text="PLA Excel file:",
                 font=("Segoe UI", 10, "bold"),
                 fg=self.FG_LABEL, bg=self.BG_MID).grid(row=0, column=0, sticky="w")
        self.xlsx_var = tk.StringVar()
        tk.Entry(frm, textvariable=self.xlsx_var, width=52,
                 font=("Segoe UI", 10),
                 bg=self.BG_ENTRY, fg=self.FG_WHITE,
                 insertbackground=self.FG_WHITE,
                 relief="flat", bd=4
                 ).grid(row=1, column=0, padx=(0, 8), pady=(4, 0), sticky="ew")
        tk.Button(frm, text="Browse…",
                  font=("Segoe UI", 10),
                  bg=self.BTN_BLUE, fg=self.FG_WHITE,
                  activebackground="#3a85c9", activeforeground=self.FG_WHITE,
                  relief="flat", padx=10, pady=4,
                  command=self._browse
                  ).grid(row=1, column=1, pady=(4, 0))
        frm.columnconfigure(0, weight=1)

        # Generate button + Open button on same row
        btn_frm = tk.Frame(self, bg=self.BG_MID, pady=8)
        btn_frm.pack()
        self.btn_gen = tk.Button(btn_frm,
                                 text="  Generate 360 PPTX  ",
                                 font=("Segoe UI", 12, "bold"),
                                 bg=self.BTN_GRN, fg=self.FG_WHITE,
                                 activebackground="#25a050",
                                 activeforeground=self.FG_WHITE,
                                 relief="flat", padx=16, pady=8,
                                 command=self._start_generate)
        self.btn_gen.pack(side="left", padx=(0, 8))

        self.btn_open = tk.Button(btn_frm,
                                  text="  Open PPTX  ",
                                  font=("Segoe UI", 12, "bold"),
                                  bg=self.BTN_BLUE, fg=self.FG_WHITE,
                                  activebackground="#3a85c9",
                                  activeforeground=self.FG_WHITE,
                                  relief="flat", padx=16, pady=8,
                                  command=self._open_pptx)
        # hidden until a file is generated
        self._last_out = None

        # Progress bar
        pb_frm = tk.Frame(self, bg=self.BG_MID, padx=20)
        pb_frm.pack(fill="x")
        self.progress = ttk.Progressbar(pb_frm, mode="indeterminate", length=540)
        self.progress.pack(fill="x")

        # Log area
        log_frm = tk.Frame(self, bg=self.BG_MID, padx=20, pady=10)
        log_frm.pack(fill="both", expand=True)
        tk.Label(log_frm, text="Log:",
                 font=("Segoe UI", 9, "bold"),
                 fg=self.FG_MUTED, bg=self.BG_MID).pack(anchor="w")
        self.log_text = tk.Text(log_frm, height=14, width=72,
                                font=("Consolas", 9),
                                bg=self.BG_DARK, fg="#c8d8e8",
                                relief="flat", bd=0, state="disabled")
        sb = tk.Scrollbar(log_frm, command=self.log_text.yview)
        self.log_text.configure(yscrollcommand=sb.set)
        self.log_text.pack(side="left", fill="both", expand=True)
        sb.pack(side="right", fill="y")

        # Status bar
        self.status_var = tk.StringVar(value="Ready.")
        tk.Label(self, textvariable=self.status_var,
                 font=("Segoe UI", 9),
                 fg=self.FG_MUTED, bg=self.BG_DARK,
                 anchor="w", padx=8, pady=4
                 ).pack(fill="x", side="bottom")

        self.geometry("600x530")

    # ── UI callbacks ──────────────────────────────────────────────────────────

    def _browse(self):
        path = filedialog.askopenfilename(
            title="Select PLA Excel file",
            filetypes=[("Excel files", "*.xlsx *.xlsm"), ("All files", "*.*")])
        if path:
            self.xlsx_var.set(path)

    def _log(self, msg):
        self.log_text.configure(state="normal")
        self.log_text.insert("end", msg + "\n")
        self.log_text.see("end")
        self.log_text.configure(state="disabled")
        self.status_var.set(msg[:90])

    def _start_generate(self):
        xlsx = self.xlsx_var.get().strip()
        if not xlsx or not os.path.exists(xlsx):
            messagebox.showerror("Error", "Please select a valid Excel file.")
            return
        if not os.path.exists(TEMPLATE_PPTX):
            messagebox.showerror(
                "Template missing",
                f"360_dummy.pptx not found at:\n{TEMPLATE_PPTX}\n\n"
                "Place the template in the same folder as this script.")
            return

        self.btn_gen.configure(state="disabled", bg=self.BTN_GREY)
        self.progress.start(12)
        self._log(f"Starting — {os.path.basename(xlsx)}\n")

        def _worker():
            try:
                out = generate_360(xlsx, lambda m: self.after(0, self._log, m))
                self.after(0, self._on_done, out)
            except Exception as e:
                self.after(0, self._on_error, str(e), traceback.format_exc())

        threading.Thread(target=_worker, daemon=True).start()

    def _on_done(self, out_path):
        self.progress.stop()
        self.btn_gen.configure(state="normal", bg=self.BTN_GRN)
        self._last_out = out_path
        self.btn_open.pack(side="left")          # show Open button
        self.status_var.set(f"Saved: {out_path}")
        messagebox.showinfo("Done",
                            f"Report generated successfully!\n\n{out_path}")

    def _open_pptx(self):
        if self._last_out and os.path.exists(self._last_out):
            os.startfile(self._last_out)

    def _on_error(self, err, tb):
        self.progress.stop()
        self.btn_gen.configure(state="normal", bg=self.BTN_GRN)
        self._log(f"\nERROR: {err}\n{tb}")
        self.status_var.set(f"Error: {err[:80]}")
        messagebox.showerror("Generation failed",
                             f"{err}\n\nSee the log panel for full details.")


# ── Entry point ───────────────────────────────────────────────────────────────
if __name__ == "__main__":
    App().mainloop()
