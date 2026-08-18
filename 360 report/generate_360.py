"""
Generate 360 Caretaker Engagement Review PPTX from xlsx data.
Usage: python generate_360.py
"""
import sys, os, io, zipfile, datetime, shutil
sys.stdout.reconfigure(encoding='utf-8')

import openpyxl
from lxml import etree
from PIL import Image, ImageGrab
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import win32com.client
import time

BASE_DIR      = os.path.dirname(os.path.abspath(__file__))
XLSX_PATH     = os.path.join(BASE_DIR, "Xignux S.A. De C.V.xlsx")
XLSX_COPY     = os.path.join(BASE_DIR, "Xignux_copy.xlsx")
PPTX_TEMPLATE = os.path.join(BASE_DIR, "360_dummy.pptx")
IMG_DIR       = os.path.join(BASE_DIR, "_imgs_tmp")
os.makedirs(IMG_DIR, exist_ok=True)

# ── 1. Fresh copy of xlsx ─────────────────────────────────────────────────────
if os.path.exists(XLSX_COPY):
    os.remove(XLSX_COPY)
shutil.copy2(XLSX_PATH, XLSX_COPY)

# ── 2. Read xlsx data ─────────────────────────────────────────────────────────
wb_data = openpyxl.load_workbook(XLSX_COPY, data_only=True)

def find_label_value(ws, label):
    """Return cell value immediately right of the first cell whose text contains label."""
    for row in ws.iter_rows(values_only=True):
        for idx, val in enumerate(row):
            if val and label.lower().rstrip(':') in str(val).lower():
                nxt = idx + 1
                if nxt < len(row) and row[nxt] is not None:
                    return str(row[nxt]).strip()
    return ""

# --- Customer Info ---
ci            = wb_data["Customer Info."]
customer_name = find_label_value(ci, "Customer Name")
crm_id        = find_label_value(ci, "CRM ID")
erp_id        = find_label_value(ci, "ERP ID")
cdm_name      = find_label_value(ci, "CDM Name")
tsm_name      = find_label_value(ci, "TSM Name")
subcontractor = find_label_value(ci, "Subcontractor")

# Systems in scope (rows after the "Ext. SID" header row)
systems_in_scope = []
in_scope = False
for row in ci.iter_rows(values_only=True):
    if row[3] and "Ext. SID" in str(row[3]):
        in_scope = True
        continue
    if in_scope and row[3] and str(row[3]).strip() and '>' not in str(row[3]):
        systems_in_scope.append({
            "ext_sid":    str(row[3]).strip(),
            "sys_num":    str(row[4]).strip() if row[4] else "",
            "sys_type":   str(row[5]).strip() if row[5] else "",
            "sys_role":   str(row[6]).strip() if row[6] else "",
            "datacenter": str(row[7]).strip() if row[7] else "",
            "eu_dp":      str(row[8]).strip() if row[8] else "",
            "gxp":        str(row[9]).strip() if row[9] else "",
        })

# --- System info ---
si = wb_data["System info."]
si_rows = list(si.iter_rows(values_only=True))

def get_si_value(label):
    for row in si_rows:
        for idx, val in enumerate(row):
            if val and label.lower().rstrip(':') in str(val).lower():
                nxt = idx + 1
                if nxt < len(row) and row[nxt]:
                    return str(row[nxt]).strip()
    return ""

sys_type   = get_si_value("System Type")
sys_role   = get_si_value("System Role")
db_type    = get_si_value("Database Type")
db_name    = get_si_value("Database Name")
net_seg    = get_si_value("Network Segment ID")
main_fqdn  = get_si_value("Main Appl Host FQDN")     # FIX 1

# --- BASIS Checks ---
bc      = wb_data["BASIS Checks"]
bc_rows = list(bc.iter_rows(values_only=True))

basis_sid = str(bc_rows[0][0]).strip() if bc_rows[0][0] else ""
basis_product = basis_version = basis_db_type = basis_os_type = basis_ewa = ""
for row in bc_rows[:8]:
    if not row[1]:
        continue
    lbl = str(row[1]).strip().rstrip(':').lower()
    val = str(row[2]).strip() if row[2] else ""
    if "lead product"  in lbl: basis_product  = val
    elif "version"     in lbl: basis_version  = val
    elif "db type"     in lbl: basis_db_type  = val
    elif "os type"     in lbl: basis_os_type  = val
    elif "ewa"         in lbl: basis_ewa       = val

basis_checks = []
for row in bc_rows:
    if not row[3]:
        continue
    lbl = str(row[3]).strip()
    if lbl.lower() == "checks:":
        continue
    result  = str(row[4]).strip() if row[4] else ""
    details = str(row[5]).strip() if row[5] else ""   # full text, no truncation
    if lbl:
        basis_checks.append((lbl, result, details))

# --- Recommendations ---
rec      = wb_data["Recommendations"]
rec_rows = list(rec.iter_rows(values_only=True))

rec_table_starts = [i for i, row in enumerate(rec_rows) if row[1] and "Item No" in str(row[1])]

def parse_rec_table(start_idx):
    items = []
    for row in rec_rows[start_idx+1:]:
        if not any(c for c in row if c is not None):
            break
        items.append(tuple(str(row[j]).strip() if row[j] else "" for j in range(1, 9)))
    return items

basis_recs = parse_rec_table(rec_table_starts[0]) if len(rec_table_starts) > 0 else []
db_recs    = parse_rec_table(rec_table_starts[1]) if len(rec_table_starts) > 1 else []
sm_recs    = parse_rec_table(rec_table_starts[2]) if len(rec_table_starts) > 2 else []

print(f"Customer: {customer_name}")
print(f"Main FQDN: {main_fqdn}")
print(f"BASIS checks: {len(basis_checks)}, Recs BASIS/DB/SM: {len(basis_recs)}/{len(db_recs)}/{len(sm_recs)}")

# ── 3. Extract embedded images from xlsx ─────────────────────────────────────
def extract_sheet_image(sheet_rels_path):
    """Extract first image referenced by a sheet drawing."""
    with zipfile.ZipFile(XLSX_COPY, 'r') as z:
        rels_xml = z.read(sheet_rels_path).decode('utf-8')
        root = etree.fromstring(rels_xml.encode())
        drawing_path = None
        for rel in root:
            t = rel.get('Target', '')
            if 'drawing' in t:
                drawing_path = t.replace('..', 'xl')
                break
        if not drawing_path:
            return None
        draw_name = drawing_path.split('/')[-1]
        dr_rels = z.read(f'xl/drawings/_rels/{draw_name}.rels').decode('utf-8')
        dr_root = etree.fromstring(dr_rels.encode())
        img_path = None
        for rel in dr_root:
            t = rel.get('Target', '')
            if 'media/image' in t:
                img_path = t.replace('../', 'xl/')
                break
        if not img_path:
            return None
        return Image.open(io.BytesIO(z.read(img_path))).copy()

sm_img = extract_sheet_image('xl/worksheets/_rels/sheet11.xml.rels')
db_img = extract_sheet_image('xl/worksheets/_rels/sheet9.xml.rels')

sm_img_path = os.path.join(IMG_DIR, "sm_checks.png")
db_img_path = os.path.join(IMG_DIR, "db_checks.png")
if sm_img:
    sm_img.save(sm_img_path)
    print(f"SM Checks image: {sm_img.size}")
if db_img:
    db_img.save(db_img_path)
    print(f"DB Checks image: {db_img.size}")

# ── 4. Excel COM screenshots for Recommendations ──────────────────────────────
def excel_range(start_0based, end_0based):
    return f"B{start_0based+1}:I{end_0based+1}"

basis_range = excel_range(rec_table_starts[0], rec_table_starts[0] + len(basis_recs)) if rec_table_starts else "B3:I11"
db_range    = excel_range(rec_table_starts[1], rec_table_starts[1] + len(db_recs))    if len(rec_table_starts) > 1 else "B15:I20"
sm_range    = excel_range(rec_table_starts[2], rec_table_starts[2] + len(sm_recs))    if len(rec_table_starts) > 2 else "B23:I26"

basis_rec_img = os.path.join(IMG_DIR, "rec_basis.png")
db_rec_img    = os.path.join(IMG_DIR, "rec_db.png")
sm_rec_img    = os.path.join(IMG_DIR, "rec_sm.png")

print(f"Taking Recommendations screenshots... BASIS:{basis_range} DB:{db_range} SM:{sm_range}")

xl = win32com.client.Dispatch("Excel.Application")
xl.Visible = True
xl.DisplayAlerts = False
time.sleep(1.5)
wb_xl = xl.Workbooks.Open(os.path.abspath(XLSX_COPY))
ws_rec = wb_xl.Sheets("Recommendations")
ws_rec.Activate()

for rng_addr, out_path in [(basis_range, basis_rec_img), (db_range, db_rec_img), (sm_range, sm_rec_img)]:
    ws_rec.Range(rng_addr).CopyPicture(Appearance=1, Format=2)
    time.sleep(1.5)
    img = ImageGrab.grabclipboard()
    if img:
        img.save(out_path)
        print(f"  Saved {os.path.basename(out_path)}: {img.size}")
    else:
        print(f"  WARNING: clipboard empty for {rng_addr}")

wb_xl.Close(False)
xl.Quit()
time.sleep(1.0)

# ── 5. Helper functions ───────────────────────────────────────────────────────
def set_cell_text(cell_obj, txt):
    """Clear all runs in cell's first paragraph, write txt preserving first run's format."""
    tf    = cell_obj.text_frame
    para  = tf.paragraphs[0]
    p_el  = para._p
    ns    = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    runs  = p_el.findall(f'{{{ns}}}r')
    if runs:
        first_rPr = runs[0].find(f'{{{ns}}}rPr')
        rPr_copy  = etree.fromstring(etree.tostring(first_rPr)) if first_rPr is not None else None
        for r in runs:
            p_el.remove(r)
        run = para.add_run()
        run.text = txt
        if rPr_copy is not None:
            r_el = run._r
            existing = r_el.find(f'{{{ns}}}rPr')
            if existing is not None:
                r_el.remove(existing)
            r_el.insert(0, rPr_copy)
    else:
        para.add_run().text = txt

def update_para_value(para, new_value):
    """Update the value run (run1) of a label+value paragraph, keep label run (run0) intact."""
    runs = para.runs
    if len(runs) >= 2:
        runs[1].text = '\xa0' + new_value
    elif len(runs) == 1:
        # Single run: replace value part after first colon
        text = runs[0].text
        if '\xa0' in text:
            label = text.split('\xa0')[0]
            runs[0].text = label + '\xa0' + new_value
        else:
            runs[0].text = text + '\xa0' + new_value

def delete_slide(prs, idx):
    """Delete slide at 0-based index from presentation."""
    NS_R     = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldIdLst = prs.slides._sldIdLst
    sldId    = sldIdLst[idx]
    rId      = sldId.get(f'{{{NS_R}}}id')
    sldIdLst.remove(sldId)
    prs.part._rels.pop(rId)

def remove_hyperlinks(shape):
    """Strip all hlinkClick elements from a shape's text runs."""
    if not hasattr(shape, 'text_frame'):
        return
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(f'{{{ns}}}rPr')
            if rPr is not None:
                for hl in rPr.findall(f'{{{ns}}}hlinkClick'):
                    rPr.remove(hl)
        p_el = para._p
        for endRPr in p_el.findall(f'{{{ns}}}endParaRPr'):
            for hl in endRPr.findall(f'{{{ns}}}hlinkClick'):
                endRPr.remove(hl)

def fit_image(img_path, avail_w_emu, avail_h_emu):
    """Return (w, h) in EMU that fits img within available area, preserving aspect."""
    img = Image.open(img_path)
    iw, ih = img.size
    asp = iw / ih
    if avail_w_emu / asp <= avail_h_emu:
        w = avail_w_emu
        h = int(w / asp)
    else:
        h = avail_h_emu
        w = int(h * asp)
    return w, h

def replace_picture(slide, old_shape, new_img_path, left=None, top=None, width=None, height=None):
    """Remove old picture shape and insert new one in same position unless overrides given."""
    l = left   if left   is not None else old_shape.left
    t = top    if top    is not None else old_shape.top
    w = width  if width  is not None else old_shape.width
    h = height if height is not None else old_shape.height
    old_shape._element.getparent().remove(old_shape._element)
    return slide.shapes.add_picture(new_img_path, l, t, w, h)

def remove_tables_and_pics(slide):
    to_rm = [s._element for s in slide.shapes if s.shape_type in (19, 13)]
    for el in to_rm:
        el.getparent().remove(el)

def insert_screenshot(slide, img_path, top_offset_inches=1.6):
    """Insert an image centered and scaled to fit the slide body area."""
    if not os.path.exists(img_path):
        return
    pad   = Inches(0.15)
    w, h  = fit_image(img_path,
                      SLIDE_W - 2*pad,
                      SLIDE_H - Inches(top_offset_inches) - Inches(0.3))
    left  = (SLIDE_W - w) // 2
    top   = Inches(top_offset_inches)
    slide.shapes.add_picture(img_path, left, top, w, h)

# ── 6. Build PPTX from template ───────────────────────────────────────────────
now        = datetime.datetime.now()
month_year = now.strftime("%B %Y")

prs     = Presentation(PPTX_TEMPLATE)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ── Slide 1: Cover ────────────────────────────────────────────────────────────
for shape in prs.slides[0].shapes:
    if "Title" in shape.name and hasattr(shape, 'text_frame'):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "Innova Sport" in run.text:
                    run.text = run.text.replace("Innova Sport S.A De C.V.", customer_name)
    if "Date" in shape.name and hasattr(shape, 'text_frame'):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = month_year

# ── Slide 3: Customer Card Overview ──────────────────────────────────────────
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.name == "TextBox 9" and hasattr(shape, 'text_frame'):
        paras = shape.text_frame.paragraphs
        line_values = [
            f"CRM ID: {crm_id} \tCDM: {cdm_name} TSM: {tsm_name} ",
            f"SUBCONTRACTOR: {subcontractor if subcontractor else 'N/A'} ",
            f"ERP ID: {erp_id}",
        ]
        for i, para in enumerate(paras):
            if i < len(line_values) and para.runs:
                para.runs[0].text = line_values[i]

    if shape.name == "TextBox 20" and hasattr(shape, 'text_frame'):
        paras = shape.text_frame.paragraphs
        # Map paragraph index -> new value  (FIX 1: index 5 = MAIN APPL HOST FQDN)
        value_map = {
            0: sys_type,
            1: db_type,
            2: sys_role,
            3: db_name,
            4: net_seg,
            5: main_fqdn,    # FIX 1
        }
        for i, para in enumerate(paras):
            if i in value_map:
                update_para_value(para, value_map[i])

# ── Slide 6: BASIS Checks ─────────────────────────────────────────────────────
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.shape_type != 19:
        continue
    tbl   = shape.table
    nrows = len(tbl.rows)

    if nrows == 6:
        # System header table: SID + product info rows
        set_cell_text(tbl.rows[0].cells[0], basis_sid)
        info_map = {1: basis_product, 2: basis_version, 3: basis_db_type, 4: basis_os_type, 5: basis_ewa}
        for r, val in info_map.items():
            if r < nrows:
                set_cell_text(tbl.rows[r].cells[2], val)

    elif nrows >= 4:
        # Checks table: row0=header, rows1+ = check items  (FIX 2: full details text)
        for ci_idx, (check_lbl, result, details) in enumerate(basis_checks):
            row_idx = ci_idx + 1
            if row_idx >= nrows:
                break
            row = tbl.rows[row_idx]
            set_cell_text(row.cells[0], check_lbl)
            set_cell_text(row.cells[1], result)
            set_cell_text(row.cells[2], details)   # FIX 2: full text, all runs cleared

# ── Slide 7: DB Checks  (FIX 3: replace screenshot with properly sized image) ─
slide7 = prs.slides[6]

# Update SID text box (TextBox 5)
for shape in slide7.shapes:
    if shape.name == "TextBox 5" and hasattr(shape, 'text_frame'):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = ""
        paras = shape.text_frame.paragraphs
        if paras:
            if paras[0].runs:
                paras[0].runs[0].text = ""
            # Write SID on what appears as last non-empty para
            for para in reversed(paras):
                if para.runs:
                    para.runs[0].text = basis_sid
                    break

# Update header text box (TextBox 3)
for shape in slide7.shapes:
    if shape.name == "TextBox 3" and hasattr(shape, 'text_frame'):
        tf = shape.text_frame
        # Build label->value pairs matching the runs structure
        replacements = [
            ("SAP HANA, platform edition", basis_product),
            ("2.0 / 079 / 08 (2.007908)",  basis_version),
            ("Linux",                       basis_os_type),
            ("HANA MDC (Multi-Tenant Database Container)", db_type),
            ("HANA MDC",                    db_type),
        ]
        for para in tf.paragraphs:
            for run in para.runs:
                for old, new in replacements:
                    if old in run.text:
                        run.text = run.text.replace(old, new)

# FIX 3: Replace existing picture with DB Checks image, properly sized
if os.path.exists(db_img_path):
    # Find the picture shape to get its bounds as reference for available area
    pic_shape = next((s for s in slide7.shapes if s.shape_type == 13), None)

    # Available area: below the lower horizontal line (Straight Connector 29 at top=2.58")
    # Footer text is at top=7.07"
    area_top    = Inches(2.68)
    area_bottom = Inches(7.0)
    area_left   = Inches(0.3)
    area_right  = SLIDE_W - Inches(0.3)
    avail_w     = area_right  - area_left
    avail_h     = area_bottom - area_top

    w, h = fit_image(db_img_path, avail_w, avail_h)
    left = area_left + (avail_w - w) // 2
    top  = area_top  + (avail_h - h) // 2

    if pic_shape:
        pic_shape._element.getparent().remove(pic_shape._element)
    slide7.shapes.add_picture(db_img_path, left, top, w, h)

# ── Slide 8: SM Checks (one slide only, FIX 4 removes slide 9 duplicate) ──────
if os.path.exists(sm_img_path):
    remove_tables_and_pics(prs.slides[7])
    insert_screenshot(prs.slides[7], sm_img_path, top_offset_inches=1.3)

# ── Slide 9: SM Checks duplicate  (FIX 4: will be deleted below) ─────────────
# ── Slide 11: Recommendations - BASIS ────────────────────────────────────────
remove_tables_and_pics(prs.slides[10])
insert_screenshot(prs.slides[10], basis_rec_img, top_offset_inches=1.6)

# ── Slide 12: Recommendations - Database ──────────────────────────────────────
remove_tables_and_pics(prs.slides[11])
insert_screenshot(prs.slides[11], db_rec_img, top_offset_inches=1.4)

# ── Slide 13: Recommendations - Database duplicate (FIX 5: deleted below) ────

# ── Slide 14: Recommendations - SM ────────────────────────────────────────────
remove_tables_and_pics(prs.slides[13])
insert_screenshot(prs.slides[13], sm_rec_img, top_offset_inches=1.6)

# ── Slide 15: Attachments - update filename + remove hyperlink  (FIX 6) ───────
slide15 = prs.slides[14]
xlsx_base = os.path.basename(XLSX_PATH)
for shape in slide15.shapes:
    if hasattr(shape, 'text_frame'):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "Innova Sport" in run.text:
                    run.text = run.text.replace("PLA_Innova Sport S.A De C.V.xlsx", xlsx_base)
        remove_hyperlinks(shape)   # FIX 6

# ── FIX 4 & 5: Delete duplicate slides (highest index first to avoid shift) ───
# Current slide layout (16 total):
#   index 12 = Slide 13 - Recommendations Database (duplicate)
#   index  8 = Slide  9 - Server Management Checks (duplicate)
print("Deleting duplicate slides...")
delete_slide(prs, 12)   # Rec Database duplicate (delete highest index first)
delete_slide(prs,  8)   # SM Checks duplicate

print(f"Slides after deletion: {len(prs.slides)}")

# ── 7. Save ───────────────────────────────────────────────────────────────────
safe_name = customer_name.replace("/", "-").replace("\\", "-")
out_name  = f"360_{safe_name}_{now.strftime('%Y%m')}.pptx"
out_path  = os.path.join(BASE_DIR, out_name)
prs.save(out_path)
print(f"\nSaved: {out_path}")
